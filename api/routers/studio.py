"""Studio 路由：/api/studio/*

为 v5 PPT 共创工作台提供专用 API：
- 缩略图渲染
- 差异预览
- 全屏预览数据
"""
import json
import base64
import tempfile
import os
from typing import List, Dict, Any, Optional

from fastapi import APIRouter, HTTPException
from pydantic import BaseModel, Field

router = APIRouter(prefix="/api/studio", tags=["studio"])

# 内存中的 Studio 会话状态（轻量级，后续可迁移到 Redis/DB）
_studio_sessions: Dict[str, Dict[str, Any]] = {}


# =============================================================================
# Pydantic 模型
# =============================================================================

class ThumbnailRequest(BaseModel):
    """缩略图渲染请求"""
    slide: Dict[str, Any] = Field(..., description="单页 slide JSON")
    width: int = Field(320, description="缩略图宽度")
    height: int = Field(180, description="缩略图高度")


class DiffPreviewRequest(BaseModel):
    """差异预览请求"""
    original_slides: List[Dict[str, Any]] = Field(..., description="原始 slides")
    modified_slides: List[Dict[str, Any]] = Field(..., description="修改后 slides")


class FullscreenRequest(BaseModel):
    """全屏预览请求"""
    session_id: str = Field(..., description="Studio 会话 ID")


class StudioSessionUpdate(BaseModel):
    """Studio 会话状态更新"""
    session_id: str = Field(..., description="会话 ID")
    slides: List[Dict[str, Any]] = Field(..., description="当前 slides 数据")
    source_text: Optional[str] = Field(None, description="原始文本")


# =============================================================================
# 端点
# =============================================================================

@router.post("/thumbnail-render")
async def render_thumbnail(request: ThumbnailRequest):
    """
    渲染单页 slide 缩略图

    使用 ppt_generator 渲染单页为图片，返回 base64 编码。
    若渲染引擎不可用，返回结构化文本描述作为降级方案。
    """
    print(f"[V5-API] POST /api/studio/thumbnail-render | slide keys={list(request.slide.keys())}")
    try:
        # 尝试用 pptx 渲染单页为图片
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from ppt_generator import _build_slide  # 复用现有 slide 构建器

        prs = Presentation()
        prs.slide_width = Emu(int(request.width * 914400 / 96))   # px → emu
        prs.slide_height = Emu(int(request.height * 914400 / 96))

        blank_layout = prs.slide_layouts[6]  # blank
        slide_obj = prs.slides.add_slide(blank_layout)
        _build_slide(slide_obj, request.slide, prs)

        # 保存为临时 pptx，用 LibreOffice 或 Pillow 转图片
        tmp_pptx = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        prs.save(tmp_pptx.name)
        tmp_pptx.close()

        # 尝试用 pdf 转换获取图片
        try:
            import subprocess
            tmp_png = tmp_pptx.name + ".png"
            # 使用 soffice 转 PDF → 再转 PNG（降级方案跳过）
            result = subprocess.run(
                ["soffice", "--headless", "--convert-to", "png", "--outdir",
                 os.path.dirname(tmp_pptx.name), tmp_pptx.name],
                capture_output=True, timeout=10
            )
            if os.path.exists(tmp_png):
                with open(tmp_png, "rb") as f:
                    img_b64 = base64.b64encode(f.read()).decode()
                os.unlink(tmp_png)
                os.unlink(tmp_pptx.name)
                print(f"[V5-API] thumbnail-render: rendered via soffice, size={len(img_b64)}")
                return {"image_base64": img_b64, "format": "png", "render_mode": "soffice"}
        except Exception:
            pass

        # soffice 不可用 → 返回结构化文本描述
        os.unlink(tmp_pptx.name)
        return _text_thumbnail(request.slide, request.width, request.height)

    except ImportError:
        # python-pptx 不可用 → 返回文本描述
        print("[V5-API] thumbnail-render: pptx not available, using text fallback")
        return _text_thumbnail(request.slide, request.width, request.height)
    except Exception as e:
        print(f"[V5-API] thumbnail-render error: {e}")
        raise HTTPException(status_code=500, detail=f"缩略图渲染失败: {str(e)}")


@router.post("/diff-preview")
async def diff_preview(request: DiffPreviewRequest):
    """
    生成 slides 差异预览

    对比 original 和 modified slides，返回逐页差异描述。
    """
    print(f"[V5-API] POST /api/studio/diff-preview | original={len(request.original_slides)}, modified={len(request.modified_slides)}")
    try:
        diffs = []
        orig_count = len(request.original_slides)
        mod_count = len(request.modified_slides)

        # 逐页对比
        max_pages = max(orig_count, mod_count)
        for i in range(max_pages):
            if i >= orig_count:
                diffs.append({
                    "page": i,
                    "status": "added",
                    "description": f"新增第 {i+1} 页",
                    "slide": request.modified_slides[i],
                })
            elif i >= mod_count:
                diffs.append({
                    "page": i,
                    "status": "removed",
                    "description": f"删除第 {i+1} 页",
                    "slide": request.original_slides[i],
                })
            else:
                orig = request.original_slides[i]
                mod = request.modified_slides[i]
                changes = _compare_slides(orig, mod)
                if changes:
                    diffs.append({
                        "page": i,
                        "status": "modified",
                        "description": f"第 {i+1} 页有 {len(changes)} 处修改",
                        "changes": changes,
                    })
                else:
                    diffs.append({
                        "page": i,
                        "status": "unchanged",
                        "description": f"第 {i+1} 页无变化",
                    })

        summary = {
            "total_pages": max_pages,
            "added": sum(1 for d in diffs if d["status"] == "added"),
            "removed": sum(1 for d in diffs if d["status"] == "removed"),
            "modified": sum(1 for d in diffs if d["status"] == "modified"),
            "unchanged": sum(1 for d in diffs if d["status"] == "unchanged"),
        }

        print(f"[V5-API] diff-preview: {summary}")
        return {"diffs": diffs, "summary": summary}
    except Exception as e:
        print(f"[V5-API] diff-preview error: {e}")
        raise HTTPException(status_code=500, detail=f"差异预览失败: {str(e)}")


@router.get("/fullscreen-preview")
async def fullscreen_preview(session_id: str):
    """
    获取全屏预览数据

    返回指定会话的当前 slides 全量 JSON，供前端全屏渲染。
    """
    print(f"[V5-API] GET /api/studio/fullscreen-preview | session_id={session_id}")
    try:
        session = _studio_sessions.get(session_id)
        if not session:
            # 降级：返回空 slides
            print(f"[V5-API] fullscreen-preview: session {session_id} not found")
            return {
                "session_id": session_id,
                "slides": [],
                "slide_count": 0,
                "status": "no_session",
            }
        slides = session.get("slides", [])
        print(f"[V5-API] fullscreen-preview: {len(slides)} slides")
        return {
            "session_id": session_id,
            "slides": slides,
            "slide_count": len(slides),
            "source_text": session.get("source_text", ""),
            "status": "ok",
        }
    except Exception as e:
        print(f"[V5-API] fullscreen-preview error: {e}")
        raise HTTPException(status_code=500, detail=f"全屏预览失败: {str(e)}")


@router.post("/session")
async def update_studio_session(request: StudioSessionUpdate):
    """
    创建或更新 Studio 会话状态

    将当前 slides 数据保存到内存会话，供全屏预览和差异对比使用。
    """
    print(f"[V5-API] POST /api/studio/session | session_id={request.session_id}, slides={len(request.slides)}")
    _studio_sessions[request.session_id] = {
        "slides": request.slides,
        "source_text": request.source_text or "",
        "updated_at": __import__("datetime").datetime.now().isoformat(),
    }
    return {"success": True, "session_id": request.session_id,
            "slide_count": len(request.slides)}


# =============================================================================
# 工具函数
# =============================================================================

def _text_thumbnail(slide: dict, width: int, height: int) -> dict:
    """生成文本描述作为缩略图降级方案"""
    title = slide.get("title", "(无标题)")
    layout = slide.get("layout", "default")
    items = slide.get("items", [])
    content = slide.get("content", "")

    lines = [f"📊 {title}"]
    if layout:
        lines.append(f"版式: {layout}")
    if items:
        for item in items[:5]:
            text = item.get("text", "")
            lines.append(f"  • {text}")
        if len(items) > 5:
            lines.append(f"  ... 等 {len(items)} 项")
    if content and not items:
        lines.append(content[:100])

    return {
        "text_preview": "\n".join(lines),
        "width": width,
        "height": height,
        "render_mode": "text",
        "slide_keys": list(slide.keys()),
    }


def _compare_slides(original: dict, modified: dict) -> list:
    """对比两页 slide 的差异"""
    changes = []

    # 对比标题
    orig_title = original.get("title", "")
    mod_title = modified.get("title", "")
    if orig_title != mod_title:
        changes.append({
            "field": "title",
            "old": orig_title,
            "new": mod_title,
        })

    # 对比内容
    orig_content = original.get("content", "")
    mod_content = modified.get("content", "")
    if orig_content != mod_content:
        changes.append({
            "field": "content",
            "old": orig_content[:100],
            "new": mod_content[:100],
        })

    # 对比 items
    orig_items = original.get("items", [])
    mod_items = modified.get("items", [])
    if orig_items != mod_items:
        changes.append({
            "field": "items",
            "old_count": len(orig_items),
            "new_count": len(mod_items),
        })

    # 对比 layout
    if original.get("layout") != modified.get("layout"):
        changes.append({
            "field": "layout",
            "old": original.get("layout"),
            "new": modified.get("layout"),
        })

    return changes
