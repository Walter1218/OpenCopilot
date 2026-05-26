"""Office 文档解析探针：支持 .docx 和 .pptx 文本提取，用于全文修订模式的上下文注入。"""
import asyncio
import os


async def read_office_file(file_path: str) -> dict:
    """从 .docx 或 .pptx 文件中提取纯文本内容。"""
    expanded = os.path.expanduser(file_path)
    if not os.path.exists(expanded):
        raise FileNotFoundError(f"文件未找到: {file_path}")

    ext = os.path.splitext(expanded)[1].lower()
    loop = asyncio.get_running_loop()

    if ext == '.docx':
        return await loop.run_in_executor(None, _extract_docx, expanded)
    elif ext == '.pptx':
        return await loop.run_in_executor(None, _extract_pptx, expanded)
    else:
        raise ValueError(f"不支持的文件格式: {ext}（当前仅支持 .docx / .pptx）")


def _extract_docx(file_path: str) -> dict:
    from docx import Document
    doc = Document(file_path)
    parts = [p.text for p in doc.paragraphs if p.text.strip()]

    # 也提取表格内容
    for table in doc.tables:
        for row in table.rows:
            row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)

    return {
        "type": "docx",
        "content": '\n\n'.join(parts),
        "paragraph_count": len(doc.paragraphs),
        "table_count": len(doc.tables),
    }


def _extract_pptx(file_path: str) -> dict:
    from pptx import Presentation
    prs = Presentation(file_path)
    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        texts.append(t)
        if texts:
            slides.append(f"--- 幻灯片 {i + 1} ---\n" + '\n'.join(texts))

    return {
        "type": "pptx",
        "content": '\n\n'.join(slides),
        "slide_count": len(prs.slides),
    }
