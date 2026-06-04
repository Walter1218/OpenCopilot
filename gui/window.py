"""AICardWindow - 核心悬浮卡片"""
import os, sys, json, uuid, time, tempfile, re, traceback
from typing import Dict, Any, Optional, List
from PyQt6.QtWidgets import (QWidget, QVBoxLayout, QHBoxLayout, QTextEdit, QLabel, QPushButton, QFrame, QGraphicsDropShadowEffect, QFileDialog, QMessageBox, QComboBox, QGroupBox, QFormLayout, QSplitter, QListWidget, QListWidgetItem)
from PyQt6.QtCore import Qt, pyqtSignal, QTimer
from PyQt6.QtGui import QCursor, QColor
import httpx
from llm_provider import load_config, save_config, ProviderFactory
from opencopilot.agent.caller import call_agent_pipeline_sync
from markdown_renderer import render as md_render
from system_probe_client import SystemProbeClient
from ppt_generator import generate_ppt_from_json, extract_json_from_text
from opencopilot.capabilities.ppt import CoCreationDialog
from core.theme_manager import ThemeManager
from core.shortcut_manager import ShortcutManager
from widgets.file_drop_zone import FileDropZone
from widgets.progress_widget import ProgressWidget, MultiStepProgressWidget
from widgets.context_menu import TextContextMenu, FileContextMenu, CodeContextMenu
from widgets.batch_dialog import BatchDialog
from widgets.terminology_dialog import TerminologyDialog
from widgets.translation_memory import TranslationMemory
from widgets.skill_panel import SkillPanel, SkillSearchWidget, SkillCommandParser
from widgets.skill_context_menu import SkillContextMenu, SkillCommandWidget
from widgets.skill_search_dialog import SkillSearchDialog, SkillQuickAccessWidget
from opencopilot.capabilities.skill import SkillRegistry, SkillContext, IntentRouter, SkillExecutor
from gui.shared import check_accessibility_permission, make_panel_persistent
from gui.workers.ai import AIWorker
from gui.workers.chat import ChatWorker
from gui.workers.broker import BrokerEventsWorker
from gui.workers.browser import BrowserReaderWorker
from gui.workers.scanner import ModelScannerWorker
from gui.dialogs.settings import SettingsDialog
from gui.dialogs.ppt_preview import PPTPreviewDialog
class AICardWindow(QWidget):
    ide_probe_result = pyqtSignal(bool)
    browser_probe_result = pyqtSignal(str)

    def __init__(self, provider):
        super().__init__()
        self.provider = provider
        self.worker = None
        self.chat_worker = None
        self.current_text = ""
        self.chat_history = []
        self.session_id = str(uuid.uuid4())
        self.active_browser = None
        self._temp_chat_pos = 0
        self._pending_hide = False
        self._user_initiated_hide = False  # 区分用户主动隐藏 vs 系统自动隐藏
        # 上下文感知
        self.context_source = "drag"
        self.context_meta = {}
        self.task_context = ""  # 工作台注入的任务上下文
        self._ide_selection_range = None  # IDE 选区范围（用于局部回写）
        self._ide_full_document = ""  # IDE 全文缓存（用于上下文和定位）
        self._agent_online = True  # 默认乐观假设在线，探活结果回来后更新
        self._allow_close = False
        self._is_dragging = False  # 是否正在拖拽
        # 文档修订模式
        self.revision_mode = False
        self.full_document = ""  # 修订模式下的全文档缓存
        
        # 初始化新组件
        self.theme_manager = ThemeManager()
        self.shortcut_manager = ShortcutManager()
        self.translation_memory = TranslationMemory()
        self.progress_widget = ProgressWidget()
        self.multi_step_progress = MultiStepProgressWidget()
        
        # 初始化 Skill 架构
        self.skill_registry = SkillRegistry()
        self.skill_router = IntentRouter(self.skill_registry)
        self.skill_executor = SkillExecutor(self.skill_registry, self.skill_router)
        self._init_skills()
        
        self.current_active_app = ""
        self.current_bundle_id = ""
        self.recent_apps = []  # 记录最近激活的应用历史
        self.broker_events_worker = BrokerEventsWorker()
        self.broker_events_worker.app_activated.connect(self._on_app_activated)
        self.broker_events_worker.start()

        self.initUI()

    def _on_app_activated(self, app_name, bundle_id):
        self.current_active_app = app_name
        self.current_bundle_id = bundle_id
        
        # 维护最近打开的应用列表（去重，保持最新在最后）
        if app_name in self.recent_apps:
            self.recent_apps.remove(app_name)
        self.recent_apps.append(app_name)
        # 仅保留最近 10 个
        self.recent_apps = self.recent_apps[-10:]
        
        # 为了调试方便，在 UI 终端打印出系统焦点切换事件
        print(f"[UI 接收] 系统焦点已切换至: {app_name} ({bundle_id}) | 最近使用: {self.recent_apps}")
    
    def _init_skills(self):
        """初始化所有 Skill"""
        try:
            # 注册所有 Skill
            skills = [
                CodingSkill(),
                KnowledgeSkill(),
                PPTSkill(),
                EvaluationSkill(),
                FileSkill(),
                FormatSkill(),
                PersonaSkill()
            ]
            
            for skill in skills:
                self.skill_registry.register(skill)
                print(f"[Skill 注册] {skill.metadata.display_name}")
            
            print(f"[Skill 注册完成] 共注册 {len(skills)} 个技能")
            
        except Exception as e:
            print(f"[Skill 注册失败] {e}")
            import traceback
            traceback.print_exc()

    def initUI(self):
        # 无边框、置顶
        # 不用 Qt.Tool：macOS 会在父窗口失焦时自动隐藏 Tool 窗口
        #   且 Tool + hidesOnDeactivate=False 后跨窗口拖拽失效
        # 用 WA_ShowWithoutActivating 控制首次显示不抢焦点
        self.setWindowFlags(
            Qt.WindowType.FramelessWindowHint |
            Qt.WindowType.WindowStaysOnTopHint
        )
        self.setAttribute(Qt.WidgetAttribute.WA_TranslucentBackground)
        self.setAttribute(Qt.WidgetAttribute.WA_ShowWithoutActivating)
        self.setAcceptDrops(True)
        self.setMouseTracking(True)  # 启用鼠标追踪，用于拖拽缩放
        self._resize_margin = 14
        self._resizing = False
        self._resize_edge = None
        self._resize_start_geo = None
        self._resize_start_pos = None

        self.resize(680, 520)
        self._drag_pos = None  # 拖动起始位置
        self.frame = QFrame(self)
        self.frame.setStyleSheet("""
            QFrame {
                background-color: rgba(30, 30, 35, 240);
                border-radius: 12px;
                border: 1px solid rgba(100, 100, 100, 100);
            }
        """)
        self.frame.resize(660, 500)
        self.frame.move(10, 10)
        
        shadow = QGraphicsDropShadowEffect()
        shadow.setBlurRadius(20)
        shadow.setColor(QColor(0, 0, 0, 180))
        shadow.setOffset(0, 5)
        self.frame.setGraphicsEffect(shadow)

        layout = QVBoxLayout(self.frame)
        layout.setContentsMargins(15, 15, 15, 15)

        title_layout = QHBoxLayout()
        self.title_label = QLabel("✨ Smart Copilot  (可拖动)", self)
        self.title_label.setStyleSheet("color: #4da6ff; font-weight: bold; font-size: 14px; background: transparent; border: none;")
        self.title_label.setCursor(Qt.CursorShape.OpenHandCursor)

        # Agent 在线状态指示灯（绿色=在线，红色=离线）
        self.agent_status_dot = QLabel("●", self)
        self.agent_status_dot.setStyleSheet("color: #4caf50; font-size: 10px; background: transparent; border: none;")
        self.agent_status_dot.setToolTip("ASU 核心守护服务在线")
        
        self.btn_persona = QPushButton("🎭", self)
        self.btn_persona.setStyleSheet("""
            QPushButton { background: transparent; border: none; font-size: 14px; }
            QPushButton:hover { color: #fff; }
        """)
        self.btn_persona.setCursor(Qt.CursorShape.PointingHandCursor)
        self.btn_persona.setToolTip("角色工坊")
        self.btn_persona.clicked.connect(self.open_persona_workshop)
        
        self.btn_settings = QPushButton("⚙️", self)
        self.btn_settings.setStyleSheet("""
            QPushButton { background: transparent; border: none; font-size: 14px; }
            QPushButton:hover { color: #fff; }
        """)
        self.btn_settings.setCursor(Qt.CursorShape.PointingHandCursor)
        self.btn_settings.clicked.connect(self.open_settings)
        
        self.btn_close = QPushButton("✕", self)
        self.btn_close.setStyleSheet("""
            QPushButton { background: transparent; border: none; font-size: 14px; color: #888; }
            QPushButton:hover { color: #ff5555; }
        """)
        self.btn_close.setCursor(Qt.CursorShape.PointingHandCursor)
        self.btn_close.clicked.connect(self.hide_card)

        title_layout.addWidget(self.title_label)
        title_layout.addWidget(self.agent_status_dot)
        title_layout.addStretch()
        title_layout.addWidget(self.btn_settings)
        title_layout.addWidget(self.btn_close)
        layout.addLayout(title_layout)

        # Agent 离线提示横幅（默认隐藏，探活失败时显示）
        self.agent_offline_banner = QLabel("⚠️ ASU 核心守护服务未启动，请运行 install_daemon.sh 或手动启动 Agent。", self)
        self.agent_offline_banner.setStyleSheet("""
            QLabel {
                background-color: rgba(200, 80, 50, 200);
                color: #fff;
                border-radius: 6px;
                padding: 6px 10px;
                font-size: 11px;
            }
        """)
        self.agent_offline_banner.setWordWrap(True)
        self.agent_offline_banner.hide()
        layout.addWidget(self.agent_offline_banner)

        # --- TabWidget ---
        self.tabs = QTabWidget(self.frame)
        self.tabs.setStyleSheet("""
            QTabWidget::pane { border: none; }
            QTabBar::tab {
                background: rgba(40, 40, 45, 200);
                color: #aaa;
                padding: 6px 12px;
                border-top-left-radius: 6px;
                border-top-right-radius: 6px;
                margin-right: 2px;
            }
            QTabBar::tab:selected {
                background: rgba(60, 60, 70, 240);
                color: #fff;
                font-weight: bold;
            }
        """)
        layout.addWidget(self.tabs)

        # ==========================
        # Tab 1: 快捷划词
        # ==========================
        self.tab_quick = QWidget()
        quick_layout = QVBoxLayout(self.tab_quick)
        quick_layout.setContentsMargins(0, 10, 0, 0)
        
        # 全局悬浮操作按钮区
        global_actions_layout = QHBoxLayout()
        global_actions_layout.setContentsMargins(0, 0, 0, 5)
        global_actions_layout.setSpacing(10)
        
        # [视觉分析] 按钮
        self.btn_vision_analyze = QPushButton("👁️ 视觉分析前台", self)
        self.btn_vision_analyze.setStyleSheet("""
            QPushButton {
                background-color: rgba(147, 112, 219, 180);
                color: #fff;
                border-radius: 8px;
                padding: 6px 12px;
                font-size: 12px;
                font-weight: bold;
                border: 1px solid rgba(147, 112, 219, 255);
            }
            QPushButton:hover {
                background-color: rgba(147, 112, 219, 255);
            }
        """)
        self.btn_vision_analyze.setCursor(Qt.CursorShape.PointingHandCursor)
        self.btn_vision_analyze.clicked.connect(self._on_vision_analyze_clicked)
        global_actions_layout.addWidget(self.btn_vision_analyze)
        global_actions_layout.addStretch()
        quick_layout.addLayout(global_actions_layout)
        
        # IDE 插件探测状态栏
        self.ide_status_layout = QHBoxLayout()
        self.ide_status_layout.setContentsMargins(0, 0, 0, 5)
        self.btn_read_ide = QPushButton("📥 极速读取当前 IDE 全文")
        self.btn_read_ide.setStyleSheet("""
            QPushButton {
                background-color: rgba(77, 166, 255, 180);
                color: #fff;
                border-radius: 8px;
                padding: 6px 12px;
                font-size: 12px;
                font-weight: bold;
                border: 1px solid rgba(77, 166, 255, 255);
            }
            QPushButton:hover {
                background-color: rgba(77, 166, 255, 255);
            }
        """)
        self.btn_read_ide.setCursor(Qt.CursorShape.PointingHandCursor)
        self.btn_read_ide.clicked.connect(self.read_from_ide_extension)
        self.btn_read_ide.hide() # 默认隐藏，探测到才显示
        self.ide_probe_result.connect(self._update_ide_btn)
        
        self.ide_status_layout.addWidget(self.btn_read_ide)
        self.btn_paste_clipboard = QPushButton("📋 粘贴剪贴板")
        self.btn_paste_clipboard.setStyleSheet("""
            QPushButton {
                background-color: rgba(120, 80, 220, 180);
                color: #fff;
                border-radius: 8px;
                padding: 6px 12px;
                font-size: 12px;
                font-weight: bold;
                border: 1px solid rgba(120, 80, 220, 255);
            }
            QPushButton:hover {
                background-color: rgba(120, 80, 220, 255);
            }
        """)
        self.btn_paste_clipboard.setCursor(Qt.CursorShape.PointingHandCursor)
        self.btn_paste_clipboard.clicked.connect(self.paste_from_clipboard)
        self.ide_status_layout.addWidget(self.btn_paste_clipboard)
        self.ide_status_layout.addStretch()
        quick_layout.addLayout(self.ide_status_layout)

        # 浏览器探测状态栏
        self.browser_status_layout = QHBoxLayout()
        self.browser_status_layout.setContentsMargins(0, 0, 0, 5)
        self.btn_read_browser = QPushButton("🌐 一键读取当前网页全文")
        self.btn_read_browser.setStyleSheet("""
            QPushButton {
                background-color: rgba(255, 140, 0, 180);
                color: #fff;
                border-radius: 8px;
                padding: 6px 12px;
                font-size: 12px;
                font-weight: bold;
                border: 1px solid rgba(255, 140, 0, 255);
            }
            QPushButton:hover {
                background-color: rgba(255, 140, 0, 255);
            }
        """)
        self.btn_read_browser.setCursor(Qt.CursorShape.PointingHandCursor)
        self.btn_read_browser.clicked.connect(self.read_from_browser)
        self.btn_read_browser.hide() # 默认隐藏
        self.browser_probe_result.connect(self._update_browser_btn)
        
        self.browser_status_layout.addWidget(self.btn_read_browser)
        self.browser_status_layout.addStretch()
        quick_layout.addLayout(self.browser_status_layout)

        # 快捷指令工具栏
        self.btn_layout = QHBoxLayout()
        self.btn_layout.setContentsMargins(0, 0, 0, 5)
        self.btn_layout.setSpacing(8)
        
        button_style = """
            QPushButton {
                background-color: rgba(60, 60, 70, 200);
                color: #ddd;
                border-radius: 8px;
                padding: 4px 8px;
                font-size: 11px;
                border: 1px solid rgba(100, 100, 100, 100);
            }
            QPushButton:hover {
                background-color: rgba(80, 80, 100, 255);
                color: #fff;
                border: 1px solid #4da6ff;
            }
        """
        
        self.btn_auto = QPushButton("✨ 自动")
        self.btn_trans = QPushButton("🌐 翻译")
        self.btn_code = QPushButton("💻 代码解析")
        self.btn_polish = QPushButton("✍️ 润色")
        self.btn_revision = QPushButton("📝 全文修订")
        self.btn_revision.setCheckable(True)
        self.btn_batch = QPushButton("📦 批量处理")
        self.btn_batch.setStyleSheet(button_style)
        self.btn_batch.setCursor(Qt.CursorShape.PointingHandCursor)
        self.btn_batch.clicked.connect(self.open_batch_dialog)
        
        for btn in [self.btn_auto, self.btn_trans, self.btn_code, self.btn_polish, self.btn_revision, self.btn_batch]:
            btn.setStyleSheet(button_style)
            btn.setCursor(Qt.CursorShape.PointingHandCursor)
            self.btn_layout.addWidget(btn)
            
        self.btn_layout.addStretch()
        quick_layout.addLayout(self.btn_layout)
        
        # 绑定点击事件
        self.btn_auto.clicked.connect(lambda: self.trigger_ai("auto"))
        self.btn_trans.clicked.connect(lambda: self.trigger_ai("translate"))
        self.btn_code.clicked.connect(lambda: self.trigger_ai("code"))
        self.btn_polish.clicked.connect(lambda: self.trigger_ai("polish"))
        self.btn_revision.clicked.connect(self._toggle_revision_mode)
        
        # 文件拖拽区
        self.file_drop_zone = FileDropZone()
        self.file_drop_zone.file_dropped.connect(self._on_file_dropped)
        self.file_drop_zone.setMaximumHeight(80)
        quick_layout.addWidget(self.file_drop_zone)
        
        # 进度条（默认隐藏）
        self.progress_widget.hide()
        self.multi_step_progress.hide()
        quick_layout.addWidget(self.progress_widget)
        quick_layout.addWidget(self.multi_step_progress)

        # 自定义指令输入栏（让用户填写修改要求）
        self.instruction_layout = QHBoxLayout()
        self.instruction_layout.setContentsMargins(0, 0, 0, 5)
        self.instruction_input = QLineEdit(self.tab_quick)
        self.instruction_input.setPlaceholderText("💬 输入修改要求（如：改为 async/await、修复 bug…）按 Enter 提交")
        self.instruction_input.setFocusPolicy(Qt.FocusPolicy.ClickFocus)
        self.instruction_input.setStyleSheet("""
            QLineEdit {
                background-color: rgba(40, 40, 50, 200);
                color: #fff; border: 1px solid rgba(100, 100, 120, 150);
                border-radius: 6px; padding: 6px 10px; font-size: 12px;
            }
            QLineEdit:focus { border: 1px solid #4da6ff; }
        """)
        self.instruction_input.returnPressed.connect(self._on_custom_instruction)

        self.btn_custom_submit = QPushButton("▶", self.tab_quick)
        self.btn_custom_submit.setToolTip("提交修改指令")
        self.btn_custom_submit.setStyleSheet("""
            QPushButton {
                background-color: rgba(77, 166, 255, 180);
                color: #fff; border-radius: 6px;
                padding: 6px 14px; font-size: 13px; font-weight: bold;
            }
            QPushButton:hover { background-color: rgba(77, 166, 255, 255); }
        """)
        self.btn_custom_submit.setCursor(Qt.CursorShape.PointingHandCursor)
        self.btn_custom_submit.clicked.connect(self._on_custom_instruction)

        self.instruction_layout.addWidget(self.instruction_input, stretch=1)
        self.instruction_layout.addWidget(self.btn_custom_submit)
        quick_layout.addLayout(self.instruction_layout)

        self.text_edit = QTextEdit(self.tab_quick)
        self.text_edit.setReadOnly(True)
        self.text_edit.setTextInteractionFlags(Qt.TextInteractionFlag.TextSelectableByMouse | Qt.TextInteractionFlag.TextSelectableByKeyboard)
        self.text_edit.setFocusPolicy(Qt.FocusPolicy.NoFocus)
        self.text_edit.setStyleSheet("""
            QTextEdit {
                background-color: transparent;
                color: #eeeeee;
                font-size: 13px;
                border: none;
                line-height: 1.5;
            }
            QScrollBar:vertical {
                width: 6px;
                background: transparent;
            }
            QScrollBar::handle:vertical {
                background: rgba(255, 255, 255, 60);
                border-radius: 3px;
            }
        """)
        self.text_edit.setContextMenuPolicy(Qt.ContextMenuPolicy.CustomContextMenu)
        self.text_edit.customContextMenuRequested.connect(self._show_text_context_menu)
        quick_layout.addWidget(self.text_edit)

        # 追问按钮
        ask_more_layout = QHBoxLayout()
        self.btn_ask_more = QPushButton("💬 进一步追问")
        self.btn_ask_more.setStyleSheet(button_style)
        self.btn_ask_more.setCursor(Qt.CursorShape.PointingHandCursor)
        self.btn_ask_more.clicked.connect(self.jump_to_chat)
        ask_more_layout.addStretch()
        ask_more_layout.addWidget(self.btn_ask_more)
        quick_layout.addLayout(ask_more_layout)

        # 回写 IDE 操作栏
        apply_layout = QHBoxLayout()
        self.btn_apply_to_ide = QPushButton("📝 回写到 IDE")
        self.btn_apply_to_ide.setStyleSheet("""
            QPushButton {
                background-color: rgba(40, 167, 69, 180);
                color: #fff; border-radius: 8px;
                padding: 5px 10px; font-size: 11px; font-weight: bold;
                border: 1px solid rgba(40, 167, 69, 255);
            }
            QPushButton:hover {
                background-color: rgba(40, 167, 69, 255);
            }
        """)
        self.btn_apply_to_ide.setCursor(Qt.CursorShape.PointingHandCursor)
        self.btn_apply_to_ide.clicked.connect(self._apply_to_ide)
        self.btn_apply_to_ide.hide()  # 初始隐藏，AI 回复后显示

        self.btn_copy_result = QPushButton("📋 复制结果")
        self.btn_copy_result.setStyleSheet("""
            QPushButton {
                background-color: rgba(120, 80, 220, 180);
                color: #fff; border-radius: 8px;
                padding: 5px 10px; font-size: 11px; font-weight: bold;
                border: 1px solid rgba(120, 80, 220, 255);
            }
            QPushButton:hover {
                background-color: rgba(120, 80, 220, 255);
            }
        """)
        self.btn_copy_result.setCursor(Qt.CursorShape.PointingHandCursor)
        self.btn_copy_result.clicked.connect(self._copy_result_to_clipboard)
        self.btn_copy_result.hide()  # 初始隐藏

        self.btn_export_ppt = QPushButton("💾 导出为 PPT")
        self.btn_export_ppt.setStyleSheet("""
            QPushButton {
                background-color: rgba(255, 153, 51, 180);
                color: #fff; border-radius: 8px;
                padding: 5px 10px; font-size: 11px; font-weight: bold;
                border: 1px solid rgba(255, 153, 51, 255);
            }
            QPushButton:hover {
                background-color: rgba(255, 153, 51, 255);
            }
        """)
        self.btn_export_ppt.setCursor(Qt.CursorShape.PointingHandCursor)
        self.btn_export_ppt.clicked.connect(self._export_to_ppt)
        self.btn_export_ppt.hide()  # 初始隐藏

        apply_layout.addWidget(self.btn_apply_to_ide)
        apply_layout.addStretch()
        apply_layout.addWidget(self.btn_export_ppt)
        apply_layout.addWidget(self.btn_copy_result)
        quick_layout.addLayout(apply_layout)

        # ==========================
        # Tab 2: 连续对话
        # ==========================
        self.tab_chat = QWidget()
        chat_layout = QVBoxLayout(self.tab_chat)
        chat_layout.setContentsMargins(0, 10, 0, 0)

        self.chat_display = QTextEdit(self.tab_chat)
        self.chat_display.setReadOnly(True)
        self.chat_display.setFocusPolicy(Qt.FocusPolicy.NoFocus)
        self.chat_display.setStyleSheet(self.text_edit.styleSheet())
        chat_layout.addWidget(self.chat_display)

        input_layout = QHBoxLayout()
        self.chat_input = QLineEdit(self.tab_chat)
        self.chat_input.setPlaceholderText("输入消息，按 Enter 发送...")
        self.chat_input.setFocusPolicy(Qt.FocusPolicy.ClickFocus)
        self.chat_input.setStyleSheet("""
            QLineEdit {
                background-color: rgba(30, 30, 35, 200);
                color: #fff;
                border: 1px solid rgba(100, 100, 100, 150);
                border-radius: 4px;
                padding: 4px 8px;
                font-size: 13px;
            }
            QLineEdit:focus {
                border: 1px solid #4da6ff;
            }
        """)
        self.chat_input.returnPressed.connect(self.send_chat_message)
        
        self.btn_send = QPushButton("发送", self.tab_chat)
        self.btn_send.setStyleSheet("""
            QPushButton {
                background-color: #4da6ff;
                color: #000;
                border-radius: 4px;
                padding: 4px 12px;
                font-weight: bold;
            }
            QPushButton:hover {
                background-color: #66b3ff;
            }
        """)
        self.btn_send.setCursor(Qt.CursorShape.PointingHandCursor)
        self.btn_send.clicked.connect(self.send_chat_message)

        input_layout.addWidget(self.chat_input)
        input_layout.addWidget(self.btn_send)
        chat_layout.addLayout(input_layout)

        self.tabs.addTab(self.tab_quick, "⚡️ 快捷划词")
        self.tabs.addTab(self.tab_chat, "💬 连续对话")
        
        # ==========================
        # Tab 3: PPT Assistant
        # ==========================
        self.tab_ppt_assistant = QWidget()
        ppt_layout = QVBoxLayout(self.tab_ppt_assistant)
        ppt_layout.setContentsMargins(0, 10, 0, 0)
        
        # PPT Assistant 标题
        ppt_title = QLabel("🎯 PPT 智能助手", self.tab_ppt_assistant)
        ppt_title.setStyleSheet("color: #4da6ff; font-weight: bold; font-size: 16px; background: transparent; border: none;")
        ppt_layout.addWidget(ppt_title)
        
        # 描述
        ppt_desc = QLabel("AI 驱动的 PPT 共创编辑器，支持智能内容生成、图表转换、风格优化", self.tab_ppt_assistant)
        ppt_desc.setStyleSheet("color: #aaa; font-size: 12px; background: transparent; border: none;")
        ppt_desc.setWordWrap(True)
        ppt_layout.addWidget(ppt_desc)
        
        # 功能按钮区
        ppt_buttons_layout = QVBoxLayout()
        ppt_buttons_layout.setSpacing(10)
        
        # 从文本创建 PPT
        self.btn_create_ppt = QPushButton("📝 从文本创建 PPT", self.tab_ppt_assistant)
        self.btn_create_ppt.setStyleSheet("""
            QPushButton {
                background-color: rgba(77, 166, 255, 180);
                color: #fff;
                border: 1px solid rgba(77, 166, 255, 255);
                border-radius: 8px;
                padding: 12px 20px;
                font-weight: bold;
                font-size: 14px;
            }
            QPushButton:hover {
                background-color: rgba(77, 166, 255, 255);
            }
        """)
        self.btn_create_ppt.setCursor(Qt.CursorShape.PointingHandCursor)
        self.btn_create_ppt.clicked.connect(self._launch_ppt_cocreation)
        ppt_buttons_layout.addWidget(self.btn_create_ppt)
        
        # 从模板创建
        self.btn_template_ppt = QPushButton("🎨 从模板创建", self.tab_ppt_assistant)
        self.btn_template_ppt.setStyleSheet("""
            QPushButton {
                background-color: rgba(120, 80, 220, 180);
                color: #fff;
                border: 1px solid rgba(120, 80, 220, 255);
                border-radius: 8px;
                padding: 12px 20px;
                font-weight: bold;
                font-size: 14px;
            }
            QPushButton:hover {
                background-color: rgba(120, 80, 220, 255);
            }
        """)
        self.btn_template_ppt.setCursor(Qt.CursorShape.PointingHandCursor)
        self.btn_template_ppt.clicked.connect(self._launch_ppt_template)
        ppt_buttons_layout.addWidget(self.btn_template_ppt)
        
        # 打开现有 PPT
        self.btn_open_ppt = QPushButton("📂 打开现有 PPT", self.tab_ppt_assistant)
        self.btn_open_ppt.setStyleSheet("""
            QPushButton {
                background-color: rgba(40, 167, 69, 180);
                color: #fff;
                border: 1px solid rgba(40, 167, 69, 255);
                border-radius: 8px;
                padding: 12px 20px;
                font-weight: bold;
                font-size: 14px;
            }
            QPushButton:hover {
                background-color: rgba(40, 167, 69, 255);
            }
        """)
        self.btn_open_ppt.setCursor(Qt.CursorShape.PointingHandCursor)
        self.btn_open_ppt.clicked.connect(self._open_existing_ppt)
        ppt_buttons_layout.addWidget(self.btn_open_ppt)
        
        ppt_layout.addLayout(ppt_buttons_layout)
        
        # 最近使用
        recent_label = QLabel("📋 最近使用", self.tab_ppt_assistant)
        recent_label.setStyleSheet("color: #aaa; font-size: 12px; background: transparent; border: none; margin-top: 10px;")
        ppt_layout.addWidget(recent_label)
        
        self.recent_ppt_list = QListWidget(self.tab_ppt_assistant)
        self.recent_ppt_list.setStyleSheet("""
            QListWidget {
                background-color: rgba(30, 30, 35, 200);
                color: #eee;
                border: 1px solid rgba(100, 100, 100, 150);
                border-radius: 6px;
                padding: 4px;
            }
            QListWidget::item {
                padding: 6px;
                border-radius: 4px;
            }
            QListWidget::item:hover {
                background-color: rgba(77, 166, 255, 80);
            }
        """)
        ppt_layout.addWidget(self.recent_ppt_list)
        
        ppt_layout.addStretch()
        
        self.tabs.addTab(self.tab_ppt_assistant, "🎯 PPT 助手")
        
        # ==========================
        # Tab 4: 技能中心
        # ==========================
        self.tab_skill_center = QWidget()
        skill_center_layout = QVBoxLayout(self.tab_skill_center)
        skill_center_layout.setContentsMargins(0, 10, 0, 0)
        
        # 技能面板
        self.skill_panel = SkillPanel(self.skill_registry)
        self.skill_panel.skill_execute.connect(self._on_skill_execute)
        skill_center_layout.addWidget(self.skill_panel)
        
        self.tabs.addTab(self.tab_skill_center, "⚡ 技能中心")
        
        self.tabs.setFocusPolicy(Qt.FocusPolicy.NoFocus)

        # 确保切换 Tab 时重置焦点
        self.tabs.currentChanged.connect(self._on_tab_changed)
        
        # 初始化技能搜索对话框
        self.skill_search_dialog = SkillSearchDialog(self.skill_registry)
        self.skill_search_dialog.skill_execute.connect(self._on_skill_execute)
        
        # 初始化技能命令解析器
        self.skill_command_parser = SkillCommandParser(self.skill_registry)
        
        # 设置快捷键
        self._setup_shortcuts()

    def _setup_shortcuts(self):
        """设置快捷键"""
        from PyQt6.QtGui import QShortcut, QKeySequence
        
        # Ctrl+K 打开技能搜索
        search_shortcut = QShortcut(QKeySequence("Ctrl+K"), self)
        search_shortcut.activated.connect(self._show_skill_search)
        
        # Ctrl+/ 打开技能命令输入
        command_shortcut = QShortcut(QKeySequence("Ctrl+/"), self)
        command_shortcut.activated.connect(self._show_skill_command)
        
        # Ctrl+Shift+S 切换到技能中心
        skill_center_shortcut = QShortcut(QKeySequence("Ctrl+Shift+S"), self)
        skill_center_shortcut.activated.connect(lambda: self.tabs.setCurrentWidget(self.tab_skill_center))
    
    def _show_skill_search(self):
        """显示技能搜索对话框"""
        self.skill_search_dialog.show_dialog()
    
    def _show_skill_command(self):
        """显示技能命令输入"""
        # 切换到对话Tab并清空输入框
        self.tabs.setCurrentIndex(1)
        self.chat_input.clear()
        self.chat_input.setPlaceholderText("输入 /技能名称 执行技能...")
        self.chat_input.setFocus()

    def jump_to_chat(self):
        # 切换到聊天 Tab，将源文本上下文带入对话
        self.chat_display.clear()
        context_info = ""
        if self.current_text:
            preview = self.current_text[:150] + ("…" if len(self.current_text) > 150 else "")
            context_info = f"\n\n📄 当前上下文（{len(self.current_text)} 字符，来源: {self.context_source}）:\n{preview}"
        self.append_chat_message("系统", f"已将上下文带入，您可以继续追问。{context_info}")
        self.tabs.setCurrentIndex(1)
        self.chat_input.setFocus()

    def _on_tab_changed(self, index):
        # 强制恢复默认光标，防止拖拽或耗时操作遗留的光标状态
        QApplication.restoreOverrideCursor()
        self.setCursor(Qt.CursorShape.ArrowCursor)
        
        if index == 1:  # 切换到对话 Tab
            self.chat_input.setFocus()
        # 注意：非聊天 Tab 不调用 self.setFocus()，
        # 否则会与 WA_ShowWithoutActivating 冲突，导致 macOS 收回窗口
    
    def _on_skill_execute(self, skill_name: str, params: Dict[str, Any]):
        """执行技能"""
        try:
            # 获取技能
            skill = self.skill_registry.get_skill(skill_name)
            if not skill:
                print(f"[技能执行失败] 未找到技能: {skill_name}")
                return
            
            # 创建执行上下文
            intent = params.get("intent", "")
            input_data = params.get("input_data", {})
            
            # 添加当前上下文
            if self.current_text:
                input_data["selected_text"] = self.current_text
            if self.context_source:
                input_data["context_source"] = self.context_source
            
            context = SkillContext(
                intent=intent,
                input_data=input_data
            )
            
            # 执行技能
            print(f"[技能执行] {skill_name}, 意图: {intent}")
            
            # 切换到对话 Tab 显示结果
            self.tabs.setCurrentIndex(1)
            self.chat_input.setFocus()
            
            # 在对话中显示执行信息
            self.append_chat_message("系统", f"正在执行技能: {skill.metadata.display_name}...")
            
            # 异步执行技能
            import asyncio
            
            async def execute_skill():
                try:
                    result = await skill.execute(context)
                    return result
                except Exception as e:
                    print(f"[技能执行异常] {e}")
                    return None
            
            # 在新线程中执行异步任务
            def run_async():
                loop = asyncio.new_event_loop()
                asyncio.set_event_loop(loop)
                result = loop.run_until_complete(execute_skill())
                loop.close()
                
                # 在主线程中更新UI
                from PyQt6.QtCore import QMetaObject, Qt
                
                if result and result.success:
                    QMetaObject.invokeMethod(
                        self, 
                        "_on_skill_result",
                        Qt.ConnectionType.QueuedConnection,
                        skill_name,
                        result.output
                    )
                else:
                    QMetaObject.invokeMethod(
                        self,
                        "_on_skill_error", 
                        Qt.ConnectionType.QueuedConnection,
                        skill_name,
                        str(result.error) if result else "执行失败"
                    )
            
            # 启动执行线程
            import threading
            thread = threading.Thread(target=run_async)
            thread.daemon = True
            thread.start()
            
        except Exception as e:
            print(f"[技能执行失败] {e}")
            import traceback
            traceback.print_exc()
    
    def _on_skill_result(self, skill_name: str, result: str):
        """技能执行结果回调"""
        self.append_chat_message("AI", f"技能 {skill_name} 执行结果:\n{result}")
    
    def _on_skill_error(self, skill_name: str, error: str):
        """技能执行错误回调"""
        self.append_chat_message("系统", f"技能 {skill_name} 执行失败: {error}")
    
    def _launch_ppt_cocreation(self):
        """启动 PPT 共创编辑器"""
        try:
            from opencopilot.capabilities.ppt import CoCreationDialog
            from PyQt6.QtWidgets import QInputDialog, QTextEdit, QDialogButtonBox
            
            # 获取当前上下文文本
            text = self.current_text or ""
            
            if not text:
                # 弹出输入对话框让用户输入内容
                dialog = QDialog(self)
                dialog.setWindowTitle("📝 输入 PPT 内容")
                dialog.setMinimumSize(500, 400)
                
                layout = QVBoxLayout(dialog)
                
                label = QLabel("请输入或粘贴 PPT 内容，AI 将自动生成大纲：")
                label.setStyleSheet("color: #aaa; font-size: 13px;")
                layout.addWidget(label)
                
                text_edit = QTextEdit()
                text_edit.setPlaceholderText("例如：\n\n一、项目背景\n- 市场需求分析\n- 竞争格局\n\n二、技术方案\n- 架构设计\n- 核心技术\n\n三、实施计划\n- 时间节点\n- 里程碑")
                text_edit.setStyleSheet("""
                    QTextEdit {
                        background-color: #2d2d2d;
                        color: #eee;
                        border: 1px solid #555;
                        border-radius: 6px;
                        padding: 8px;
                        font-size: 13px;
                    }
                """)
                layout.addWidget(text_edit)
                
                # 按钮
                btn_box = QDialogButtonBox(QDialogButtonBox.StandardButton.Ok | QDialogButtonBox.StandardButton.Cancel)
                btn_box.accepted.connect(dialog.accept)
                btn_box.rejected.connect(dialog.reject)
                layout.addWidget(btn_box)
                
                if dialog.exec() != QDialog.DialogCode.Accepted:
                    return
                
                text = text_edit.toPlainText().strip()
                if not text:
                    QMessageBox.warning(self, "提示", "请输入内容后再创建 PPT")
                    return
            
            # 显示加载提示
            loading_msg = QMessageBox(self)
            loading_msg.setWindowTitle("AI 生成中")
            loading_msg.setText("正在分析内容，生成 PPT 大纲...")
            loading_msg.setStandardButtons(QMessageBox.StandardButton.NoButton)
            loading_msg.show()
            QApplication.processEvents()
            
            # 调用 AI 生成初始大纲
            json_data = self._generate_ppt_outline_with_ai(text)
            
            loading_msg.close()
            
            if not json_data:
                # AI 生成失败，使用智能兜底模板
                print("[PPT] AI 生成失败，使用兜底模板")
                json_data = self._generate_fallback_outline(text)
            
            dialog = CoCreationDialog(
                original_text=text,
                json_data=json_data,
                agent_url="http://127.0.0.1:18888",
                parent=self
            )
            
            if dialog.exec() == QDialog.DialogCode.Accepted:
                output_path = dialog.get_output_path()
                if output_path:
                    QMessageBox.information(self, "导出成功", f"PPT 已成功导出至：\n{output_path}")
        except ImportError as e:
            QMessageBox.warning(self, "加载失败", f"无法加载 PPT 共创编辑器：{e}")
        except Exception as e:
            QMessageBox.warning(self, "错误", f"启动 PPT 助手时出错：{e}")
    
    def _generate_ppt_outline_with_ai(self, text: str) -> dict:
        """使用 AI 生成 PPT 大纲
        
        使用API已有的 ppt_generator context_source，避免prompt冲突
        
        Returns:
            dict: {"title": "PPT标题", "slides": [...]} 或 None（生成失败）
        """
        try:
            import requests
            import re
            
            print(f"[PPT] 开始生成大纲，内容长度: {len(text)} 字符")
            print(f"[PPT] 内容前100字: {text[:100]}")
            
            # 使用API已有的 ppt_generator context_source
            # 通过 ASUCustomAgentClient 统一调用，确保 web_search 等参数一致传递
            from llm_provider import ASUCustomAgentClient
            agent_client = ASUCustomAgentClient()
            full_text = ""
            for chunk in agent_client.stream_agent_task(
                text[:3000],
                action_type="chat",
                session_id=f"ppt_gen_{int(time.time())}",
                is_new_task=True,
                context_source="ppt_generator",
            ):
                if isinstance(chunk, tuple):
                    continue  # 跳过 annotations 等 metadata
                full_text += chunk
            
            print(f"[PPT] AI 返回内容长度: {len(full_text)} 字符")
            print(f"[PPT] AI 返回前300字: {full_text[:300]}")
            
            # 清理模型输出中的思考过程和代码块标记
            # 移除 <think>...</think> 标签
            cleaned = re.sub(r'<think>.*?</think>', '', full_text, flags=re.DOTALL)
            # 移除未闭合的 <think> 标签
            if '<think>' in cleaned:
                cleaned = cleaned.split('<think>')[0]
            # 移除 ```json ... ``` 标记
            cleaned = re.sub(r'```(?:json)?\s*', '', cleaned)
            cleaned = cleaned.replace('```', '')
            
            # 提取 JSON
            from ppt_generator import extract_json_from_text
            json_data = extract_json_from_text(cleaned)
            
            if json_data:
                # 兼容两种返回格式：数组或字典
                if isinstance(json_data, list):
                    result = {"title": "演示文稿", "slides": json_data}
                elif isinstance(json_data, dict) and "slides" in json_data:
                    result = json_data
                else:
                    print(f"[PPT] JSON 格式不符合预期: {type(json_data)}")
                    return None
                
                print(f"[PPT] 解析成功，共 {len(result.get('slides', []))} 页幻灯片")
                return result
            else:
                print("[PPT] extract_json_from_text 返回空")
                print(f"[PPT] 清理后输出: {cleaned[:500]}")
            
            return None
        except Exception as e:
            print(f"[PPT] AI 生成大纲失败: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def _generate_fallback_outline(self, text: str) -> dict:
        """AI生成失败时的智能兜底模板
        
        根据文本内容自动提取要点，生成合理的PPT大纲
        """
        import re
        
        # 提取文本中的要点
        lines = text.split('\n')
        sections = []  # (标题, 要点列表)
        current_title = ""
        current_items = []
        
        for line in lines:
            line = line.strip()
            if not line:
                continue
            
            # 检测标题行（#、##、数字编号、中文序号）
            is_header = False
            title_text = ""
            
            if line.startswith('#'):
                # Markdown 标题
                title_text = re.sub(r'^#+\s*', '', line).strip()
                is_header = True
            elif re.match(r'^[一二三四五六七八九十]+[、.]', line):
                # 中文序号
                title_text = line
                is_header = True
            elif re.match(r'^\d+[、.\s]', line):
                # 数字序号
                title_text = line
                is_header = True
            elif re.match(r'^【.+】', line):
                # 中文括号标题
                title_text = line
                is_header = True
            
            if is_header and title_text:
                if current_title and current_items:
                    sections.append((current_title, current_items[:5]))  # 最多5个要点
                current_title = title_text
                current_items = []
            elif line.startswith('- ') or line.startswith('* ') or line.startswith('• '):
                # 列表项
                item_text = line[2:].strip()
                if item_text:
                    current_items.append(item_text)
            elif len(line) > 10 and len(line) < 100:
                # 普通文本行（长度适中）
                current_items.append(line)
        
        # 添加最后一个section
        if current_title and current_items:
            sections.append((current_title, current_items[:5]))
        
        # 如果没有提取到section，用简单分割
        if not sections:
            chunk_size = 200
            chunks = [text[i:i+chunk_size] for i in range(0, min(len(text), 1000), chunk_size)]
            for i, chunk in enumerate(chunks[:4]):
                sections.append((f"内容要点 {i+1}", [chunk[:80]]))
        
        # 构建幻灯片
        slides = []
        
        # 封面页
        title_text = sections[0][0] if sections else "演示文稿"
        slides.append({
            "type": "title",
            "layout": "center",
            "title": title_text,
            "subtitle": "内容概览"
        })
        
        # 内容页（最多8页）
        for section_title, items in sections[:8]:
            if not items:
                continue
            slide_items = [{"level": 0, "text": item[:60]} for item in items[:5]]
            slides.append({
                "type": "content",
                "layout": "text_only",
                "title": section_title[:30],
                "items": slide_items
            })
        
        # 总结页
        slides.append({
            "type": "content",
            "layout": "center",
            "title": "总结",
            "items": [
                {"level": 0, "text": "以上为主要内容概览"},
                {"level": 0, "text": "可在编辑器中进一步调整"}
            ]
        })
        
        print(f"[PPT] 兜底模板生成完成，共 {len(slides)} 页")
        return {"title": title_text, "slides": slides}
    
    def _launch_ppt_template(self):
        """从模板创建 PPT"""
        # TODO: 实现模板选择功能
        QMessageBox.information(self, "敬请期待", "模板功能正在开发中...")
    
    def _open_existing_ppt(self):
        """打开现有 PPT 文件"""
        file_path, _ = QFileDialog.getOpenFileName(
            self, "打开 PPT 文件", 
            os.path.expanduser("~/Desktop"),
            "PowerPoint Files (*.pptx *.ppt)"
        )
        if file_path:
            # TODO: 实现加载现有 PPT 的功能
            QMessageBox.information(self, "功能开发中", "加载现有 PPT 功能正在开发中...")

    def _get_ide_port(self):
        """从临时文件读取当前激活的 IDE 插件端口"""
        port_file = os.path.join(tempfile.gettempdir(), 'asu_ide_port.txt')
        if os.path.exists(port_file):
            try:
                with open(port_file, 'r') as f:
                    port = f.read().strip()
                    if port.isdigit():
                        return port
            except Exception:
                print(f"[ASU] 读取 IDE 端口文件失败: {traceback.format_exc()}")
        return None

    def _on_custom_instruction(self):
        """用户在指令输入栏按 Enter 或点击 ▶ 按钮时触发。"""
        instruction = self.instruction_input.text().strip()
        if not instruction:
            return
            
        # 如果是视觉分析模式且带有截图缓存
        if self.context_source == "vision" and hasattr(self, 'ai_card_image_base64') and self.ai_card_image_base64:
            self.trigger_ai("custom", custom_instruction=instruction, image_base64=self.ai_card_image_base64)
            # 发送后清除图片缓存
            self.ai_card_image_base64 = None
            return

        if not self.current_text:
            # 修改点：如果用户没有提供任何文本上下文，直接将其作为自由对话处理
            self.instruction_input.clear()
            self.tabs.setCurrentIndex(1)  # 自动切换到对话 Tab
            self.chat_input.setText(instruction)
            self.send_chat_message()
            return
            
        # trigger_ai 会自动读取 instruction_input 中的内容
        self.trigger_ai("custom")

    def _apply_to_ide(self):
        """将 AI 回复结果回写到 IDE 当前文件。
        
        策略：
        - 有选区范围(_ide_selection_range)：只替换选中片段（局部回写）
        - 无选区但有全文(_ide_full_document)：在全文中找到原始片段并替换
        - 都没有：全文替换（兜底）
        """
        port = self._get_ide_port()
        if not port:
            self.text_edit.setPlainText("❌ 无法连接 IDE 插件，请确认已在 VSCode/Trae 中安装并激活插件。")
            return

        # 获取 AI 回复的纯文本内容
        result_text = self.text_edit.toPlainText()
        if not result_text.strip():
            return

        try:
            # 策略 1: 有选区范围 → 精确局部替换
            if hasattr(self, '_ide_selection_range') and self._ide_selection_range:
                payload = {
                    "replace": result_text,
                    "range": self._ide_selection_range
                }
                mode_label = "局部替换(选区)"
            # 策略 2: 无选区但有全文和原始文本 → 在全文中定位并替换片段
            elif hasattr(self, '_ide_full_document') and self._ide_full_document and self.current_text:
                full_doc = self._ide_full_document
                original = self.current_text
                # 在全文中查找原始片段的位置
                idx = full_doc.find(original)
                if idx >= 0:
                    # 计算行号和列号
                    before = full_doc[:idx]
                    start_line = before.count('\n')
                    start_col = len(before) - before.rfind('\n') - 1 if '\n' in before else len(before)
                    
                    after_original = full_doc[:idx + len(original)]
                    end_line = after_original.count('\n')
                    end_col = len(after_original) - after_original.rfind('\n') - 1 if '\n' in after_original else len(after_original)
                    
                    payload = {
                        "replace": result_text,
                        "range": {
                            "startLine": start_line,
                            "startCol": start_col,
                            "endLine": end_line,
                            "endCol": end_col
                        }
                    }
                    mode_label = "局部替换(定位)"
                else:
                    # 在全文中找不到原文片段 → 全文替换
                    payload = {"content": result_text}
                    mode_label = "全文替换(未定位到原文)"
            # 策略 3: 兜底全文替换
            else:
                payload = {"content": result_text}
                mode_label = "全文替换"

            response = httpx.post(
                f"http://127.0.0.1:{port}/apply",
                json=payload, timeout=3.0
            )
            if response.status_code == 200:
                data = response.json()
                mode = data.get("mode", "unknown")
                self.btn_apply_to_ide.setText(f"✅ 已回写({mode_label})")
                self.btn_apply_to_ide.setStyleSheet(self.btn_apply_to_ide.styleSheet().replace("rgba(40, 167, 69, 180)", "rgba(100, 100, 110, 180)").replace("rgba(40, 167, 69, 255)", "rgba(100, 100, 110, 255)"))
                self.btn_apply_to_ide.setEnabled(False)
                print(f"[ASU] 回写IDE成功 | mode={mode} | strategy={mode_label}")
            else:
                self.text_edit.setPlainText(f"❌ 回写失败: HTTP {response.status_code}")
                print(f"[ASU] 回写IDE失败 | status={response.status_code}")
        except Exception as e:
            self.text_edit.setPlainText(f"❌ 回写失败: {str(e)}\n\n请确认 IDE 插件正在运行。")
            print(f"[ASU] 回写IDE异常 | error={str(e)}")

    def _copy_result_to_clipboard(self):
        """将 AI 回复结果复制到系统剪贴板。"""
        result_text = self.text_edit.toPlainText()
        if result_text.strip():
            clipboard = QApplication.clipboard()
            clipboard.setText(result_text)
            self.btn_copy_result.setText("✅ 已复制")
            QTimer.singleShot(2000, lambda: self.btn_copy_result.setText("📋 复制结果"))

    def _export_to_ppt(self):
        """将 AI 生成的 Markdown 大纲，通过人机共创界面，最终导出为 PPT"""
        text = self.text_edit.toPlainText()
        if not text:
            return
            
        json_data = extract_json_from_text(text)
        if not json_data:
            QMessageBox.warning(self, "解析失败", "大模型输出的内容不符合预期的 JSON 大纲格式。")
            return
            
        # 唤起人机共创编辑器
        dialog = PPTPreviewDialog(json_data, self)
        if dialog.exec() == QDialog.DialogCode.Accepted:
            final_json = dialog.get_final_json()
            try:
                # 调用底层排版引擎生成 PPT
                ppt_path = generate_ppt_from_json(final_json)
                QMessageBox.information(self, "导出成功", f"PPT 已成功导出至：\n{ppt_path}")
                # macOS 自动打开文件
                subprocess.run(["open", ppt_path])
            except Exception as e:
                import traceback
                print(traceback.format_exc())
                QMessageBox.critical(self, "导出失败", f"生成 PPT 时发生错误：\n{e}")

    def _try_get_ide_selection(self):
        """尝试自动读取 IDE 选区范围（用于回写时的局部替换定位）。"""
        port = self._get_ide_port()
        if not port:
            return
        try:
            response = httpx.get(f"http://127.0.0.1:{port}/selection", timeout=1.0)
            if response.status_code == 200:
                data = response.json()
                sel_range = data.get("range", None)
                sel_text = data.get("text", "")
                if sel_range and sel_text:
                    self._ide_selection_range = sel_range
                    # 如果还没有全文缓存，顺便读取
                    if not (hasattr(self, '_ide_full_document') and self._ide_full_document):
                        ctx_resp = httpx.get(f"http://127.0.0.1:{port}/context", timeout=1.0)
                        if ctx_resp.status_code == 200:
                            self._ide_full_document = ctx_resp.json().get("content", "")
        except Exception:
            pass  # 静默失败，回写时会降级到全文替换

    def paste_from_clipboard(self):
        """从系统剪贴板粘贴文本，等待用户输入修改要求或点击快捷指令。"""
        clipboard = QApplication.clipboard()
        text = clipboard.text()
        if not text or not text.strip():
            self.text_edit.setPlainText("❌ 剪贴板为空，请先在 IDE 中复制文本（Cmd+C）。")
            return
        self.current_text = text
        self.context_source = "clipboard"
        self.context_meta = {}
        self.tabs.setCurrentIndex(0)
        # 不自动触发 AI，展示内容后等待用户指令
        preview = text[:200] + ("…" if len(text) > 200 else "")
        self.text_edit.setPlainText(
            f"📄 已获取内容（{len(text)} 字符）:\n\n{preview}\n\n"
            f"👇 请输入修改要求后按 Enter，或点击下方快捷指令。"
        )
        self.instruction_input.setFocus()

    def _on_vision_analyze_clicked(self):
        """处理视觉分析按钮点击：通过 Broker 获取截图并设置状态"""
        self.btn_vision_analyze.setEnabled(False)
        self.btn_vision_analyze.setText("⏳ 正在截取...")
        
        # 使用 QThread 防止网络请求阻塞 UI
        class VisionWorker(QThread):
            finished = pyqtSignal(str, str) # status, base64_or_err
            
            def run(self):
                try:
                    probe = SystemProbeClient()
                    if not probe.is_broker_alive():
                        self.finished.emit("error", "无法连接到 Privileged Broker，请确保它已启动。")
                        return
                    
                    b64_img = probe.get_front_window_screenshot()
                    if b64_img:
                        self.finished.emit("success", b64_img)
                    else:
                        self.finished.emit("error", "获取到的截图数据为空。")
                except Exception as e:
                    self.finished.emit("error", str(e))
                    
        self._vision_worker = VisionWorker()
        self._vision_worker.finished.connect(self._on_vision_capture_result)
        self._vision_worker.start()
        
    def _on_vision_capture_result(self, status, result):
        self.btn_vision_analyze.setEnabled(True)
        self.btn_vision_analyze.setText("👁️ 视觉分析前台")
        
        if status == "success":
            self.context_source = "vision"
            self.context_meta = {"has_image": True}
            # 将 base64 存入特殊的变量中供 _on_send_clicked 提取
            self.ai_card_image_base64 = result
            
            # 显示缩略提示
            self.text_edit.setPlainText(
                "📸 已成功捕获前台窗口截图！\n\n"
                "图像已作为多模态输入准备就绪。\n"
                "请在下方输入框中输入您对该图像的分析指令（例如：'提取图中的表格数据'，'解释图中的架构逻辑'等），然后按 Enter 发送。"
            )
            self.instruction_input.setFocus()
        else:
            QMessageBox.warning(self, "视觉分析失败", f"截图提取失败：\n{result}")

    def read_from_ide_extension(self):
        port = self._get_ide_port()
        if not port:
            self.text_edit.setPlainText("❌ 无法找到 IDE 插件端口信息，请确认已在当前 VSCode/Trae 窗口中激活了插件。")
            return
            
        try:
            # 先尝试读取选中文本（如果有选区则只处理选中部分）
            sel_response = httpx.get(f"http://127.0.0.1:{port}/selection", timeout=2.0)
            sel_text = ""
            sel_range = None
            if sel_response.status_code == 200:
                sel_data = sel_response.json()
                sel_text = sel_data.get("text", "")
                sel_range = sel_data.get("range", None)

            # 读取诊断、diff、symbol等高级上下文
            diagnostics_data = []
            git_diff_data = ""
            symbol_data = None
            try:
                diag_resp = httpx.get(f"http://127.0.0.1:{port}/diagnostics", timeout=1.0)
                if diag_resp.status_code == 200:
                    diagnostics_data = diag_resp.json().get("diagnostics", [])
            except Exception: pass
            try:
                diff_resp = httpx.get(f"http://127.0.0.1:{port}/git-diff", timeout=1.0)
                if diff_resp.status_code == 200:
                    git_diff_data = diff_resp.json().get("diff", "")
            except Exception: pass
            try:
                sym_resp = httpx.get(f"http://127.0.0.1:{port}/symbol", timeout=1.0)
                if sym_resp.status_code == 200:
                    symbol_data = sym_resp.json() # symbol_data is the whole object with name, kind, text, range
                    if "error" in symbol_data or ("symbol" in symbol_data and symbol_data["symbol"] is None):
                        symbol_data = None
            except Exception: pass

            # 读取全文
            response = httpx.get(f"http://127.0.0.1:{port}/context", timeout=2.0)
            if response.status_code == 200:
                data = response.json()
                content = data.get("content", "")
                filename = data.get("fileName", "Unknown")
                
                meta_dict = {
                    "file_name": filename,
                    "language": data.get("languageId", ""),
                    "diagnostics": diagnostics_data,
                    "git_diff": git_diff_data,
                    "symbol": symbol_data
                }
                
                if content:
                    # 如果有选中文本，优先使用选中文本作为操作目标
                    if sel_text and sel_text.strip():
                        self.current_text = sel_text
                        self.context_source = "ide"
                        self._ide_selection_range = sel_range  # 记录选区范围，用于局部回写
                        self._ide_full_document = content       # 记录全文，用于上下文
                        self.context_meta = meta_dict
                        self.text_edit.clear()
                        preview = sel_text[:200] + ("…" if len(sel_text) > 200 else "")
                        
                        diag_count = len(diagnostics_data)
                        diag_msg = f" | 发现 {diag_count} 个诊断报错 🎯" if diag_count > 0 else ""
                        
                        self.text_edit.setPlainText(
                            f"✅ 已读取 IDE 选中文本 [{filename}]{diag_msg}\n\n"
                            f"选中: {len(sel_text)} 字符 / 全文: {len(content)} 字符\n\n"
                            f"📄 预览:\n{preview}\n\n"
                            f"👇 请输入修改要求后按 Enter，或点击下方快捷指令。\n"
                            f"💡 回写时将只替换选中的文本片段。"
                        )
                    else:
                        # 没有选区，判断是否有光标下的 symbol
                        if symbol_data and symbol_data.get("text"):
                            self.current_text = symbol_data["text"]
                            self.context_source = "ide"
                            self._ide_selection_range = symbol_data["range"]
                            self._ide_full_document = content
                            self.context_meta = meta_dict
                            self.text_edit.clear()
                            
                            sym_name = symbol_data.get("name", "Unknown")
                            sym_text = symbol_data.get("text", "")
                            preview = sym_text[:200] + ("…" if len(sym_text) > 200 else "")
                            
                            diag_count = len(diagnostics_data)
                            diag_msg = f" | 发现 {diag_count} 个诊断报错 🎯" if diag_count > 0 else ""
                            
                            self.text_edit.setPlainText(
                                f"✅ 已智能截取光标所在代码块 [{sym_name}]{diag_msg}\n\n"
                                f"块大小: {len(sym_text)} 字符 / 全文: {len(content)} 字符\n\n"
                                f"📄 预览:\n{preview}\n\n"
                                f"👇 请输入修改要求后按 Enter，或点击下方快捷指令。\n"
                                f"💡 回写时将只替换该代码块片段。"
                            )
                        else:
                            self.current_text = content
                            self.context_source = "ide"
                            self._ide_selection_range = None
                            self._ide_full_document = content
                            self.context_meta = meta_dict
                            self.text_edit.clear()
                            
                            diag_count = len(diagnostics_data)
                            diag_msg = f" | 发现 {diag_count} 个诊断报错 🎯" if diag_count > 0 else ""
                            
                            self.text_edit.setPlainText(f"✅ 已成功从 IDE 插件读取全文 [{filename}]{diag_msg}\n\n文件大小: {len(content)} 字符\n\n请点击下方快捷指令进行分析。")
                    self.tabs.setCurrentIndex(0)
                    self.btn_read_ide.setText("✅ 已读取")
                    self.btn_read_ide.setStyleSheet(self.btn_read_ide.styleSheet().replace("rgba(77, 166, 255, 180)", "rgba(40, 167, 69, 180)").replace("rgba(77, 166, 255, 255)", "rgba(40, 167, 69, 255)"))
                    self.instruction_input.setFocus()
                else:
                    self.text_edit.setPlainText("❌ 从 IDE 读取的文件内容为空")
            else:
                self.text_edit.setPlainText(f"❌ 读取失败，插件返回状态码: {response.status_code}")
        except Exception as e:
            self.text_edit.setPlainText(f"❌ 无法连接到 IDE 伴生插件，请确认已在 VSCode/Trae 中安装并激活插件。\n\n错误信息: {e}")

    def send_chat_message(self):
        user_text = self.chat_input.text().strip()
        if not user_text:
            return
        
        # 检查是否是技能命令（以/开头）
        if user_text.startswith('/'):
            self._handle_skill_command(user_text)
            self.chat_input.clear()
            return
            
        self.chat_input.clear()
        self.append_chat_message("你", user_text)
        
        self.append_chat_message("AI", "正在思考...", is_temp=True)
        
        if self.chat_worker and self.chat_worker.isRunning():
            self.chat_worker.stop()
            self.chat_worker.wait()

        # 合并工作台任务上下文 + 源文本内容
        meta = dict(self.context_meta)
        if self.task_context:
            meta["task"] = self.task_context
        # 将源文本内容传入 context_meta，让 Agent 知道用户在讨论什么
        if self.current_text:
            meta["source_text"] = self.current_text
            meta["source_type"] = self.context_source

        # 核心修复：把底层的系统状态（当前焦点、最近打开的应用等）注入到对话的上下文元数据中
        meta["current_active_app"] = self.current_active_app
        meta["recent_apps"] = self.recent_apps

        print(f"[ASU] Chat请求 | session={self.session_id[:8]}... | text_len={len(user_text)} | has_source={bool(self.current_text)} | meta_keys={list(meta.keys())}")

        self.chat_worker = ChatWorker(self.provider, user_text, self.session_id,
                                      self.context_source, meta)
        self.chat_worker.text_updated.connect(self.on_chat_updated)
        self.chat_worker.finished_signal.connect(self.on_chat_finished)
        self.chat_worker.start()
    
    def _handle_skill_command(self, command: str):
        """处理技能命令"""
        try:
            # 解析命令
            result = self.skill_command_parser.parse(command)
            
            if result:
                skill_name = result["skill_name"]
                intent = result["intent"]
                params = result["params"]
                
                # 显示执行信息
                self.append_chat_message("系统", f"正在执行技能: {result['metadata'].display_name}...")
                
                # 构建执行参数
                execute_params = {
                    "intent": intent,
                    "input_data": params
                }
                
                # 执行技能
                self._on_skill_execute(skill_name, execute_params)
            else:
                # 显示帮助信息
                self._show_skill_help(command)
                
        except Exception as e:
            self.append_chat_message("系统", f"命令解析失败: {e}")
    
    def _show_skill_help(self, command: str):
        """显示技能帮助信息"""
        # 获取所有可用技能
        all_metadata = self.skill_registry.get_all_metadata()
        
        help_text = f"可用技能列表 (输入 /技能名称 执行):\n\n"
        
        for skill_name, metadata in all_metadata.items():
            help_text += f"• /{skill_name} - {metadata.display_name}\n"
            help_text += f"  {metadata.description[:50]}...\n\n"
        
        help_text += "\n示例:\n"
        help_text += "• /coding - 执行编程技能\n"
        help_text += "• /knowledge:query - 查询知识库\n"
        help_text += "• /ppt:generate - 生成PPT\n"
        
        self.append_chat_message("系统", help_text)

    def append_chat_message(self, role, text, is_temp=False):
        color = "#4da6ff" if role == "你" else "#42f554" if role == "AI" else "#aaaaaa"
        
        cursor = self.chat_display.textCursor()
        cursor.movePosition(cursor.MoveOperation.End)
        self.chat_display.setTextCursor(cursor)
        
        if is_temp:
            self._temp_chat_pos = cursor.position()
            self.chat_display.insertHtml(f'<b style="color:{color};">{role}:</b> {text}<br><br>')
        else:
            self.chat_display.insertHtml(f'<b style="color:{color};">{role}:</b> {text}<br><br>')
            
        scrollbar = self.chat_display.verticalScrollBar()
        scrollbar.setValue(scrollbar.maximum())

    def on_chat_updated(self, text):
        cursor = self.chat_display.textCursor()
        cursor.setPosition(self._temp_chat_pos)
        cursor.movePosition(cursor.MoveOperation.End, cursor.MoveMode.KeepAnchor)
        cursor.removeSelectedText()
        cursor.insertHtml(md_render(text))
        scrollbar = self.chat_display.verticalScrollBar()
        scrollbar.setValue(scrollbar.maximum())

    def on_chat_finished(self):
        pass

    # ---- 新增UI组件处理方法 ----

    def open_batch_dialog(self):
        """打开批量处理对话框"""
        dialog = BatchDialog(self)
        dialog.exec()

    def _on_file_dropped(self, file_path, file_info):
        """处理文件拖入（信号发两个参数: file_path, file_info）"""
        import os
        file_name = os.path.basename(file_path)
        file_ext = os.path.splitext(file_name)[1].lower()
        
        print(f"[ASU] 文件拖入: {file_path}, info={file_info}")
        # 直接交给 _handle_file_drop 处理
        self._handle_file_drop(file_path)

    def _show_text_context_menu(self, position):
        """显示文本右键菜单（增强版，支持Skill）"""
        selected_text = self.text_edit.textCursor().selectedText()
        
        # 使用增强版右键菜单
        menu = SkillContextMenu(self.skill_registry, self)
        menu.skill_execute.connect(self._on_skill_execute)
        menu.action_triggered.connect(self._on_context_menu_action)
        
        # 显示菜单
        menu.show_for_text(selected_text, self.text_edit.mapToGlobal(position), {
            "source": "text_edit",
            "context_source": self.context_source
        })
    
    def _on_context_menu_action(self, action_id: str, data):
        """右键菜单动作处理"""
        if action_id == "copy":
            selected_text = self.text_edit.textCursor().selectedText()
            QApplication.clipboard().setText(selected_text)
        elif action_id == "translate":
            self.trigger_ai("translate")
        elif action_id == "polish":
            self.trigger_ai("polish")
        elif action_id == "revise":
            self.trigger_ai("revise")
        elif action_id == "explain":
            self.trigger_ai("explain")
        elif action_id == "summarize":
            self.trigger_ai("summarize")
        elif action_id == "open_terminology":
            self._open_terminology_dialog()
        elif action_id == "open_translation_memory":
            self._open_translation_memory_dialog()
        elif action_id == "show_all_skills":
            # 切换到技能中心Tab
            self.tabs.setCurrentWidget(self.tab_skill_center)
        elif action_id == "open_command":
            # 打开技能搜索对话框
            self.skill_search_dialog.show_dialog()

    def _open_terminology_dialog(self):
        """打开术语库管理对话框"""
        dialog = TerminologyDialog(self)
        dialog.exec()

    def _open_translation_memory_dialog(self):
        """打开翻译记忆对话框"""
        count = len(self.translation_memory.units) if hasattr(self.translation_memory, 'units') else 0
        QMessageBox.information(self, "翻译记忆", f"翻译记忆中共有 {count} 条记录")

    def _apply_settings(self, settings):
        """应用新设置"""
        if 'theme' in settings:
            self._apply_theme(settings['theme'])
        print(f"[ASU] 设置已更新: {settings}")

    def _apply_theme(self, theme_name):
        """应用主题到UI"""
        theme = self.theme_manager.get_theme_config(theme_name)
        if theme:
            # Theme 是 dataclass 对象，直接访问属性
            bg_color = theme.background if hasattr(theme, 'background') else 'rgba(30, 30, 35, 240)'
            self.frame.setStyleSheet(f"""
                QFrame {{
                    background-color: {bg_color};
                    border-radius: 12px;
                    border: 1px solid rgba(100, 100, 100, 100);
                }}
            """)

    def open_settings(self):
        dialog = SettingsDialog(self)
        dialog.config_updated.connect(self.reload_provider)
        dialog.exec()

    def reload_provider(self):
        self.provider = ProviderFactory.create_provider()

    def open_persona_workshop(self):
        from persona_gui import PersonaManagerDialog
        self.persona_dialog = PersonaManagerDialog(self)
        self.persona_dialog.show()

    def set_agent_status(self, is_online: bool):
        """根据 Agent 守护服务的探活结果更新 UI 状态灯和离线横幅。"""
        self._agent_online = is_online
        if is_online:
            self.agent_status_dot.setStyleSheet("color: #4caf50; font-size: 10px; background: transparent; border: none;")
            self.agent_status_dot.setToolTip("ASU 核心守护服务在线")
            self.agent_offline_banner.hide()
        else:
            self.agent_status_dot.setStyleSheet("color: #f44336; font-size: 10px; background: transparent; border: none;")
            self.agent_status_dot.setToolTip("ASU 核心守护服务离线")
            self.agent_offline_banner.show()

    # ---- 拖拽缩放支持 ----
    def _get_resize_edge(self, pos):
        m = self._resize_margin
        x, y = pos.x(), pos.y()
        w, h = self.width(), self.height()
        b, r = y > h - m, x > w - m
        t, l = y < m, x < m
        if b and r: return 'br'
        if b and l: return 'bl'
        if t and r: return 'tr'
        if t and l: return 'tl'
        if b: return 'b'
        if r: return 'r'
        if t: return 't'
        if l: return 'l'
        return None

    _EDGE_CURSORS = {
        'l': Qt.CursorShape.SizeHorCursor, 'r': Qt.CursorShape.SizeHorCursor,
        't': Qt.CursorShape.SizeVerCursor, 'b': Qt.CursorShape.SizeVerCursor,
        'tl': Qt.CursorShape.SizeFDiagCursor, 'br': Qt.CursorShape.SizeFDiagCursor,
        'tr': Qt.CursorShape.SizeBDiagCursor, 'bl': Qt.CursorShape.SizeBDiagCursor,
    }

    def mousePressEvent(self, event):
        if event.button() == Qt.MouseButton.LeftButton:
            # 检查是否在边缘（缩放）
            edge = self._get_resize_edge(event.pos())
            if edge:
                self._resizing = True
                self._resize_edge = edge
                self._resize_start_geo = self.geometry()
                self._resize_start_pos = event.globalPosition().toPoint()
                QApplication.setOverrideCursor(self._EDGE_CURSORS[edge])
                return
            
            # 检查是否在标题栏区域（拖动）
            # 标题栏大约在顶部40像素内
            if event.pos().y() < 40:
                self._drag_pos = event.globalPosition().toPoint() - self.frameGeometry().topLeft()
                event.accept()
                return
        super().mousePressEvent(event)

    def mouseMoveEvent(self, event):
        # 处理窗口拖动
        if self._drag_pos is not None:
            self.move(event.globalPosition().toPoint() - self._drag_pos)
            event.accept()
            return
        
        if self._resizing:
            delta = event.globalPosition().toPoint() - self._resize_start_pos
            g = self._resize_start_geo
            e = self._resize_edge
            x, y, w, h = g.x(), g.y(), g.width(), g.height()
            min_w, min_h = 400, 300

            if 'r' in e: w = max(min_w, g.width() + delta.x())
            if 'l' in e: x = g.x() + delta.x(); w = max(min_w, g.width() - delta.x())
            if 'b' in e: h = max(min_h, g.height() + delta.y())
            if 't' in e: y = g.y() + delta.y(); h = max(min_h, g.height() - delta.y())
            self.setGeometry(x, y, w, h)
            self.frame.resize(w - 20, h - 22)
            return

        edge = self._get_resize_edge(event.pos())
        if edge:
            self.setCursor(self._EDGE_CURSORS.get(edge, Qt.CursorShape.ArrowCursor))
        elif event.pos().y() < 40:
            self.setCursor(Qt.CursorShape.OpenHandCursor)
        else:
            self.setCursor(Qt.CursorShape.ArrowCursor)
        super().mouseMoveEvent(event)

    def mouseReleaseEvent(self, event):
        if self._resizing:
            QApplication.restoreOverrideCursor()
            self._resizing = False
            self._resize_edge = None
        if self._drag_pos is not None:
            self._drag_pos = None
        super().mouseReleaseEvent(event)
    # ---- 缩放支持结束 ----

    def dragEnterEvent(self, event):
        # 拖拽进入时，标记正在拖拽
        self._is_dragging = True
        self._allow_close = False
        if event.mimeData().hasText():
            event.acceptProposedAction()
        elif event.mimeData().hasUrls():
            event.acceptProposedAction()

    def dragMoveEvent(self, event):
        # 拖拽移动时，保持接受状态
        event.acceptProposedAction()

    def dragLeaveEvent(self, event):
        # 拖拽离开时，清除拖拽标记
        self._is_dragging = False

    def dropEvent(self, event):
        # 拖拽放下时，清除拖拽标记
        self._is_dragging = False
        self._allow_close = False
        
        # 优先处理文件拖入
        if event.mimeData().hasUrls():
            urls = event.mimeData().urls()
            if urls:
                file_path = urls[0].toLocalFile()
                if file_path:
                    self._handle_file_drop(file_path)
                    event.acceptProposedAction()
                    return
        
        # 处理文本拖入
        text = event.mimeData().text()
        if text:
            self.current_text = text
            if self.revision_mode:
                # 修订模式：拖入的文本作为待修订目标
                self.context_source = "revision"
                self.context_meta = {"revision_target": text}
                self.tabs.setCurrentIndex(0)
                self.trigger_ai("revision")
            else:
                self.context_source = "drag"
                self.context_meta = {}
                self.tabs.setCurrentIndex(0)
                # 不自动触发，等待用户指定修改方向
                preview = text[:200] + ("…" if len(text) > 200 else "")
                self.text_edit.setPlainText(
                    f"📄 已获取内容（{len(text)} 字符）:\n\n{preview}\n\n"
                    f"👇 请输入修改要求后按 Enter，或点击下方快捷指令。"
                )
                self.instruction_input.setFocus()

    def _handle_file_drop(self, file_path):
        """处理文件拖入"""
        import os
        file_name = os.path.basename(file_path)
        file_ext = os.path.splitext(file_name)[1].lower()
        
        print(f"[ASU] 文件拖入: {file_path}")
        
        # 显示进度条
        self.progress_widget.show()
        self.progress_widget.start(100)
        self.progress_widget.update(10, f"正在读取: {file_name}")
        
        try:
            # 读取文件内容
            content = None
            if file_ext in ['.txt', '.md', '.py', '.js', '.ts', '.json', '.xml', '.html', '.css', '.yaml', '.yml']:
                with open(file_path, 'r', encoding='utf-8') as f:
                    content = f.read()
            elif file_ext == '.docx':
                try:
                    import docx
                    doc = docx.Document(file_path)
                    content = '\n'.join([p.text for p in doc.paragraphs])
                except ImportError:
                    self.text_edit.setPlainText("❌ 需要安装 python-docx 库才能读取 .docx 文件\n\npip install python-docx")
                    self.progress_widget.hide()
                    return
            elif file_ext == '.pptx':
                try:
                    from pptx import Presentation
                    prs = Presentation(file_path)
                    content = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                content.append(shape.text)
                    content = '\n'.join(content)
                except ImportError:
                    self.text_edit.setPlainText("❌ 需要安装 python-pptx 库才能读取 .pptx 文件\n\npip install python-pptx")
                    self.progress_widget.hide()
                    return
            else:
                self.text_edit.setPlainText(f"❌ 不支持的文件格式: {file_ext}\n\n支持的格式: .txt, .md, .py, .js, .ts, .json, .xml, .html, .css, .yaml, .yml, .docx, .pptx")
                self.progress_widget.hide()
                return
            
            if content:
                self.current_text = content
                self.context_source = "file"
                self.context_meta = {
                    "file_name": file_name,
                    "file_path": file_path,
                    "file_ext": file_ext,
                }
                self.tabs.setCurrentIndex(0)
                
                # 显示文件内容预览
                preview = content[:300] + ("…" if len(content) > 300 else "")
                self.text_edit.setPlainText(
                    f"📄 已读取文件: {file_name}\n"
                    f"📊 内容长度: {len(content)} 字符\n\n"
                    f"---\n{preview}\n---\n\n"
                    f"👇 请输入修改要求后按 Enter，或点击下方快捷指令。"
                )
                self.instruction_input.setFocus()
                self.progress_widget.update(100, "文件读取完成!")
                QTimer.singleShot(1000, self.progress_widget.hide)
            else:
                self.text_edit.setPlainText(f"❌ 文件内容为空: {file_name}")
                self.progress_widget.hide()
                
        except Exception as e:
            self.text_edit.setPlainText(f"❌ 读取文件失败: {str(e)}")
            self.progress_widget.hide()
            import traceback
            traceback.print_exc()

    def show_card(self, x, y, selected_text=""):
        self.current_text = selected_text
        self.context_source = "drag" if selected_text else ""
        self.context_meta = {}
        self.session_id = str(uuid.uuid4())
        self._ide_selection_range = None
        self._ide_full_document = ""
        # 重置修订模式（每次唤出卡片都是普通模式）
        if self.revision_mode:
            self.revision_mode = False
            self.full_document = ""
            self.btn_revision.setChecked(False)
            self.btn_revision.setText("📝 全文修订")
        
        self.text_edit.clear()
        if selected_text:
            self.text_edit.setPlainText(f"📌 已自动提取系统选区内容：\n\n{selected_text}")
            self.tabs.setCurrentIndex(0)
        else:
            self.text_edit.setPlainText("✨ Smart Copilot 已就绪。\n\n你可以拖拽/粘贴文本到这里进行快捷处理，\n或者直接在下方输入文字与智能体对话。")
            # 默认切换到连续对话模式，方便直接开聊
            self.tabs.setCurrentIndex(1)
            
        self.instruction_input.clear()
        self.btn_apply_to_ide.hide()
        self.btn_copy_result.hide()
        self.btn_export_ppt.hide()
        
        # 重置 IDE 按钮状态
        self.btn_read_ide.setText("📥 极速读取当前 IDE 全文")
        self.btn_read_ide.setStyleSheet("""
            QPushButton {
                background-color: rgba(77, 166, 255, 180);
                color: #fff;
                border-radius: 8px;
                padding: 6px 12px;
                font-size: 12px;
                font-weight: bold;
                border: 1px solid rgba(77, 166, 255, 255);
            }
            QPushButton:hover {
                background-color: rgba(77, 166, 255, 255);
            }
        """)

        # 重置浏览器按钮状态
        self.btn_read_browser.setText("🌐 一键读取当前网页全文")
        self.btn_read_browser.setStyleSheet("""
            QPushButton {
                background-color: rgba(255, 140, 0, 180);
                color: #fff;
                border-radius: 8px;
                padding: 6px 12px;
                font-size: 12px;
                font-weight: bold;
                border: 1px solid rgba(255, 140, 0, 255);
            }
            QPushButton:hover {
                background-color: rgba(255, 140, 0, 255);
            }
        """)

        # 异步探测
        threading.Thread(target=self._probe_ide_extension, daemon=True).start()
        threading.Thread(target=self._probe_browser, daemon=True).start()

        # 考虑到高DPI，使用 QCursor.pos() 获得准确逻辑坐标
        pos = QCursor.pos()
        
        # 获取当前鼠标所在的屏幕
        screen = QApplication.screenAt(pos)
        if not screen:
            screen = QApplication.primaryScreen()
        screen_rect = screen.geometry()
        
        card_w = self.width()
        card_h = self.height()
        
        # 默认显示在鼠标右下方
        target_x = pos.x() + 15
        target_y = pos.y() + 15
        
        # 边缘碰撞检测：如果右侧超出屏幕，则翻转显示在鼠标左侧
        if target_x + card_w > screen_rect.right():
            target_x = pos.x() - card_w - 15
            
        # 边缘碰撞检测：如果底部超出屏幕，则翻转显示在鼠标上方
        if target_y + card_h > screen_rect.bottom():
            target_y = pos.y() - card_h - 15
            
        # 终极安全边界保护：确保无论如何卡片都不会超出当前屏幕的可视区域
        target_x = max(screen_rect.left(), min(target_x, screen_rect.right() - card_w))
        target_y = max(screen_rect.top(), min(target_y, screen_rect.bottom() - card_h))
        
        self.move(target_x, target_y)
        self.setAttribute(Qt.WidgetAttribute.WA_ShowWithoutActivating, True)
        self.show()
        self.raise_()  # macOS: 确保浮窗在所有窗口最前面
        make_panel_persistent(self)  # 设置 NSPanel 不随失焦隐藏

    def _probe_browser(self):
        try:
            client = SystemProbeClient()
            if not client.is_broker_alive():
                self.browser_probe_result.emit("")
                return
                
            front_app = client.get_frontmost_app()
            supported_browsers = ["Google Chrome", "Safari", "Brave Browser", "Microsoft Edge", "Arc"]
            
            if front_app in supported_browsers:
                self.browser_probe_result.emit(front_app)
            else:
                self.browser_probe_result.emit("")
        except Exception:
            import traceback
            print(f"[ASU] 浏览器探测失败: {traceback.format_exc()}")
            self.browser_probe_result.emit("")

    def _probe_ide_extension(self):
        port = self._get_ide_port()
        if not port:
            self.ide_probe_result.emit(False)
            return
            
        try:
            # 使用 GET 请求进行探活
            response = httpx.get(f"http://127.0.0.1:{port}/context", timeout=0.3)
            # 只要能连上（不管是不是404没打开文件），都说明插件在运行
            if response.status_code in [200, 404]:
                self.ide_probe_result.emit(True)
                return
        except Exception:
            print(f"[ASU] IDE 插件探测失败: {traceback.format_exc()}")
        self.ide_probe_result.emit(False)

    def _update_ide_btn(self, is_active):
        if is_active:
            self.btn_read_ide.show()
        else:
            self.btn_read_ide.hide()

    def _update_browser_btn(self, browser_name):
        if browser_name:
            self.active_browser = browser_name
            self.btn_read_browser.setText(f"🌐 一键读取当前网页 ({browser_name})")
            self.btn_read_browser.show()
        else:
            self.active_browser = None
            self.btn_read_browser.hide()

    def read_from_browser(self):
        if not self.active_browser:
            return

        browser = self.active_browser
        self.btn_read_browser.setEnabled(False)
        self.btn_read_browser.setText("⏳ 读取中...")
        self.text_edit.setPlainText(f"正在从 {browser} 读取网页内容...")
        
        # 强制设置等待光标
        QApplication.setOverrideCursor(Qt.CursorShape.WaitCursor)

        self._browser_worker = BrowserReaderWorker(browser)
        self._browser_worker.finished.connect(self._on_browser_text_ready)
        self._browser_worker.error.connect(self._on_browser_error)
        self._browser_worker.start()

    def _on_browser_text_ready(self, browser, text):
        QApplication.restoreOverrideCursor()
        self.btn_read_browser.setEnabled(True)
        if text:
            self.current_text = text
            self.context_source = "browser"
            self.context_meta = {"app_name": browser}
            self.text_edit.clear()
            self.text_edit.setPlainText(f"✅ 已成功从 {browser} 读取网页全文\n\n网页大小: {len(text)} 字符\n\n请点击下方快捷指令进行分析。")
            self.tabs.setCurrentIndex(0)
            self.btn_read_browser.setText("✅ 已读取全文")
            self.btn_read_browser.setStyleSheet(self.btn_read_browser.styleSheet().replace("rgba(255, 140, 0, 180)", "rgba(40, 167, 69, 180)").replace("rgba(255, 140, 0, 255)", "rgba(40, 167, 69, 255)"))
        else:
            self.text_edit.setPlainText(f"❌ 从 {browser} 读取的内容为空。")
            self.btn_read_browser.setText("🌐 一键读取当前网页全文")

    def _on_browser_error(self, err_msg):
        QApplication.restoreOverrideCursor()
        self.btn_read_browser.setEnabled(True)
        self.text_edit.setPlainText(err_msg)
        self.btn_read_browser.setText("🌐 一键读取当前网页全文")

    def trigger_ai(self, action_type, custom_instruction=None, image_base64=None):
        if not self.current_text and not image_base64:
            return
            
        # 读取自定义指令（如果有）
        if not custom_instruction:
            custom_instruction = self.instruction_input.text().strip()
        if custom_instruction:
            action_type = "custom"

        self.text_edit.clear()
        self.text_edit.setPlainText("正在思考...\n")
        
        # 隐藏回写/复制按钮（等回复完成再显示）
        self.btn_apply_to_ide.hide()
        self.btn_copy_result.hide()
        # 重置回写按钮状态
        self.btn_apply_to_ide.setEnabled(True)
        self.btn_apply_to_ide.setText("📝 回写到 IDE")
        self.btn_apply_to_ide.setStyleSheet("""
            QPushButton {
                background-color: rgba(40, 167, 69, 180);
                color: #fff; border-radius: 8px;
                padding: 5px 10px; font-size: 11px; font-weight: bold;
                border: 1px solid rgba(40, 167, 69, 255);
            }
            QPushButton:hover {
                background-color: rgba(40, 167, 69, 255);
            }
        """)
        
        if self.worker and self.worker.isRunning():
            self.worker.stop()
            self.worker.wait()

        # 合并工作台任务上下文到来源上下文中
        meta = dict(self.context_meta)
        if self.task_context:
            meta["task"] = self.task_context
        # 注入自定义指令
        if custom_instruction:
            meta["custom_instruction"] = custom_instruction
            print(f"[ASU] 自定义指令: {custom_instruction}")

        # 构建 envelope meta（包含所有必要字段，避免 custom_instruction 丢失）
        envelope_meta = {
            "file_name": meta.get("file_name", ""),
            "language": meta.get("language", ""),
        }
        if custom_instruction:
            envelope_meta["custom_instruction"] = custom_instruction
        if self.task_context:
            envelope_meta["task"] = self.task_context

        # 文档修订模式或 IDE 选中文本：构建 context_envelope，注入全文上下文
        envelope = None
        has_ide_selection = hasattr(self, '_ide_selection_range') and self._ide_selection_range
        if action_type == "revision":
            if not self.full_document:
                # 尝试自动读取 IDE 全文
                self._read_ide_silent()
            if self.full_document:
                envelope = {
                    "source": "ide",
                    "content": self.full_document,       # 全文 → Agent 用于交叉扫描
                    "selection": self.current_text,       # 拖拽文本 → 需要修改的目标
                    "task": self.task_context or "",
                    "meta": envelope_meta,
                    "timestamp": time.time(),
                }
                self.context_source = "ide"
            else:
                # 降级：无全文时仅做普通修订分析
                meta["revision_target"] = self.current_text
        elif has_ide_selection and hasattr(self, '_ide_full_document') and self._ide_full_document:
            # IDE 选中文本模式：将全文作为上下文，选中文本作为修改目标
            envelope = {
                "source": "ide",
                "content": self._ide_full_document,      # 全文 → Agent 用于上下文理解
                "selection": self.current_text,            # 选中文本 → 需要修改的目标
                "task": self.task_context or "",
                "meta": envelope_meta,
                "timestamp": time.time(),
            }
            self.context_source = "ide"

        print(f"[ASU] AI请求 | action={action_type} | source={self.context_source} | text_len={len(self.current_text)} | meta_keys={list(meta.keys())}")

        self.worker = AIWorker(self.provider, self.current_text, action_type,
                               self.session_id, self.context_source, meta,
                               context_envelope=envelope,
                               image_base64=image_base64)
        self.worker.text_updated.connect(self.on_text_updated)
        self.worker.finished_signal.connect(self._on_ai_finished)
        self.worker.start()

        # 清空指令输入
        if custom_instruction:
            self.instruction_input.clear()

    def _toggle_revision_mode(self, checked):
        """切换文档修订模式：ON 时自动读取 IDE 全文，失败则提供文件选择（支持 .docx/.pptx）"""
        base_style = self.btn_revision.styleSheet()
        self.revision_mode = checked
        if checked:
            self._read_ide_silent()
            if not self.full_document:
                # IDE 无全文 → 尝试通过 Broker 读 Office 文件
                self._try_read_office()
            if self.full_document:
                self.btn_revision.setText("📝 修订 ON")
                self.btn_revision.setStyleSheet(
                    base_style.replace(
                        "rgba(60, 60, 70, 200)", "rgba(40, 167, 69, 200)"
                    ).replace(
                        "rgba(80, 80, 100, 255)", "rgba(40, 167, 69, 255)"
                    )
                )
                doc_type = self.context_meta.get("file_name", "").split(".")[-1].upper()
                self.text_edit.setPlainText(
                    f"📝 全文修订模式已激活\n\n已读取文档全文（{len(self.full_document)} 字符，{doc_type}）\n\n请将需要修改的文本拖拽到此窗口。"
                )
            else:
                self.btn_revision.setText("📝 修订 (无全文)")
                self.btn_revision.setStyleSheet(
                    base_style.replace(
                        "rgba(60, 60, 70, 200)", "rgba(180, 130, 30, 200)"
                    ).replace(
                        "rgba(80, 80, 100, 255)", "rgba(180, 130, 30, 255)"
                    )
                )
                self.text_edit.setPlainText(
                    "⚠️ 未获取到文档全文。\n\n"
                    "• 在 IDE 中打开文档后重试，或\n"
                    "• 点击「📂 选择文件」加载 .docx/.pptx 文件\n\n"
                    "当前降级模式：仅对拖拽文本做局部修订。"
                )
                self.revision_mode = True
        else:
            self.btn_revision.setText("📝 全文修订")
            self.btn_revision.setStyleSheet(base_style)
            self.full_document = ""
            self.revision_mode = False
            self.text_edit.setPlainText("📝 全文修订模式已关闭。")

    def _try_read_office(self):
        """尝试通过 Broker 读取 .docx/.pptx 文件"""
        # 弹出文件选择对话框
        file_path, _ = QFileDialog.getOpenFileName(
            self, "选择文档进行全文修订",
            os.path.expanduser("~"),
            "文档文件 (*.docx *.pptx);;所有文件 (*)"
        )
        if not file_path:
            return

        try:
            probe = SystemProbeClient()
            if not probe.is_broker_alive():
                self.text_edit.setPlainText(
                    "⚠️ Broker 未启动，无法解析 Office 文件。\n请在原生终端运行：cd asu_broker && python run.py"
                )
                return
            result = probe.read_office_file(file_path)
            content = result.get("content", "")
            if content:
                self.full_document = content
                fname = os.path.basename(file_path)
                self.context_meta = {"file_name": fname, "language": result.get("type", "docx")}
                self.context_source = "ide"  # 统一走 IDE 上下文路径
        except Exception as e:
            self.text_edit.setPlainText(f"❌ 解析 Office 文件失败: {str(e)}")

    def _read_ide_silent(self):
        """静默读取 IDE 全文（不更新 UI 显示）"""
        port = self._get_ide_port()
        if not port:
            return
        try:
            response = httpx.get(f"http://127.0.0.1:{port}/context", timeout=2.0)
            if response.status_code == 200:
                data = response.json()
                self.full_document = data.get("content", "")
                self.context_meta = {
                    "file_name": data.get("fileName", "Unknown"),
                    "language": data.get("languageId", ""),
                }
        except Exception:
            pass  # 静默失败，不打断用户体验

    def on_text_updated(self, text):
        if not text:
            self.text_edit.setPlainText("正在思考...\n")
        else:
            self.text_edit.setHtml(md_render(text))
        scrollbar = self.text_edit.verticalScrollBar()
        scrollbar.setValue(scrollbar.maximum())

    def _on_ai_finished(self):
        """AI 回复完成后，显示回写/复制按钮。"""
        # 检查是否有 IDE 插件可用
        has_ide = bool(self._get_ide_port())
        if has_ide:
            # 如果还没有选区范围，尝试自动获取 IDE 选区（用于局部回写）
            if not (hasattr(self, '_ide_selection_range') and self._ide_selection_range):
                self._try_get_ide_selection()
            self.btn_apply_to_ide.show()
        self.btn_copy_result.show()
        
        # 如果包含标题或者是 JSON 数组格式，则显示生成PPT按钮
        if self.worker and hasattr(self.worker, 'full_text'):
            text = self.worker.full_text
            if '# ' in text or '## ' in text or extract_json_from_text(text):
                self.btn_export_ppt.show()
        elif self.chat_worker and hasattr(self.chat_worker, 'full_text'):
            text = self.chat_worker.full_text
            if '# ' in text or '## ' in text or extract_json_from_text(text):
                self.btn_export_ppt.show()

    def _export_to_ppt(self):
        """将 AI 生成的 Markdown 大纲，通过人机共创界面，最终导出为 PPT"""
        text = ""
        # 根据当前在哪个tab，获取内容
        if self.tabs.currentIndex() == 0 and self.worker and hasattr(self.worker, 'full_text'):
            text = self.worker.full_text
        elif self.tabs.currentIndex() == 1 and self.chat_worker and hasattr(self.chat_worker, 'full_text'):
            text = self.chat_worker.full_text
            
        if not text:
            QMessageBox.warning(self, "导出失败", "没有可导出的内容。")
            return
            
        json_data = extract_json_from_text(text)
        if not json_data:
            QMessageBox.warning(self, "解析失败", "大模型输出的内容不符合预期的 JSON 大纲格式。")
            return
            
        # 唤起三面板人机共创编辑器
        try:
            from opencopilot.capabilities.ppt import CoCreationDialog
            dialog = CoCreationDialog(
                original_text=text,
                json_data=json_data,
                agent_url="http://127.0.0.1:18888",
                parent=self
            )
            
            if dialog.exec() == QDialog.DialogCode.Accepted:
                output_path = dialog.get_output_path()
                if output_path:
                    QMessageBox.information(self, "导出成功", f"PPT 已成功导出至：\n{output_path}")
        except ImportError as e:
            # 降级到旧版编辑器
            print(f"[WARN] 无法加载新共创编辑器，降级到旧版: {e}")
            dialog = PPTPreviewDialog(json_data, self)
            if dialog.exec() == QDialog.DialogCode.Accepted:
                final_json = dialog.get_final_json()
                try:
                    save_path, _ = QFileDialog.getSaveFileName(
                        self, "导出为 PPT",
                        os.path.expanduser("~/Desktop/generated_presentation.pptx"),
                        "PowerPoint Files (*.pptx)"
                    )
                    if save_path:
                        generate_ppt_from_json(final_json, save_path)
                        QMessageBox.information(self, "导出成功", f"PPT 已保存至:\n{save_path}")
                        subprocess.run(["open", save_path])
                except Exception as e:
                    import traceback
                    print(traceback.format_exc())
                    QMessageBox.critical(self, "导出错误", f"导出失败: {e}")
        except Exception as e:
            import traceback
            print(traceback.format_exc())
            QMessageBox.critical(self, "导出错误", f"启动编辑器失败: {e}")

    def hide_card(self):
        self._pending_hide = False
        self._user_initiated_hide = True
        self.hide()
        if self.worker and self.worker.isRunning():
            self.worker.stop()

    def hideEvent(self, event):
        """拦截 macOS 对 Tool 窗口的自动隐藏：只在用户主动隐藏时才允许。"""
        # 正在拖拽时，不允许隐藏
        if self._is_dragging:
            event.ignore()
            return
        
        if self._user_initiated_hide or self._allow_close:
            self._user_initiated_hide = False
            super().hideEvent(event)
        else:
            # 系统/焦点变化导致的自动隐藏 → 忽略，并重新显示
            event.ignore()
            QTimer.singleShot(0, self._force_reshow)

    def _force_reshow(self):
        """macOS 失焦后重新显示并置顶。"""
        if not self._user_initiated_hide and not self._allow_close:
            self.show()
            self.raise_()
            make_panel_persistent(self)

    def closeEvent(self, event):
        if self._allow_close:
            event.accept()
            return
        self.hide_card()
        event.ignore()


# ==========================================
# 4. 任务工作台 (三击右键唤出)
# ==========================================
