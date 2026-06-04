"""
FormatSkill - 格式转换技能
封装 MarkdownToDocxTool、MarkdownToPptxTool、TextToTableTool
"""

import os
from typing import Dict, Any, Optional
from .base import BaseSkill
from .models import (
    SkillMetadata, SkillContext, SkillResult,
    SkillStatus, ExecutionMode
)


class FormatSkill(BaseSkill):
    """格式转换技能
    
    封装格式转换相关的工具，提供统一的技能接口。
    支持 Markdown 转 Word、Markdown 转 PPT、文本转表格等转换功能。
    """
    
    @property
    def metadata(self) -> SkillMetadata:
        return SkillMetadata(
            name="format_skill",
            description="格式转换技能，支持Markdown转Word、Markdown转PPT、文本转表格等多种格式转换",
            version="1.0.0",
            author="OpenCopilot",
            intents=[
                "format_convert",
                "md_to_docx",
                "md_to_pptx",
                "text_to_table",
                "markdown_convert",
                "document_convert"
            ],
            tags=["format", "convert", "markdown", "docx", "pptx", "table"],
            dependencies=["python-docx", "python-pptx"],
            config_schema={
                "max_content_size": 10 * 1024 * 1024,  # 10MB
                "supported_formats": {
                    "markdown_to_docx": {"input": ["md", "markdown"], "output": "docx"},
                    "markdown_to_pptx": {"input": ["md", "markdown"], "output": "pptx"},
                    "text_to_table": {"input": ["text", "csv", "tsv"], "output": ["markdown", "html", "csv"]}
                }
            }
        )
    
    async def can_handle(self, context: SkillContext) -> float:
        """判断是否能处理该上下文
        
        Args:
            context: 执行上下文
        
        Returns:
            float: 置信度 (0-1)
        """
        # 检查意图
        if context.intent in self.metadata.intents:
            return 0.9
        
        # 检查输入数据
        if "action" in context.input_data:
            action = context.input_data["action"]
            if action in ["md_to_docx", "md_to_pptx", "text_to_table", "convert"]:
                return 0.8
        
        # 检查内容类型
        content = context.input_data.get("content", "")
        if isinstance(content, str):
            content_lower = content.lower()
            format_keywords = ["转换", "格式", "markdown", "word", "docx", "pptx", 
                             "表格", "convert", "format", "table"]
            if any(keyword in content_lower for keyword in format_keywords):
                return 0.7
        
        return 0.0
    
    async def execute(self, context: SkillContext) -> SkillResult:
        """执行格式转换
        
        Args:
            context: 执行上下文，包含:
                - action: 操作类型 (md_to_docx/md_to_pptx/text_to_table)
                - content: 源内容
                - output_path: 输出路径 (可选，用于文件转换)
                - format: 输出格式 (可选，用于表格转换)
                - template: 模板路径 (可选)
        
        Returns:
            SkillResult: 执行结果
        """
        try:
            action = context.input_data.get("action", "md_to_docx")
            content = context.input_data.get("content", "")
            
            # 根据操作类型分发
            if action == "md_to_docx":
                return await self._convert_md_to_docx(context)
            elif action == "md_to_pptx":
                return await self._convert_md_to_pptx(context)
            elif action == "text_to_table":
                return await self._convert_text_to_table(context)
            else:
                # 尝试自动检测
                return await self._auto_convert(context)
                
        except Exception as e:
            return SkillResult(
                success=False,
                data={},
                error=f"格式转换失败: {str(e)}",
                status=SkillStatus.FAILED
            )
    
    async def _convert_md_to_docx(self, context: SkillContext) -> SkillResult:
        """将 Markdown 转换为 Word 文档
        
        Args:
            context: 执行上下文
        
        Returns:
            SkillResult: 转换结果
        """
        content = context.input_data.get("content", "")
        output_path = context.input_data.get("output_path")
        template = context.input_data.get("template")
        
        if not content:
            return SkillResult(
                success=False,
                data={},
                error="缺少 Markdown 内容",
                status=SkillStatus.FAILED
            )
        
        try:
            from docx import Document
            from docx.shared import Pt
            import re
            
            # 创建文档
            if template and os.path.exists(template):
                doc = Document(template)
            else:
                doc = Document()
            
            # 解析Markdown内容
            lines = content.split('\n')
            
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                
                # 处理标题
                if line.startswith('#'):
                    match = re.match(r'^(#+)', line)
                    if match:
                        level = len(match.group(1))
                        title_text = line.lstrip('#').strip()
                        
                        if level == 1:
                            doc.add_heading(title_text, 0)
                        elif level == 2:
                            doc.add_heading(title_text, 1)
                        elif level == 3:
                            doc.add_heading(title_text, 2)
                        else:
                            doc.add_heading(title_text, 3)
                
                # 处理列表
                elif line.startswith('- ') or line.startswith('* '):
                    item_text = line[2:].strip()
                    doc.add_paragraph(item_text, style='List Bullet')
                
                elif re.match(r'^\d+\.\s', line):
                    item_text = re.sub(r'^\d+\.\s', '', line).strip()
                    doc.add_paragraph(item_text, style='List Number')
                
                # 处理引用
                elif line.startswith('>'):
                    quote_text = line[1:].strip()
                    p = doc.add_paragraph()
                    p.style = 'Quote'
                    p.add_run(quote_text)
                
                # 处理普通段落
                else:
                    text = line
                    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
                    text = re.sub(r'\*(.*?)\*', r'\1', text)
                    doc.add_paragraph(text)
            
            # 如果指定了输出路径，保存文件
            if output_path:
                expanded_output = os.path.expanduser(output_path)
                os.makedirs(os.path.dirname(expanded_output), exist_ok=True)
                doc.save(expanded_output)
                
                return SkillResult(
                    success=True,
                    data={
                        "output_path": expanded_output,
                        "format": "docx",
                        "size": os.path.getsize(expanded_output)
                    },
                    status=SkillStatus.COMPLETED
                )
            else:
                # 返回文档对象的文本内容
                full_text = '\n'.join([p.text for p in doc.paragraphs])
                return SkillResult(
                    success=True,
                    data={
                        "content": full_text,
                        "format": "docx",
                        "paragraphs": len(doc.paragraphs)
                    },
                    status=SkillStatus.COMPLETED
                )
                
        except ImportError:
            return SkillResult(
                success=False,
                data={},
                error="需要安装 python-docx: pip install python-docx",
                status=SkillStatus.FAILED
            )
        except Exception as e:
            return SkillResult(
                success=False,
                data={},
                error=f"Markdown 转 Word 失败: {str(e)}",
                status=SkillStatus.FAILED
            )
    
    async def _convert_md_to_pptx(self, context: SkillContext) -> SkillResult:
        """将 Markdown 转换为 PPT 演示文稿
        
        Args:
            context: 执行上下文
        
        Returns:
            SkillResult: 转换结果
        """
        content = context.input_data.get("content", "")
        output_path = context.input_data.get("output_path")
        template = context.input_data.get("template")
        
        if not content:
            return SkillResult(
                success=False,
                data={},
                error="缺少 Markdown 内容",
                status=SkillStatus.FAILED
            )
        
        try:
            from pptx import Presentation
            from pptx.util import Inches
            import re
            
            # 创建演示文稿
            if template and os.path.exists(template):
                prs = Presentation(template)
            else:
                prs = Presentation()
            
            # 解析Markdown内容
            lines = content.split('\n')
            current_slide = None
            current_content = []
            
            def add_content_to_slide(slide, content_items):
                """向幻灯片添加内容"""
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:
                        shape.text = '\n'.join(content_items)
                        break
            
            for line in lines:
                line = line.strip()
                if not line:
                    continue
                
                # 处理标题（创建新幻灯片）
                if line.startswith('#'):
                    # 保存之前的内容
                    if current_slide and current_content:
                        add_content_to_slide(current_slide, current_content)
                        current_content = []
                    
                    match = re.match(r'^(#+)', line)
                    if match:
                        title_text = line.lstrip('#').strip()
                        
                        # 创建新幻灯片
                        slide_layout = prs.slide_layouts[1]
                        current_slide = prs.slides.add_slide(slide_layout)
                        
                        # 设置标题
                        title = current_slide.shapes.title
                        title.text = title_text
                
                # 处理列表项
                elif line.startswith('- ') or line.startswith('* '):
                    item_text = line[2:].strip()
                    current_content.append(f"• {item_text}")
                
                elif re.match(r'^\d+\.\s', line):
                    current_content.append(line)
                
                # 处理普通文本
                else:
                    current_content.append(line)
            
            # 保存最后的内容
            if current_slide and current_content:
                add_content_to_slide(current_slide, current_content)
            
            # 如果指定了输出路径，保存文件
            if output_path:
                expanded_output = os.path.expanduser(output_path)
                os.makedirs(os.path.dirname(expanded_output), exist_ok=True)
                prs.save(expanded_output)
                
                return SkillResult(
                    success=True,
                    data={
                        "output_path": expanded_output,
                        "format": "pptx",
                        "slides": len(prs.slides),
                        "size": os.path.getsize(expanded_output)
                    },
                    status=SkillStatus.COMPLETED
                )
            else:
                return SkillResult(
                    success=True,
                    data={
                        "format": "pptx",
                        "slides": len(prs.slides)
                    },
                    status=SkillStatus.COMPLETED
                )
                
        except ImportError:
            return SkillResult(
                success=False,
                data={},
                error="需要安装 python-pptx: pip install python-pptx",
                status=SkillStatus.FAILED
            )
        except Exception as e:
            return SkillResult(
                success=False,
                data={},
                error=f"Markdown 转 PPT 失败: {str(e)}",
                status=SkillStatus.FAILED
            )
    
    async def _convert_text_to_table(self, context: SkillContext) -> SkillResult:
        """将文本转换为表格
        
        Args:
            context: 执行上下文
        
        Returns:
            SkillResult: 转换结果
        """
        content = context.input_data.get("content", "")
        output_format = context.input_data.get("format", "markdown")
        delimiter = context.input_data.get("delimiter", ",")
        
        if not content:
            return SkillResult(
                success=False,
                data={},
                error="缺少文本内容",
                status=SkillStatus.FAILED
            )
        
        try:
            import re
            
            # 解析文本
            rows = []
            for line in content.split('\n'):
                line = line.strip()
                if not line:
                    continue
                
                # 尝试按分隔符分割
                if delimiter in line:
                    cells = [cell.strip() for cell in line.split(delimiter)]
                    rows.append(cells)
                # 尝试按制表符分割
                elif '\t' in line:
                    cells = [cell.strip() for cell in line.split('\t')]
                    rows.append(cells)
                # 尝试按多个空格分割
                elif re.search(r'\s{2,}', line):
                    cells = [cell.strip() for cell in re.split(r'\s{2,}', line)]
                    rows.append(cells)
                else:
                    rows.append([line])
            
            if not rows:
                return SkillResult(
                    success=False,
                    data={},
                    error="无法解析文本为表格",
                    status=SkillStatus.FAILED
                )
            
            # 确保所有行有相同的列数
            max_cols = max(len(row) for row in rows)
            for row in rows:
                while len(row) < max_cols:
                    row.append("")
            
            # 生成表格
            if output_format == "markdown":
                table_text = self._to_markdown(rows)
            elif output_format == "html":
                table_text = self._to_html(rows)
            elif output_format == "csv":
                table_text = self._to_csv(rows, delimiter)
            else:
                return SkillResult(
                    success=False,
                    data={},
                    error=f"不支持的格式: {output_format}",
                    status=SkillStatus.FAILED
                )
            
            return SkillResult(
                success=True,
                data={
                    "content": table_text,
                    "format": output_format,
                    "rows": len(rows),
                    "columns": max_cols
                },
                status=SkillStatus.COMPLETED
            )
                
        except Exception as e:
            return SkillResult(
                success=False,
                data={},
                error=f"文本转表格失败: {str(e)}",
                status=SkillStatus.FAILED
            )
    
    async def _auto_convert(self, context: SkillContext) -> SkillResult:
        """自动检测并转换
        
        Args:
            context: 执行上下文
        
        Returns:
            SkillResult: 转换结果
        """
        content = context.input_data.get("content", "")
        target_format = context.input_data.get("target_format", "docx")
        
        if not content:
            return SkillResult(
                success=False,
                data={},
                error="缺少内容",
                status=SkillStatus.FAILED
            )
        
        # 根据目标格式选择转换方法
        if target_format in ["docx", "word"]:
            return await self._convert_md_to_docx(context)
        elif target_format in ["pptx", "ppt"]:
            return await self._convert_md_to_pptx(context)
        elif target_format in ["markdown", "html", "csv"]:
            context.input_data["format"] = target_format
            return await self._convert_text_to_table(context)
        else:
            return SkillResult(
                success=False,
                data={},
                error=f"不支持的目标格式: {target_format}",
                status=SkillStatus.FAILED
            )
    
    def _to_markdown(self, rows: list) -> str:
        """转换为 Markdown 表格"""
        if not rows:
            return ""
        
        # 计算每列的最大宽度
        col_widths = []
        for col_idx in range(len(rows[0])):
            max_width = 0
            for row in rows:
                if col_idx < len(row):
                    max_width = max(max_width, len(row[col_idx]))
            col_widths.append(max_width)
        
        # 生成Markdown表格
        lines = []
        
        # 表头
        header = "| " + " | ".join(
            cell.ljust(col_widths[i]) for i, cell in enumerate(rows[0])
        ) + " |"
        lines.append(header)
        
        # 分隔线
        separator = "| " + " | ".join(
            "-" * col_widths[i] for i in range(len(rows[0]))
        ) + " |"
        lines.append(separator)
        
        # 数据行
        for row in rows[1:]:
            line = "| " + " | ".join(
                cell.ljust(col_widths[i]) for i, cell in enumerate(row)
            ) + " |"
            lines.append(line)
        
        return "\n".join(lines)
    
    def _to_html(self, rows: list) -> str:
        """转换为 HTML 表格"""
        if not rows:
            return ""
        
        lines = ["<table>"]
        
        # 表头
        lines.append("  <thead>")
        lines.append("    <tr>")
        for cell in rows[0]:
            lines.append(f"      <th>{cell}</th>")
        lines.append("    </tr>")
        lines.append("  </thead>")
        
        # 表体
        lines.append("  <tbody>")
        for row in rows[1:]:
            lines.append("    <tr>")
            for cell in row:
                lines.append(f"      <td>{cell}</td>")
            lines.append("    </tr>")
        lines.append("  </tbody>")
        
        lines.append("</table>")
        
        return "\n".join(lines)
    
    def _to_csv(self, rows: list, delimiter: str = ",") -> str:
        """转换为 CSV 格式"""
        if not rows:
            return ""
        
        lines = []
        for row in rows:
            # 处理包含分隔符的单元格
            escaped_cells = []
            for cell in row:
                if delimiter in cell or '"' in cell or '\n' in cell:
                    escaped_cell = '"' + cell.replace('"', '""') + '"'
                    escaped_cells.append(escaped_cell)
                else:
                    escaped_cells.append(cell)
            
            lines.append(delimiter.join(escaped_cells))
        
        return "\n".join(lines)
