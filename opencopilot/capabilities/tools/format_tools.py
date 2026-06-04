"""
格式转换工具
"""

import os
import asyncio
import re
from typing import Any, Dict, List
from .base import BaseTool


class MarkdownToDocxTool(BaseTool):
    """Markdown转Word工具"""
    
    @property
    def name(self) -> str:
        return "md_to_docx"
    
    @property
    def description(self) -> str:
        return "将Markdown转换为Word文档"
    
    @property
    def parameters(self) -> Dict[str, Any]:
        return {
            "md_content": {
                "type": "string",
                "description": "Markdown内容",
                "required": True
            },
            "output_path": {
                "type": "string",
                "description": "输出文件路径",
                "required": True
            },
            "template": {
                "type": "string",
                "description": "模板文件路径（可选）"
            }
        }
    
    async def execute(self, **kwargs) -> Dict[str, Any]:
        """执行转换"""
        md_content = kwargs.get("md_content")
        output_path = kwargs.get("output_path")
        template = kwargs.get("template")
        
        if not md_content or not output_path:
            return {"error": "md_content and output_path are required"}
        
        expanded_output = os.path.expanduser(output_path)
        
        try:
            from docx import Document
            from docx.shared import Pt, Inches
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            
            loop = asyncio.get_event_loop()
            
            def convert_sync():
                # 创建文档
                if template and os.path.exists(template):
                    doc = Document(template)
                else:
                    doc = Document()
                
                # 解析Markdown内容
                lines = md_content.split('\n')
                current_paragraph = None
                
                for line in lines:
                    line = line.strip()
                    if not line:
                        continue
                    
                    # 处理标题
                    if line.startswith('#'):
                        level = len(re.match(r'^(#+)', line).group(1))
                        title_text = line.lstrip('#').strip()
                        
                        if level == 1:
                            doc.add_heading(title_text, 0)
                        elif level == 2:
                            doc.add_heading(title_text, 1)
                        elif level == 3:
                            doc.add_heading(title_text, 2)
                        else:
                            doc.add_heading(title_text, 3)
                    
                    # 处理列表
                    elif line.startswith('- ') or line.startswith('* '):
                        item_text = line[2:].strip()
                        doc.add_paragraph(item_text, style='List Bullet')
                    
                    elif re.match(r'^\d+\.\s', line):
                        item_text = re.sub(r'^\d+\.\s', '', line).strip()
                        doc.add_paragraph(item_text, style='List Number')
                    
                    # 处理引用
                    elif line.startswith('>'):
                        quote_text = line[1:].strip()
                        p = doc.add_paragraph()
                        p.style = 'Quote'
                        p.add_run(quote_text)
                    
                    # 处理代码块
                    elif line.startswith('```'):
                        # 简单处理，实际应用中需要更复杂的逻辑
                        pass
                    
                    # 处理普通段落
                    else:
                        # 处理加粗和斜体
                        text = line
                        text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)  # 移除加粗标记
                        text = re.sub(r'\*(.*?)\*', r'\1', text)  # 移除斜体标记
                        
                        p = doc.add_paragraph(text)
                
                # 保存文档
                os.makedirs(os.path.dirname(expanded_output), exist_ok=True)
                doc.save(expanded_output)
            
            await loop.run_in_executor(None, convert_sync)
            
            return {
                "success": True,
                "output_path": expanded_output,
                "format": "docx"
            }
        except ImportError:
            return {"error": "需要安装 python-docx: pip install python-docx"}
        except Exception as e:
            return {"error": f"转换失败: {str(e)}"}


class MarkdownToPptxTool(BaseTool):
    """Markdown转PPT工具"""
    
    @property
    def name(self) -> str:
        return "md_to_pptx"
    
    @property
    def description(self) -> str:
        return "将Markdown转换为PPT演示文稿"
    
    @property
    def parameters(self) -> Dict[str, Any]:
        return {
            "md_content": {
                "type": "string",
                "description": "Markdown内容",
                "required": True
            },
            "output_path": {
                "type": "string",
                "description": "输出文件路径",
                "required": True
            },
            "template": {
                "type": "string",
                "description": "模板文件路径（可选）"
            }
        }
    
    async def execute(self, **kwargs) -> Dict[str, Any]:
        """执行转换"""
        md_content = kwargs.get("md_content")
        output_path = kwargs.get("output_path")
        template = kwargs.get("template")
        
        if not md_content or not output_path:
            return {"error": "md_content and output_path are required"}
        
        expanded_output = os.path.expanduser(output_path)
        
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
            
            loop = asyncio.get_event_loop()
            
            def convert_sync():
                # 创建演示文稿
                if template and os.path.exists(template):
                    prs = Presentation(template)
                else:
                    prs = Presentation()
                
                # 解析Markdown内容
                lines = md_content.split('\n')
                current_slide = None
                current_content = []
                
                for line in lines:
                    line = line.strip()
                    if not line:
                        continue
                    
                    # 处理标题（创建新幻灯片）
                    if line.startswith('#'):
                        # 保存之前的内容
                        if current_slide and current_content:
                            self._add_content_to_slide(current_slide, current_content)
                            current_content = []
                        
                        level = len(re.match(r'^(#+)', line).group(1))
                        title_text = line.lstrip('#').strip()
                        
                        # 创建新幻灯片
                        slide_layout = prs.slide_layouts[1]  # 标题和内容布局
                        current_slide = prs.slides.add_slide(slide_layout)
                        
                        # 设置标题
                        title = current_slide.shapes.title
                        title.text = title_text
                    
                    # 处理列表项
                    elif line.startswith('- ') or line.startswith('* '):
                        item_text = line[2:].strip()
                        current_content.append(f"• {item_text}")
                    
                    elif re.match(r'^\d+\.\s', line):
                        current_content.append(line)
                    
                    # 处理普通文本
                    else:
                        current_content.append(line)
                
                # 保存最后的内容
                if current_slide and current_content:
                    self._add_content_to_slide(current_slide, current_content)
                
                # 保存演示文稿
                os.makedirs(os.path.dirname(expanded_output), exist_ok=True)
                prs.save(expanded_output)
            
            def _add_content_to_slide(slide, content):
                """向幻灯片添加内容"""
                # 查找内容占位符
                for shape in slide.placeholders:
                    if shape.placeholder_format.idx == 1:  # 内容占位符
                        shape.text = '\n'.join(content)
                        break
            
            await loop.run_in_executor(None, convert_sync)
            
            return {
                "success": True,
                "output_path": expanded_output,
                "format": "pptx"
            }
        except ImportError:
            return {"error": "需要安装 python-pptx: pip install python-pptx"}
        except Exception as e:
            return {"error": f"转换失败: {str(e)}"}
    
    def _add_content_to_slide(self, slide, content):
        """向幻灯片添加内容"""
        # 查找内容占位符
        for shape in slide.placeholders:
            if shape.placeholder_format.idx == 1:  # 内容占位符
                shape.text = '\n'.join(content)
                break


class TextToTableTool(BaseTool):
    """文本转表格工具"""
    
    @property
    def name(self) -> str:
        return "text_to_table"
    
    @property
    def description(self) -> str:
        return "将结构化文本转换为表格"
    
    @property
    def parameters(self) -> Dict[str, Any]:
        return {
            "text": {
                "type": "string",
                "description": "结构化文本",
                "required": True
            },
            "format": {
                "type": "string",
                "description": "输出格式",
                "enum": ["markdown", "html", "csv"],
                "default": "markdown"
            },
            "delimiter": {
                "type": "string",
                "description": "分隔符",
                "default": ","
            }
        }
    
    async def execute(self, **kwargs) -> Dict[str, Any]:
        """执行转换"""
        text = kwargs.get("text")
        output_format = kwargs.get("format", "markdown")
        delimiter = kwargs.get("delimiter", ",")
        
        if not text:
            return {"error": "text is required"}
        
        try:
            # 解析文本
            rows = self._parse_text(text, delimiter)
            
            if not rows:
                return {"error": "无法解析文本为表格"}
            
            # 生成表格
            if output_format == "markdown":
                table_text = self._to_markdown(rows)
            elif output_format == "html":
                table_text = self._to_html(rows)
            elif output_format == "csv":
                table_text = self._to_csv(rows, delimiter)
            else:
                return {"error": f"不支持的格式: {output_format}"}
            
            return {
                "type": "table",
                "format": output_format,
                "content": table_text,
                "rows": len(rows),
                "columns": len(rows[0]) if rows else 0
            }
        except Exception as e:
            return {"error": f"转换失败: {str(e)}"}
    
    def _parse_text(self, text: str, delimiter: str) -> List[List[str]]:
        """解析文本为表格数据"""
        rows = []
        
        for line in text.split('\n'):
            line = line.strip()
            if not line:
                continue
            
            # 尝试按分隔符分割
            if delimiter in line:
                cells = [cell.strip() for cell in line.split(delimiter)]
                rows.append(cells)
            # 尝试按制表符分割
            elif '\t' in line:
                cells = [cell.strip() for cell in line.split('\t')]
                rows.append(cells)
            # 尝试按多个空格分割
            elif re.search(r'\s{2,}', line):
                cells = [cell.strip() for cell in re.split(r'\s{2,}', line)]
                rows.append(cells)
            else:
                # 单行文本，作为单列处理
                rows.append([line])
        
        # 确保所有行有相同的列数
        if rows:
            max_cols = max(len(row) for row in rows)
            for row in rows:
                while len(row) < max_cols:
                    row.append("")
        
        return rows
    
    def _to_markdown(self, rows: List[List[str]]) -> str:
        """转换为Markdown表格"""
        if not rows:
            return ""
        
        # 计算每列的最大宽度
        col_widths = []
        for col_idx in range(len(rows[0])):
            max_width = 0
            for row in rows:
                if col_idx < len(row):
                    max_width = max(max_width, len(row[col_idx]))
            col_widths.append(max_width)
        
        # 生成Markdown表格
        lines = []
        
        # 表头
        header = "| " + " | ".join(
            cell.ljust(col_widths[i]) for i, cell in enumerate(rows[0])
        ) + " |"
        lines.append(header)
        
        # 分隔线
        separator = "| " + " | ".join(
            "-" * col_widths[i] for i in range(len(rows[0]))
        ) + " |"
        lines.append(separator)
        
        # 数据行
        for row in rows[1:]:
            line = "| " + " | ".join(
                cell.ljust(col_widths[i]) for i, cell in enumerate(row)
            ) + " |"
            lines.append(line)
        
        return "\n".join(lines)
    
    def _to_html(self, rows: List[List[str]]) -> str:
        """转换为HTML表格"""
        if not rows:
            return ""
        
        lines = ["<table>"]
        
        # 表头
        lines.append("  <thead>")
        lines.append("    <tr>")
        for cell in rows[0]:
            lines.append(f"      <th>{cell}</th>")
        lines.append("    </tr>")
        lines.append("  </thead>")
        
        # 表体
        lines.append("  <tbody>")
        for row in rows[1:]:
            lines.append("    <tr>")
            for cell in row:
                lines.append(f"      <td>{cell}</td>")
            lines.append("    </tr>")
        lines.append("  </tbody>")
        
        lines.append("</table>")
        
        return "\n".join(lines)
    
    def _to_csv(self, rows: List[List[str]], delimiter: str) -> str:
        """转换为CSV格式"""
        if not rows:
            return ""
        
        lines = []
        for row in rows:
            # 处理包含分隔符的单元格
            escaped_cells = []
            for cell in row:
                if delimiter in cell or '"' in cell or '\n' in cell:
                    escaped_cell = '"' + cell.replace('"', '""') + '"'
                    escaped_cells.append(escaped_cell)
                else:
                    escaped_cells.append(cell)
            
            lines.append(delimiter.join(escaped_cells))
        
        return "\n".join(lines)