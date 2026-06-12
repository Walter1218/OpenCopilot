import re
import os
import json
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# 配置日志
logger = logging.getLogger(__name__)

def clean_markdown(text):
    """清理 Markdown 标记，返回纯文本（兼容旧调用）"""
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    return text.strip()


def parse_inline_formatting(text):
    """
    解析内联格式，返回 (clean_text, runs_info)
    runs_info = [{"text": "...", "bold": True/False, "italic": True/False}, ...]
    """
    runs = []
    # 正则匹配 **bold** 或 *italic* 片段
    pattern = r'(\*\*(.+?)\*\*|\*(.+?)\*)'
    last_end = 0
    for match in re.finditer(pattern, text):
        # 前置普通文本
        if match.start() > last_end:
            plain = text[last_end:match.start()]
            if plain:
                runs.append({"text": plain, "bold": False, "italic": False})
        # 格式化文本
        if match.group(2):  # **bold**
            runs.append({"text": match.group(2), "bold": True, "italic": False})
        elif match.group(3):  # *italic*
            runs.append({"text": match.group(3), "bold": False, "italic": True})
        last_end = match.end()
    # 尾部普通文本
    if last_end < len(text):
        runs.append({"text": text[last_end:], "bold": False, "italic": False})
    
    if not runs:
        runs.append({"text": text, "bold": False, "italic": False})
    
    clean = "".join(r["text"] for r in runs)
    return clean, runs


def _apply_formatted_paragraph(paragraph, text):
    """将含 **bold** / *italic* 的文本应用到 paragraph，正确设置 run 样式"""
    _, runs = parse_inline_formatting(text)
    # 清空默认 run
    for run in paragraph.runs:
        run.text = ""
    for i, r in enumerate(runs):
        if i == 0 and paragraph.runs:
            run = paragraph.runs[0]
            run.text = r["text"]
        else:
            run = paragraph.add_run()
            run.text = r["text"]
        run.font.bold = r["bold"]
        run.font.italic = r["italic"]

def apply_corporate_theme(slide, prs, is_title_slide=False):
    """应用现代化极简商务主题"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 252)
    
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 
        Inches(0), prs.slide_height - Inches(0.1), 
        prs.slide_width, Inches(0.1)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 82, 204)
    shape.line.fill.background()

    if is_title_slide:
        dec1 = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_TRIANGLE, 
            prs.slide_width - Inches(5), Inches(0), 
            Inches(5), prs.slide_height
        )
        dec1.fill.solid()
        dec1.fill.fore_color.rgb = RGBColor(235, 240, 248)
        dec1.line.fill.background()
        
        dec2 = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_TRIANGLE, 
            prs.slide_width - Inches(3.5), Inches(0), 
            Inches(3.5), prs.slide_height
        )
        dec2.fill.solid()
        dec2.fill.fore_color.rgb = RGBColor(0, 82, 204)
        dec2.line.fill.background()

def format_title_slide(slide, title_text, subtitle_text=""):
    """格式化封面"""
    title_shape = slide.shapes.title
    title_shape.text = clean_markdown(title_text)
    
    title_shape.width = Inches(8)
    title_shape.left = Inches(1)
    title_shape.top = Inches(2.5)
    
    for p in title_shape.text_frame.paragraphs:
        p.font.bold = True
        p.font.size = Pt(44)
        p.font.color.rgb = RGBColor(15, 25, 45)
        p.alignment = PP_ALIGN.LEFT
        
    if subtitle_text and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = clean_markdown(subtitle_text)
        subtitle_shape.width = Inches(8)
        subtitle_shape.left = Inches(1)
        subtitle_shape.top = Inches(3.8)
        for p in subtitle_shape.text_frame.paragraphs:
            p.font.size = Pt(24)
            p.font.color.rgb = RGBColor(100, 100, 110)
            p.alignment = PP_ALIGN.LEFT

def add_placeholder_image(slide, prs):
    """为 image_right 版式添加右侧配图占位"""
    try:
        img_width = Inches(4.5)
        img_height = prs.slide_height
        left = prs.slide_width - img_width
        top = Inches(0)
        
        # 使用几何图形拼接出具有设计感的占位图
        shape1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, img_width, img_height)
        shape1.fill.solid()
        shape1.fill.fore_color.rgb = RGBColor(235, 240, 245)
        shape1.line.fill.background()
        
        shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left + Inches(0.5), top + Inches(1.5), img_width - Inches(1.0), Inches(4.5))
        shape2.fill.solid()
        shape2.fill.fore_color.rgb = RGBColor(0, 102, 204)
        shape2.line.fill.background()
        
        shape3 = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(3.0), top + Inches(4.5), Inches(2.0), Inches(2.0))
        shape3.fill.solid()
        shape3.fill.fore_color.rgb = RGBColor(255, 153, 51)
        shape3.line.fill.background()
    except Exception as e:
        pass

def format_content_slide(slide, title_text, items, layout_type="text_only", prs=None, slide_data=None):
    """格式化内容页，支持多种版式：text_only / table / code / image_right / three_columns"""
    title_shape = slide.shapes.title
    
    title_shape.text = clean_markdown(title_text)
    title_shape.width = Inches(11.333)
    title_shape.left = Inches(1)
    title_shape.top = Inches(0.5)
    title_shape.height = Inches(1.2)
    
    for p in title_shape.text_frame.paragraphs:
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = RGBColor(0, 82, 204)
        p.alignment = PP_ALIGN.LEFT

    # ---- 表格布局 ----
    if layout_type == "table" and slide_data and "table_data" in slide_data:
        table_rows = slide_data["table_data"]
        if len(table_rows) >= 2:
            # 解析表头（第一行）和数据行
            header_cells = [c.strip() for c in table_rows[0].strip('|').split('|')]
            data_rows = []
            for row in table_rows[1:]:
                # 跳过分隔行 (|---|---|)
                if re.match(r'^[\|\s\-:]+$', row):
                    continue
                cells = [c.strip() for c in row.strip('|').split('|')]
                data_rows.append(cells)
            
            num_cols = len(header_cells)
            num_rows = len(data_rows) + 1  # +1 for header
            
            if num_cols > 0 and data_rows:
                table_left = Inches(1)
                table_top = Inches(1.8)
                table_width = Inches(11.333)
                table_height = Inches(min(4.5, 0.4 * num_rows + 0.5))
                
                table_shape = slide.shapes.add_table(
                    num_rows, num_cols, table_left, table_top, table_width, table_height
                )
                table = table_shape.table
                
                # 设置列宽
                col_width = table_width / num_cols
                for c in range(num_cols):
                    table.columns[c].width = int(col_width)
                
                # 表头行
                for c, cell_text in enumerate(header_cells):
                    cell = table.cell(0, c)
                    cell.text = cell_text
                    for p in cell.text_frame.paragraphs:
                        p.font.bold = True
                        p.font.size = Pt(16)
                        p.font.color.rgb = RGBColor(255, 255, 255)
                        p.alignment = PP_ALIGN.CENTER
                    # 表头背景色
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = RGBColor(0, 82, 204)
                
                # 数据行
                for r_idx, row in enumerate(data_rows):
                    for c_idx, cell_text in enumerate(row):
                        if c_idx < num_cols:
                            cell = table.cell(r_idx + 1, c_idx)
                            cell.text = clean_markdown(cell_text)
                            for p in cell.text_frame.paragraphs:
                                p.font.size = Pt(14)
                                p.font.color.rgb = RGBColor(15, 25, 45)
                            # 交替行背景
                            if r_idx % 2 == 1:
                                cell.fill.solid()
                                cell.fill.fore_color.rgb = RGBColor(240, 244, 250)
                return
    
    # ---- 代码布局 ----
    if layout_type == "code" and slide_data and "code_text" in slide_data:
        code_text = slide_data["code_text"]
        code_lang = slide_data.get("code_lang", "")
        
        # 代码块背景
        bg_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(1.8),
            Inches(11.6), Inches(5.0)
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(40, 44, 52)
        bg_shape.line.fill.background()
        
        # 代码文本
        code_box = slide.shapes.add_textbox(
            Inches(1.2), Inches(2.0),
            Inches(10.8), Inches(4.5)
        )
        tf = code_box.text_frame
        tf.word_wrap = True
        
        for i, codeline in enumerate(code_text.split('\n')):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = codeline
            p.font.name = "Courier New"
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(200, 210, 220)
            p.space_after = Pt(2)
        
        # 语言标签
        if code_lang:
            lang_label = slide.shapes.add_textbox(
                Inches(9.5), Inches(1.85),
                Inches(2.8), Inches(0.3)
            )
            lp = lang_label.text_frame.paragraphs[0]
            lp.text = code_lang
            lp.font.size = Pt(10)
            lp.font.color.rgb = RGBColor(150, 160, 170)
            lp.font.italic = True
            lp.alignment = PP_ALIGN.RIGHT
        return

    # ---- 其他布局 ----
    body_shape = slide.placeholders[1]
    
    if layout_type == "image_right":
        body_shape.width = Inches(7.5)
        if prs:
            add_placeholder_image(slide, prs)
    elif layout_type == "three_columns":
        body_shape.text_frame.clear()
        
        # 按照level=0的项目进行逻辑分组
        columns = []
        current_col = []
        for item in items:
            if item.get("level", 0) == 0 and current_col:
                columns.append(current_col)
                current_col = [item]
            else:
                current_col.append(item)
        if current_col:
            columns.append(current_col)
        
        # 如果第1列只有一个概述项，将其合并到标题中，其余列作为内容列
        if len(columns) > 3 and len(columns[0]) == 1:
            overview_text = columns[0][0].get("text", "")
            columns = columns[1:]  # 移除概述列
            # 更新标题，添加概述
            title_shape = slide.shapes.title
            title_shape.text = f"{clean_markdown(title_text)}\n{overview_text}"
        
        # 限制最多显示3列
        display_columns = columns[:3]
        
        # 计算列宽和间距
        num_cols = len(display_columns)
        total_width = Inches(11.333)
        col_width = Inches(3.2)
        col_spacing = (total_width - col_width * num_cols) / (num_cols + 1)
        
        # 显示每一列
        for col_idx, col_items in enumerate(display_columns):
            left_pos = Inches(1) + col_spacing * (col_idx + 1) + col_width * col_idx
            
            # 创建文本框
            txBox = slide.shapes.add_textbox(left_pos, Inches(2.0), col_width, Inches(4.5))
            tf = txBox.text_frame
            tf.word_wrap = True
            
            # 添加列标题（第一个 item）
            p = tf.paragraphs[0]
            _apply_formatted_paragraph(p, col_items[0].get("text", ""))
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 82, 204)
            p.alignment = PP_ALIGN.CENTER
            p.space_after = Pt(12)
            
            # 添加该列的其他 items
            for item in col_items[1:]:
                p = tf.add_paragraph()
                _apply_formatted_paragraph(p, item.get("text", ""))
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(60, 64, 67)
                p.space_before = Pt(8)
            
            # 添加列图标
            icon = slide.shapes.add_shape(MSO_SHAPE.OVAL, left_pos + (col_width - Inches(1)) / 2, Inches(1.2), Inches(1.0), Inches(1.0))
            icon.fill.solid()
            icon.fill.fore_color.rgb = RGBColor(235, 240, 248)
            icon.line.fill.background()
        return
    else:
        body_shape.width = Inches(11.333)
        
    body_shape.left = Inches(1)
    body_shape.top = Inches(1.8)
    body_shape.height = Inches(5.0)
    body_shape.text_frame.word_wrap = True  # 启用自动换行
    body_shape.text_frame.clear()
    
    for item in items:
        level = item.get("level", 0)
        text = item.get("text", "")
        if not text:
            continue
            
        p = body_shape.text_frame.paragraphs[0] if not body_shape.text_frame.text.strip() else body_shape.text_frame.add_paragraph()
        _apply_formatted_paragraph(p, text)
        p.level = level
        p.space_after = Pt(8)
        
        if level == 0:
            p.font.size = Pt(20)  # 缩小字体
            p.font.color.rgb = RGBColor(15, 25, 45)
        else:
            p.font.size = Pt(16)  # 缩小字体
            p.font.color.rgb = RGBColor(60, 64, 67)


def _repair_json_string(s: str) -> str:
    """修复 AI 生成 JSON 中常见的语法错误（预处理 + 迭代式，最多 20 轮）

    策略：
    0. 预处理：全局清除尾随逗号（,} → }、,] → ]）—— 用正则一次完成，避免位置偏移
    1. 迭代：解析 → 失败 → 定位错误位置 → 针对性修复 → 重试
    处理的问题：
    - 缺少冒号分隔符: "level:1" → "level":1
    - 对象/数组之间缺少逗号: }{"key" → },{"key"
    - 字符串值之间缺少逗号: "val""key" → "val","key"
    - 单引号替换为双引号
    """
    import re as _re

    # ---- 预处理阶段：全局清理 ----
    # 0a. 将单引号键/值替换为双引号（AI 有时输出 Python 风格的 dict）
    # 匹配 'key': 或 :'value' 模式（不处理字符串内部的单引号）
    s = _re.sub(r"(?<=[\[{,])\s*'([^']*?)'\s*:", r' "\1":', s)

    # 0b. 全局清除尾随逗号：,} → } 和 ,] → ]（含空白）
    s = _re.sub(r',(\s*[}\]])', r'\1', s)

    # ---- 迭代修复阶段 ----
    for attempt in range(20):
        try:
            json.loads(s)
            return s  # 解析成功
        except json.JSONDecodeError as e:
            pos = e.pos
            msg = str(e)

            if pos >= len(s):
                break

            fixed = False

            # --- 策略 1: 缺少冒号 "key:value" → "key":value ---
            fix_result = _try_fix_missing_colon(s, pos)
            if fix_result[0] is not None:
                fix_start, fix_end, fix_repl = fix_result
                s = s[:fix_start] + fix_repl + s[fix_end:]
                fixed = True

            if not fixed:
                # --- 策略 2: 对象/数组间缺少分隔符 ---
                # }{ },[ ]{ ][
                if pos > 0:
                    prev = s[pos - 1]
                    curr = s[pos]
                    if prev in '}]' and curr in '{[':
                        s = s[:pos] + ',' + s[pos:]
                        fixed = True

            if not fixed:
                # --- 策略 3: "value""key" → "value","key" ---
                if pos > 0 and s[pos] == '"' and s[pos - 1] == '"':
                    s = s[:pos] + ',' + s[pos:]
                    fixed = True

            if not fixed:
                # --- 策略 4: 再次尝试全局清除尾随逗号（可能在修复过程中产生新的） ---
                new_s = _re.sub(r',(\s*[}\]])', r'\1', s)
                if new_s != s:
                    s = new_s
                    fixed = True

            if not fixed:
                # --- 策略 5: 通用——在错误位置前插入逗号 ---
                if pos > 0 and s[pos - 1] not in ',{[:':
                    s = s[:pos] + ',' + s[pos:]
                    fixed = True

            if not fixed:
                # --- 策略 6: 跳过当前字符（不可修复的非法字符）---
                # 尝试删除错误位置的单个字符后是否能继续
                if pos < len(s):
                    s = s[:pos] + s[pos + 1:]
                    fixed = True

            if not fixed:
                logger.warning(f"[repair_json] 第{attempt+1}轮无法修复: pos={pos}, msg={msg}")
                break

            logger.debug(f"[repair_json] 第{attempt+1}轮修复 pos={pos}")

    return s


def _try_fix_missing_colon(s: str, error_pos: int):
    """尝试修复 "word:value" 缺少冒号的问题

    从错误位置向前扫描，找到 "word: 模式，返回 (start, end, replacement)
    """
    import re
    # 向前搜索最近的 "word: 模式（在 error_pos 前 50 字符范围内）
    search_start = max(0, error_pos - 50)
    region = s[search_start:error_pos + 30]

    # 匹配: "word: 其中 word 是 \w+
    matches = list(re.finditer(r'"(\w+):', region))
    if not matches:
        return (None, None, None)

    # 取最后一个匹配（最接近错误位置的）
    m = matches[-1]
    abs_start = search_start + m.start()
    abs_end = search_start + m.end()

    # 验证：确认这个 "word: 不是在字符串值内部
    # word 前面应该是 { , [ 或空白
    if abs_start > 0:
        prev_char = s[abs_start - 1]
        if prev_char not in '{[, \t\n\r':
            return (None, None, None)

    # 构造修复: "word: → "word":
    replacement = f'"{m.group(1)}":'
    return (abs_start, abs_end, replacement)


def extract_json_from_text(text):
    """从大模型的混合输出中提取 JSON 数据，或降级解析 Markdown 为 PPT 大纲"""
    # 记录原始输入长度
    logger.info(f"[extract_json] 输入长度: {len(text)} 字符")
    
    # 0. 清理中文引号（AI有时会返回中文引号，导致JSON解析失败）
    # 注意：只替换作为JSON结构元素的中文引号，不要破坏字符串内部的内容
    original_text = text
    
    # 替换中文双引号为英文双引号（仅在JSON结构位置）
    # 匹配模式：JSON键值对的引号（前面是{,或:，后面是}或:或,）
    text = re.sub(r'(?<=[\{,:])\s*[\u201c\u201d\u300c\u300d\u300e\u300f]\s*', '"', text)
    text = re.sub(r'\s*[\u201c\u201d\u300c\u300d\u300e\u300f]\s*(?=[\},:])', '"', text)
    
    # 替换中文单引号为英文单引号（仅在字符串内部）
    text = text.replace('\u2018', "'").replace('\u2019', "'")
    
    # 记录清理后的长度
    if len(text) != len(original_text):
        logger.info(f"[extract_json] 清理中文引号后长度: {len(text)} 字符 (减少 {len(original_text) - len(text)} 字符)")
    
    # 1. 尝试匹配 ```json ... ``` 块（支持对象 {} 或数组 []）
    match = re.search(r'```(?:json)?\s*(\{.*?\}|\[.*?\])\s*```', text, re.DOTALL)
    json_str = match.group(1) if match else text
    
    if match:
        logger.info(f"[extract_json] 找到 ```json 代码块，长度: {len(json_str)} 字符")
    else:
        logger.info(f"[extract_json] 未找到 ```json 代码块，使用完整文本")

    # 2. 尝试直接解析 JSON
    clean_str = json_str  # 默认使用完整文本，后面可能被截取
    try:
        start_obj = json_str.find('{')
        start_arr = json_str.find('[')
        start_idx = -1
        
        if start_obj != -1 and start_arr != -1:
            start_idx = min(start_obj, start_arr)
        else:
            start_idx = max(start_obj, start_arr)
            
        if start_idx != -1:
            end_idx = -1
            if json_str[start_idx] == '{':
                end_idx = json_str.rfind('}') + 1
            else:
                end_idx = json_str.rfind(']') + 1
                
            if end_idx > start_idx:
                clean_str = json_str[start_idx:end_idx]
                logger.info(f"[extract_json] 提取JSON字符串，长度: {len(clean_str)} 字符")
                
                # 记录JSON字符串的前200字符和后200字符
                logger.info(f"[extract_json] JSON前200字符: {clean_str[:200]}")
                logger.info(f"[extract_json] JSON后200字符: {clean_str[-200:]}")
                
                parsed = json.loads(clean_str)
                
                # 兼容两种格式：
                # 1. {"title": "...", "slides": [...]} - 完整格式，直接返回
                # 2. {"slides": [...]} - 部分格式，返回完整对象
                # 3. [...] - 数组格式，直接返回
                if isinstance(parsed, dict) and "slides" in parsed:
                    logger.info(f"[extract_json] JSON解析成功，slides数量: {len(parsed.get('slides', []))}")
                    return parsed  # 返回完整字典，保留 title 等字段
                elif isinstance(parsed, list):
                    logger.info(f"[extract_json] JSON解析成功，数组长度: {len(parsed)}")
                    return parsed
    except json.JSONDecodeError as e:
        logger.error(f"[extract_json] JSON解析失败: {e}")
        logger.error(f"[extract_json] 错误位置: {e.pos}")
        logger.error(f"[extract_json] 错误行: {e.lineno}, 错误列: {e.colno}")
        # 记录错误位置附近的内容
        if e.pos < len(json_str):
            start = max(0, e.pos - 100)
            end = min(len(json_str), e.pos + 100)
            logger.error(f"[extract_json] 错误位置附近内容: {json_str[start:end]}")
        
        # 2.0.1 尝试迭代修复 AI 常见的 JSON 语法错误
        try:
            repaired = _repair_json_string(clean_str)
            if repaired != clean_str:
                logger.info("[extract_json] 迭代修复完成，尝试重新解析...")
                print(f"[extract_json] 尝试迭代修复 JSON 语法错误 (修复后长度={len(repaired)})")
                parsed = json.loads(repaired)
                if isinstance(parsed, dict) and "slides" in parsed:
                    logger.info(f"[extract_json] 修复后解析成功！slides数量: {len(parsed.get('slides', []))}")
                    print(f"[extract_json] ✅ 自动修复 JSON 语法错误，解析出 {len(parsed.get('slides', []))} 个 slides")
                    return parsed
                elif isinstance(parsed, list):
                    logger.info(f"[extract_json] 修复后解析成功，数组长度: {len(parsed)}")
                    return parsed
        except (json.JSONDecodeError, Exception) as repair_e:
            logger.error(f"[extract_json] 迭代修复后仍然失败: {repair_e}")
            print(f"[extract_json] ❌ 迭代修复后仍然失败: {repair_e}")
            # 输出修复后仍然出错的位置附近内容
            if isinstance(repair_e, json.JSONDecodeError) and repair_e.pos < len(repaired):
                ctx_start = max(0, repair_e.pos - 80)
                ctx_end = min(len(repaired), repair_e.pos + 80)
                print(f"[extract_json] 修复后错误位置附近: ...{repaired[ctx_start:ctx_end]}...")
    except Exception as e:
        logger.error(f"[extract_json] JSON解析异常: {type(e).__name__}: {e}")
    
    # 2.0.5 json_repair 专业库兜底（在截断修复前尝试，处理复杂组合错误）
    try:
        import json_repair
        repaired_obj = json_repair.repair_json(json_str, return_objects=True)
        if isinstance(repaired_obj, dict) and "slides" in repaired_obj:
            logger.info(f"[extract_json] json_repair 修复成功，slides数量: {len(repaired_obj.get('slides', []))}")
            print(f"[extract_json] ✅ json_repair 自动修复 JSON，解析出 {len(repaired_obj.get('slides', []))} 个 slides")
            return repaired_obj
        elif isinstance(repaired_obj, list) and len(repaired_obj) > 0:
            logger.info(f"[extract_json] json_repair 修复成功，数组长度: {len(repaired_obj)}")
            return repaired_obj
    except ImportError:
        pass
    except Exception as jr_e:
        logger.warning(f"[extract_json] json_repair 修复失败: {jr_e}")

    # 2.1 尝试修复被截断的 JSON
    try:
        # 检查是否是被截断的 JSON（包含 slides 但不完整）
        if '"slides"' in json_str and '"slides":[' in json_str:
            logger.info("[extract_json] 检测到可能被截断的JSON，尝试修复")
            # 找到 slides 数组的开始
            slides_start = json_str.find('"slides":[')
            if slides_start > 0:
                # 提取 slides 数组内容（跳过 "slides":[ 部分）
                slides_content = json_str[slides_start + 10:]  # 跳过 "slides":[
                
                # 找到最后一个完整的 slide 对象
                # 通过匹配花括号来找到完整的对象
                brace_count = 0
                last_complete_slide_end = -1
                in_slide = False
                
                for i, char in enumerate(slides_content):
                    if char == '{':
                        if not in_slide:
                            in_slide = True
                        brace_count += 1
                    elif char == '}':
                        brace_count -= 1
                        if brace_count == 0 and in_slide:
                            last_complete_slide_end = i
                            in_slide = False
                
                if last_complete_slide_end > 0:
                    # 截取到最后一个完整的 slide
                    slides_content = slides_content[:last_complete_slide_end + 1]
                    
                    # 构建完整的 JSON
                    # 提取 title 字段
                    title_match = re.search(r'"title"\s*:\s*"([^"]*)"', json_str[:slides_start])
                    if title_match:
                        fixed_json = '{"title":"' + title_match.group(1) + '","slides":[' + slides_content + ']}'
                    else:
                        fixed_json = '{"slides":[' + slides_content + ']}'
                    
                    try:
                        parsed = json.loads(fixed_json)
                        if isinstance(parsed, dict) and "slides" in parsed:
                            logger.info(f"[extract_json] 修复被截断的 JSON，保留 {len(parsed['slides'])} 个 slides")
                            print(f"[ppt_generator] 修复被截断的 JSON，保留 {len(parsed['slides'])} 个 slides")
                            return parsed
                    except json.JSONDecodeError as e:
                        logger.error(f"[extract_json] 修复后的JSON解析失败: {e}")
                        # 截断修复后的 JSON 可能仍有语法错误，尝试迭代修复
                        try:
                            repaired = _repair_json_string(fixed_json)
                            parsed = json.loads(repaired)
                            if isinstance(parsed, dict) and "slides" in parsed:
                                logger.info(f"[extract_json] 截断+迭代修复 JSON，保留 {len(parsed['slides'])} 个 slides")
                                print(f"[ppt_generator] 截断+迭代修复 JSON，保留 {len(parsed['slides'])} 个 slides")
                                return parsed
                        except (json.JSONDecodeError, Exception) as e2:
                            logger.error(f"[extract_json] 截断+迭代修复仍失败: {e2}")
                            pass
                else:
                    logger.warning("[extract_json] 未找到完整的slide对象")
            else:
                logger.warning("[extract_json] 未找到slides数组开始位置")
        else:
            logger.info("[extract_json] 未检测到被截断的JSON特征")
    except Exception as e:
        logger.error(f"[extract_json] 修复截断JSON异常: {type(e).__name__}: {e}")
        
    # 3. 如果 JSON 解析失败，但包含 Markdown 标题，则降级将 Markdown 转换为 JSON 结构
    if '# ' in text or '## ' in text:
        logger.info("[extract_json] JSON解析失败，尝试Markdown降级处理")
        slides = []
        current_slide = None
        table_buffer = []       # 累积表格行
        in_table = False        # 是否在表格中
        in_code_block = False   # 是否在代码块中
        code_buffer = []        # 累积代码行
        code_lang = ""          # 代码语言

        def flush_table():
            """将累积的表格行写入当前 slide"""
            nonlocal table_buffer, in_table
            if current_slide and table_buffer:
                current_slide["layout"] = "table"
                current_slide["table_data"] = table_buffer.copy()
                current_slide["items"] = current_slide.get("items", [])
            table_buffer = []
            in_table = False

        def flush_code():
            """将累积的代码块写入当前 slide"""
            nonlocal code_buffer, code_lang, in_code_block
            if current_slide and code_buffer:
                current_slide["layout"] = "code"
                current_slide["code_lang"] = code_lang
                current_slide["code_text"] = "\n".join(code_buffer)
                current_slide["items"] = current_slide.get("items", [])
            code_buffer = []
            code_lang = ""
            in_code_block = False

        for line in text.split('\n'):
            stripped = line.strip()
            
            # 代码块检测
            if stripped.startswith('```'):
                if in_code_block:
                    flush_code()
                else:
                    flush_table()  # 先刷新可能的表格
                    code_lang = stripped[3:].strip()
                    in_code_block = True
                continue
            if in_code_block:
                code_buffer.append(line)
                continue
            
            # 表格检测（包含 | 且下一行是 |---| 分隔符）
            if '|' in stripped and stripped.strip('|').strip():
                if not in_table:
                    flush_code()  # 刷新可能的代码块
                    # 检查是否真的是表格（需要后续行也存在）
                    in_table = True
                table_buffer.append(stripped)
                continue
            elif in_table and stripped.strip():
                # 可能是表格的 continuation
                if '|' in stripped:
                    table_buffer.append(stripped)
                    continue
                else:
                    flush_table()
            
            # 空行处理
            if not stripped:
                flush_table()
                flush_code()
                continue
            
            flush_table()
            flush_code()
            
            if stripped.startswith('# '):
                if current_slide: slides.append(current_slide)
                current_slide = {"type": "title", "layout": "center", "title": stripped[2:].strip(), "subtitle": "", "items": []}
            elif stripped.startswith('## '):
                if current_slide: slides.append(current_slide)
                current_slide = {"type": "content", "layout": "text_only", "title": stripped[3:].strip(), "items": []}
            elif stripped.startswith('- ') or stripped.startswith('* '):
                if current_slide:
                    if "items" not in current_slide: current_slide["items"] = []
                    current_slide["items"].append({"level": 0, "text": stripped[2:].strip()})
            else:
                if current_slide and current_slide.get("type") == "content":
                    if "items" not in current_slide: current_slide["items"] = []
                    current_slide["items"].append({"level": 0, "text": stripped})

        # 刷新残留缓冲区
        flush_table()
        flush_code()
        if current_slide:
            slides.append(current_slide)
        if slides:
            logger.info(f"[extract_json] Markdown降级处理成功，生成 {len(slides)} 个slides")
            return slides
        else:
            logger.warning("[extract_json] Markdown降级处理未生成任何slides")
            
    logger.error("[extract_json] 所有解析方式均失败，返回None")
    return None

def format_chart_slide(slide, slide_data, prs):
    """格式化图表页"""
    from pptx.chart.data import CategoryChartData
    from pptx.enum.chart import XL_CHART_TYPE
    
    title_text = slide_data.get("title", "图表")
    chart_type_str = slide_data.get("chart_type", "bar")
    chart_data = slide_data.get("chart_data", {})
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = clean_markdown(title_text)
    title_shape.width = Inches(11.333)
    title_shape.left = Inches(1)
    title_shape.top = Inches(0.5)
    title_shape.height = Inches(1.2)
    
    for p in title_shape.text_frame.paragraphs:
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = RGBColor(0, 82, 204)
        p.alignment = PP_ALIGN.LEFT
    
    # 准备图表数据
    labels = chart_data.get("labels", [])
    datasets = chart_data.get("datasets", [])
    
    if not labels or not datasets:
        # 如果没有图表数据，显示占位符
        body_shape = slide.placeholders[1]
        body_shape.text = "图表数据加载中..."
        return
    
    # 映射图表类型
    chart_type_map = {
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
    }
    xl_chart_type = chart_type_map.get(chart_type_str, XL_CHART_TYPE.BAR_CLUSTERED)
    
    # 创建图表数据
    chart_data_obj = CategoryChartData()
    chart_data_obj.categories = labels
    for dataset in datasets:
        chart_data_obj.add_series(dataset.get("label", ""), dataset.get("data", []))
    
    # 添加图表到幻灯片
    chart_left = Inches(1)
    chart_top = Inches(1.8)
    chart_width = Inches(11.333)
    chart_height = Inches(5.0)
    
    chart_frame = slide.shapes.add_chart(
        xl_chart_type, chart_left, chart_top, chart_width, chart_height, chart_data_obj
    )
    
    # 设置图表样式
    chart = chart_frame.chart
    chart.has_legend = True
    chart.legend.include_in_layout = False
    
    # 设置系列颜色
    colors = [RGBColor(220, 53, 69), RGBColor(40, 167, 69), RGBColor(0, 123, 255), RGBColor(255, 193, 7)]
    for i, series in enumerate(chart.series):
        if i < len(datasets) and "color" in datasets[i]:
            color_hex = datasets[i]["color"].lstrip('#')
            r, g, b = int(color_hex[0:2], 16), int(color_hex[2:4], 16), int(color_hex[4:6], 16)
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = RGBColor(r, g, b)
        elif i < len(colors):
            series.format.fill.solid()
            series.format.fill.fore_color.rgb = colors[i]


def format_flowchart_slide(slide, slide_data, prs):
    """格式化流程图页 - 使用形状和连接线绘制流程图"""
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    
    title_text = slide_data.get("title", "流程图")
    flowchart_data = slide_data.get("flowchart_data", {})
    
    # 设置标题
    title_shape = slide.shapes.title
    title_shape.text = clean_markdown(title_text)
    title_shape.width = Inches(11.333)
    title_shape.left = Inches(1)
    title_shape.top = Inches(0.5)
    title_shape.height = Inches(1.2)
    
    for p in title_shape.text_frame.paragraphs:
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = RGBColor(0, 82, 204)
        p.alignment = PP_ALIGN.LEFT
    
    nodes = flowchart_data.get("nodes", [])
    edges = flowchart_data.get("edges", [])
    
    if not nodes:
        # 如果没有节点数据，显示占位符
        body_shape = slide.placeholders[1]
        body_shape.text = "流程图数据加载中..."
        return
    
    # 计算布局参数
    num_nodes = len(nodes)
    start_left = Inches(1.5)
    start_top = Inches(2.5)
    node_width = Inches(1.8)
    node_height = Inches(0.8)
    h_spacing = Inches(0.5)  # 水平间距
    v_spacing = Inches(0.3)  # 垂直间距
    
    # 根据节点数量决定布局方向
    if num_nodes <= 4:
        # 水平布局
        total_width = num_nodes * node_width + (num_nodes - 1) * h_spacing
        start_left = (Inches(13.333) - total_width) / 2
        
        node_positions = {}
        for i, node in enumerate(nodes):
            left = start_left + i * (node_width + h_spacing)
            top = start_top
            
            # 根据形状类型选择形状
            shape_type = node.get("shape", "process")
            if shape_type == "start" or shape_type == "end":
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, node_width, node_height)
                shape.fill.solid()
                # 使用更柔和的颜色
                shape.fill.fore_color.rgb = RGBColor(46, 139, 87) if shape_type == "start" else RGBColor(178, 34, 34)
                shape.line.color.rgb = RGBColor(34, 100, 60) if shape_type == "start" else RGBColor(130, 25, 25)
            else:
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, node_width, node_height)
                shape.fill.solid()
                # 使用渐变蓝色系
                shape.fill.fore_color.rgb = RGBColor(65, 105, 225)
                shape.line.color.rgb = RGBColor(50, 80, 180)
            
            shape.line.width = Pt(1.5)
            
            # 设置文本
            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(8)
            tf.margin_right = Pt(8)
            p = tf.paragraphs[0]
            p.text = node.get("text", "")
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            node_positions[node["id"]] = {
                "left": left,
                "top": top,
                "right": left + node_width,
                "bottom": top + node_height,
                "center_x": left + node_width / 2,
                "center_y": top + node_height / 2
            }
    else:
        # 垂直布局（节点较多时）
        cols = min(3, (num_nodes + 1) // 2)
        rows = (num_nodes + cols - 1) // cols
        
        total_width = cols * node_width + (cols - 1) * h_spacing
        start_left = (Inches(13.333) - total_width) / 2
        
        node_positions = {}
        for i, node in enumerate(nodes):
            col = i % cols
            row = i // cols
            left = start_left + col * (node_width + h_spacing)
            top = start_top + row * (node_height + v_spacing)
            
            # 根据形状类型选择形状
            shape_type = node.get("shape", "process")
            if shape_type == "start" or shape_type == "end":
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, node_width, node_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(46, 139, 87) if shape_type == "start" else RGBColor(178, 34, 34)
                shape.line.color.rgb = RGBColor(34, 100, 60) if shape_type == "start" else RGBColor(130, 25, 25)
            else:
                shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, node_width, node_height)
                shape.fill.solid()
                shape.fill.fore_color.rgb = RGBColor(65, 105, 225)
                shape.line.color.rgb = RGBColor(50, 80, 180)
            
            shape.line.width = Pt(1.5)
            
            # 设置文本
            tf = shape.text_frame
            tf.word_wrap = True
            tf.margin_left = Pt(8)
            tf.margin_right = Pt(8)
            p = tf.paragraphs[0]
            p.text = node.get("text", "")
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = PP_ALIGN.CENTER
            
            node_positions[node["id"]] = {
                "left": left,
                "top": top,
                "right": left + node_width,
                "bottom": top + node_height,
                "center_x": left + node_width / 2,
                "center_y": top + node_height / 2
            }
    
    # 绘制连接线
    from pptx.enum.shapes import MSO_CONNECTOR_TYPE
    from pptx.oxml.ns import qn
    
    for edge in edges:
        from_id = edge.get("from")
        to_id = edge.get("to")
        
        if from_id in node_positions and to_id in node_positions:
            from_pos = node_positions[from_id]
            to_pos = node_positions[to_id]
            
            # 计算连接线起点和终点
            # 水平布局：从右侧中点连接到左侧中点
            # 垂直布局：从底部中点连接到顶部中点
            if num_nodes <= 4:
                # 水平布局
                start_x = from_pos["right"]
                start_y = from_pos["center_y"]
                end_x = to_pos["left"]
                end_y = to_pos["center_y"]
                
                # 使用肘形连接线（带角度），上下偏移避免重叠
                offset_y = Inches(0.3) if from_pos["center_y"] < Inches(4) else Inches(-0.3)
                
                # 添加肘形连接线
                connector = slide.shapes.add_connector(
                    MSO_CONNECTOR_TYPE.ELBOW,
                    start_x, start_y,
                    end_x, end_y
                )
                connector.line.color.rgb = RGBColor(70, 130, 180)  # 钢蓝色
                connector.line.width = Pt(2)
                
                # 通过 XML 调整连接线路径，添加中间转折点
                # 创建自定义路径使连接线带角度
                sp_elem = connector._element
                cxnSp = sp_elem
                
                # 添加中间点使连接线呈现 Z 字形
                # python-pptx 的 connector 会自动根据起点终点计算肘形路径
            else:
                # 垂直布局
                start_x = from_pos["center_x"]
                start_y = from_pos["bottom"]
                end_x = to_pos["center_x"]
                end_y = to_pos["top"]
                
                # 使用肘形连接线
                connector = slide.shapes.add_connector(
                    MSO_CONNECTOR_TYPE.ELBOW,
                    start_x, start_y,
                    end_x, end_y
                )
                connector.line.color.rgb = RGBColor(70, 130, 180)  # 钢蓝色
                connector.line.width = Pt(2)
            
            # 设置末端箭头样式
            line_elem = connector.line._ln
            tailEnd = line_elem.makeelement(qn('a:tailEnd'), {})
            tailEnd.set('type', 'triangle')
            tailEnd.set('w', 'med')
            tailEnd.set('len', 'med')
            line_elem.append(tailEnd)


def generate_ppt_from_json(json_data, output_path="output.pptx"):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for slide_data in json_data:
        slide_type = slide_data.get("type", "content")
        content_type = slide_data.get("content_type", "text")
        
        if slide_type == "title":
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            apply_corporate_theme(slide, prs, is_title_slide=True)
            format_title_slide(slide, slide_data.get("title", "演示文稿"), slide_data.get("subtitle", ""))
            
        elif slide_type == "ending":
            # 结尾页：复用 title 布局，居中展示 "谢谢" + "Q & A"
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            apply_corporate_theme(slide, prs, is_title_slide=True)
            format_title_slide(slide, slide_data.get("title", "谢谢"), slide_data.get("subtitle", "Q & A"))
            
        elif slide_type == "content":
            # 检查是否是图表类型
            if content_type == "chart" and "chart_data" in slide_data:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                apply_corporate_theme(slide, prs, is_title_slide=False)
                format_chart_slide(slide, slide_data, prs)
            # 检查是否是流程图类型
            elif content_type == "flowchart" and "flowchart_data" in slide_data:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                apply_corporate_theme(slide, prs, is_title_slide=False)
                format_flowchart_slide(slide, slide_data, prs)
            # 检查 items 中是否包含流程图数据
            elif any(item.get("content_type") == "flowchart" for item in slide_data.get("items", [])):
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                apply_corporate_theme(slide, prs, is_title_slide=False)
                # 提取第一个流程图 item
                flowchart_item = next(item for item in slide_data.get("items", []) if item.get("content_type") == "flowchart")
                flowchart_slide_data = {
                    "title": slide_data.get("title", "流程图"),
                    "flowchart_data": flowchart_item.get("flowchart_data", {})
                }
                format_flowchart_slide(slide, flowchart_slide_data, prs)
            else:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                apply_corporate_theme(slide, prs, is_title_slide=False)
                layout_type = slide_data.get("layout", "text_only")
                format_content_slide(slide, slide_data.get("title", "内容"), slide_data.get("items", []), layout_type, prs, slide_data)
            
    if len(prs.slides) == 0:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        apply_corporate_theme(slide, prs, is_title_slide=True)
        format_title_slide(slide, "自动生成的幻灯片")
        
    prs.save(output_path)
    return os.path.abspath(output_path)

def generate_ppt_from_text(text, output_path="output.pptx"):
    """入口函数，兼容 JSON 提取与旧版降级处理"""
    json_data = extract_json_from_text(text)
    
    if json_data and isinstance(json_data, list):
        return generate_ppt_from_json(json_data, output_path)
    
    # 如果没提取到 JSON，说明大模型没有按规范输出（或者走的是旧流程），抛出异常让上层处理
    raise ValueError("无法从内容中提取出有效的幻灯片 JSON 结构。")

if __name__ == "__main__":
    pass