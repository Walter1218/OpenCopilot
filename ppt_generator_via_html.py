#!/usr/bin/env python3
"""
PPT生成器 - 先生成HTML预览，再转换为PPT
"""

import json
import os
import sys
import re

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


def clean_markdown(text):
    """清理Markdown标记"""
    if not text:
        return ""
    text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
    text = re.sub(r'\*(.*?)\*', r'\1', text)
    text = re.sub(r'`(.*?)`', r'\1', text)
    text = re.sub(r'#{1,6}\s+', '', text)
    return text.strip()


def generate_ppt_via_html(slides_data, output_path="output.pptx"):
    """
    先生成HTML预览，再转换为PPT
    
    Args:
        slides_data: 幻灯片数据列表
        output_path: 输出PPT文件路径
    
    Returns:
        生成的PPT文件路径
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for i, slide_data in enumerate(slides_data):
        slide_type = slide_data.get("type", "content")
        title = slide_data.get("title", "")
        subtitle = slide_data.get("subtitle", "")
        items = slide_data.get("items", [])
        layout = slide_data.get("layout", "text_only")
        
        # 创建幻灯片
        if slide_type == "title":
            slide_layout = prs.slide_layouts[0]
            slide_obj = prs.slides.add_slide(slide_layout)
            _format_title_slide(slide_obj, title, subtitle)
        elif slide_type == "ending":
            slide_layout = prs.slide_layouts[0]
            slide_obj = prs.slides.add_slide(slide_layout)
            _format_ending_slide(slide_obj, title, subtitle)
        else:
            slide_layout = prs.slide_layouts[1]
            slide_obj = prs.slides.add_slide(slide_layout)
            _format_content_slide(slide_obj, title, items, layout, i+1, len(slides_data))
    
    prs.save(output_path)
    return os.path.abspath(output_path)


def _format_title_slide(slide, title, subtitle):
    """格式化标题页"""
    # 设置背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 82, 204)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = clean_markdown(title)
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.8))
        tf = subtitle_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = clean_markdown(subtitle)
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER


def _format_ending_slide(slide, title, subtitle):
    """格式化结尾页"""
    # 设置背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 26, 46)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = clean_markdown(title)
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.8))
        tf = subtitle_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = clean_markdown(subtitle)
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER


def _format_content_slide(slide, title, items, layout, page_num, total_pages):
    """格式化内容页"""
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = clean_markdown(title)
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 82, 204)
    p.alignment = PP_ALIGN.LEFT
    
    # 内容区域
    if layout == "three_columns":
        _format_three_columns(slide, items, title_box)
    else:
        _format_text_only(slide, items)
    
    # 页码
    page_box = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3))
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{page_num} / {total_pages}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(153, 153, 153)
    p.alignment = PP_ALIGN.RIGHT


def _format_three_columns(slide, items, title_box):
    """格式化三列布局"""
    # 按level=0的项目进行逻辑分组
    columns = []
    current_col = []
    for item in items:
        if item.get("level", 0) == 0 and current_col:
            columns.append(current_col)
            current_col = [item]
        else:
            current_col.append(item)
    if current_col:
        columns.append(current_col)
    
    # 如果第1列只有一个概述项，将其合并到标题中
    if len(columns) > 3 and len(columns[0]) == 1:
        overview_text = columns[0][0].get("text", "")
        # 更新标题
        title_box.text_frame.paragraphs[0].text = f"{title_box.text_frame.paragraphs[0].text}\n{overview_text}"
        columns = columns[1:]
    
    # 限制最多3列
    display_columns = columns[:3]
    
    # 计算列宽和间距
    num_cols = len(display_columns)
    total_width = Inches(11.333)
    col_width = Inches(3.2)
    col_spacing = (total_width - col_width * num_cols) / (num_cols + 1)
    
    # 显示每一列
    for col_idx, col_items in enumerate(display_columns):
        left_pos = Inches(1) + col_spacing * (col_idx + 1) + col_width * col_idx
        
        # 创建列容器
        col_box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left_pos, Inches(2.0), col_width, Inches(4.5)
        )
        col_box.fill.solid()
        col_box.fill.fore_color.rgb = RGBColor(248, 249, 250)
        col_box.line.fill.background()
        
        # 列标题
        col_title_box = slide.shapes.add_textbox(left_pos + Inches(0.2), Inches(2.2), col_width - Inches(0.4), Inches(0.5))
        tf = col_title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = clean_markdown(col_items[0].get("text", ""))
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 82, 204)
        p.alignment = PP_ALIGN.CENTER
        
        # 列内容
        content_box = slide.shapes.add_textbox(left_pos + Inches(0.2), Inches(2.8), col_width - Inches(0.4), Inches(3.5))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for item_idx, item in enumerate(col_items[1:]):
            if item_idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = clean_markdown(item.get("text", ""))
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(60, 64, 67)
            p.space_before = Pt(8)


def _format_text_only(slide, items):
    """格式化纯文本布局"""
    content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(5.0))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for item_idx, item in enumerate(items):
        level = item.get("level", 0)
        text = item.get("text", "")
        
        if item_idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = clean_markdown(text)
        p.level = level
        p.space_after = Pt(8)
        
        if level == 0:
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(15, 25, 45)
        else:
            p.font.size = Pt(16)
            p.font.color.rgb = RGBColor(60, 64, 67)


def generate_html_preview(slides_data, output_path="preview.html"):
    """
    生成HTML预览文件
    
    Args:
        slides_data: 幻灯片数据列表
        output_path: 输出HTML文件路径
    
    Returns:
        生成的HTML文件路径
    """
    # 这里可以调用之前的HTML预览生成逻辑
    # 为了简化，暂时返回None
    return None


if __name__ == "__main__":
    # 测试用例
    test_slides = [
        {
            "type": "title",
            "title": "测试标题",
            "subtitle": "测试副标题"
        },
        {
            "type": "content",
            "title": "测试内容页",
            "layout": "text_only",
            "items": [
                {"level": 0, "text": "标题1"},
                {"level": 1, "text": "内容1"},
                {"level": 0, "text": "标题2"},
                {"level": 1, "text": "内容2"}
            ]
        }
    ]
    
    output_path = os.path.expanduser("~/Documents/trae_projects/OpenCopilot/output/test_via_html.pptx")
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    
    generate_ppt_via_html(test_slides, output_path)
    print(f"✅ 测试PPT已生成: {output_path}")
