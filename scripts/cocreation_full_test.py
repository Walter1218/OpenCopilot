#!/usr/bin/env python3
"""
使用完整的 AI Agent 深度研究报告生成 PPT
"""

import os
import sys
import json

# 添加项目根目录到路径
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from opencopilot.capabilities.ppt.pipeline import PPTGenerationPipeline
from ppt_generator import generate_ppt_from_json


def main():
    # 1. 读取完整文档
    doc_path = "/Users/onetwo/WorkBuddy/2026-05-28-21-09-15/AI_Agent_深度研究报告.md"
    
    with open(doc_path, 'r', encoding='utf-8') as f:
        content = f.read()
    
    print(f"📄 已读取文档: {len(content)} 字符")
    print(f"   文件: {doc_path}")
    
    # 2. 生成 PPT
    print("\n🔄 开始生成 PPT...")
    pipeline = PPTGenerationPipeline()
    result = pipeline.run(content)
    
    print(f"\n✅ PPT 生成完成: {result.total_pages} 页")
    print(f"   主题: {[t.title for t in result.topics]}")
    
    # 3. 查看各页内容
    print("\n📊 各页详情:")
    for i, slide in enumerate(result.slides):
        title = slide.get('title', 'N/A')
        content_type = slide.get('content_type', 'text')
        items_count = len(slide.get('items', []))
        print(f"   第 {i+1} 页: {title} ({content_type}, {items_count} 项)")
    
    # 4. 将第三页转换为图表
    if len(result.slides) >= 3:
        print("\n🔄 将第三页转换为图表...")
        result.slides[2]['content_type'] = 'chart'
        result.slides[2]['chart_type'] = 'bar'
        result.slides[2]['chart_data'] = {
            "title": "AI Agent 市场规模与增长",
            "labels": ["2024年", "2025年", "2026年(预测)", "2030年(预测)"],
            "datasets": [
                {"label": "全球市场(亿美元)", "data": [52.9, 120, 250, 471], "color": "#007bff"},
                {"label": "中国市场(亿人民币)", "data": [100, 190, 400, 800], "color": "#28a745"}
            ]
        }
        print("✅ 已转换为图表")
    
    # 5. 生成 PPT 文件
    output_dir = os.path.expanduser("~/Documents/trae_projects/OpenCopilot/output")
    os.makedirs(output_dir, exist_ok=True)
    
    output_path = os.path.join(output_dir, "AI_Agent_深度研究报告_完整版.pptx")
    print(f"\n📁 生成 PPT 文件: {output_path}")
    
    try:
        generate_ppt_from_json(result.slides, output_path)
        print(f"✅ PPT 已生成: {output_path}")
        
        # 验证生成结果
        from pptx import Presentation
        prs = Presentation(output_path)
        print(f"\n📊 最终 PPT 统计:")
        print(f"   总页数: {len(prs.slides)}")
        
        # 检查第三页是否有图表
        if len(prs.slides) >= 3:
            slide_3 = prs.slides[2]
            has_chart = any(shape.has_chart for shape in slide_3.shapes if hasattr(shape, 'has_chart'))
            print(f"   第三页是否有图表: {'是 ✅' if has_chart else '否 ❌'}")
        
        return output_path
        
    except Exception as e:
        print(f"⚠️ PPT 生成警告: {e}")
        import traceback
        traceback.print_exc()
        
        # 保存 JSON 作为备选
        json_path = os.path.join(output_dir, "AI_Agent_深度研究报告_完整版.json")
        with open(json_path, 'w', encoding='utf-8') as f:
            json.dump({"slides": result.slides}, f, ensure_ascii=False, indent=2)
        print(f"✅ JSON 已保存: {json_path}")
        return json_path


if __name__ == "__main__":
    output_path = main()
    print(f"\n" + "="*60)
    print(f"📁 文件位置: {output_path}")
    print("="*60)
