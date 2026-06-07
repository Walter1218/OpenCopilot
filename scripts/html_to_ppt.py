#!/usr/bin/env python3
"""
将HTML预览转换为PPT - 使用python-pptx的富文本功能
"""

import json
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import re


def clean_html(html_text):
    """清理HTML标签，保留纯文本"""
    clean = re.sub(r'<[^>]+>', '', html_text)
    return clean.strip()


def create_ppt_from_html_data(slides_data, output_path="output.pptx"):
    """根据HTML数据创建PPT"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    for i, slide in enumerate(slides_data):
        slide_type = slide.get("type", "content")
        title = slide.get("title", "")
        subtitle = slide.get("subtitle", "")
        items = slide.get("items", [])
        layout = slide.get("layout", "text_only")
        
        # 创建幻灯片
        if slide_type == "title":
            slide_layout = prs.slide_layouts[0]
            slide_obj = prs.slides.add_slide(slide_layout)
            format_title_slide_html(slide_obj, title, subtitle)
        elif slide_type == "ending":
            slide_layout = prs.slide_layouts[0]
            slide_obj = prs.slides.add_slide(slide_layout)
            format_ending_slide_html(slide_obj, title, subtitle)
        else:
            slide_layout = prs.slide_layouts[1]
            slide_obj = prs.slides.add_slide(slide_layout)
            format_content_slide_html(slide_obj, title, items, layout, i+1, len(slides_data))
    
    prs.save(output_path)
    return os.path.abspath(output_path)


def format_title_slide_html(slide, title, subtitle):
    """格式化标题页"""
    # 设置背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 82, 204)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.8))
        tf = subtitle_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER


def format_ending_slide_html(slide, title, subtitle):
    """格式化结尾页"""
    # 设置背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(26, 26, 46)
    
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # 副标题
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(4.2), Inches(11.333), Inches(0.8))
        tf = subtitle_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(28)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.alignment = PP_ALIGN.CENTER


def format_content_slide_html(slide, title, items, layout, page_num, total_pages):
    """格式化内容页"""
    # 标题
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(11.333), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 82, 204)
    p.alignment = PP_ALIGN.LEFT
    
    # 内容区域
    if layout == "three_columns":
        # 三列布局
        columns = []
        current_col = []
        for item in items:
            if item.get("level", 0) == 0 and current_col:
                columns.append(current_col)
                current_col = [item]
            else:
                current_col.append(item)
        if current_col:
            columns.append(current_col)
        
        # 如果第1列只有一个概述项，将其合并到标题中
        if len(columns) > 3 and len(columns[0]) == 1:
            overview_text = columns[0][0].get("text", "")
            # 更新标题
            title_box.text_frame.paragraphs[0].text = f"{title}\n{overview_text}"
            columns = columns[1:]
        
        # 限制最多3列
        display_columns = columns[:3]
        
        # 计算列宽和间距
        num_cols = len(display_columns)
        total_width = Inches(11.333)
        col_width = Inches(3.2)
        col_spacing = (total_width - col_width * num_cols) / (num_cols + 1)
        
        # 显示每一列
        for col_idx, col_items in enumerate(display_columns):
            left_pos = Inches(1) + col_spacing * (col_idx + 1) + col_width * col_idx
            
            # 创建列容器
            col_box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left_pos, Inches(2.0), col_width, Inches(4.5)
            )
            col_box.fill.solid()
            col_box.fill.fore_color.rgb = RGBColor(248, 249, 250)
            col_box.line.fill.background()
            
            # 列标题
            title_box = slide.shapes.add_textbox(left_pos + Inches(0.2), Inches(2.2), col_width - Inches(0.4), Inches(0.5))
            tf = title_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = col_items[0].get("text", "")
            p.font.size = Pt(18)
            p.font.bold = True
            p.font.color.rgb = RGBColor(0, 82, 204)
            p.alignment = PP_ALIGN.CENTER
            
            # 列内容
            content_box = slide.shapes.add_textbox(left_pos + Inches(0.2), Inches(2.8), col_width - Inches(0.4), Inches(3.5))
            tf = content_box.text_frame
            tf.word_wrap = True
            
            for item_idx, item in enumerate(col_items[1:]):
                if item_idx == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = item.get("text", "")
                p.font.size = Pt(14)
                p.font.color.rgb = RGBColor(60, 64, 67)
                p.space_before = Pt(8)
    else:
        # 普通布局
        content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(11.333), Inches(5.0))
        tf = content_box.text_frame
        tf.word_wrap = True
        
        for item_idx, item in enumerate(items):
            level = item.get("level", 0)
            text = item.get("text", "")
            
            if item_idx == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            
            p.text = text
            p.level = level
            p.space_after = Pt(8)
            
            if level == 0:
                p.font.size = Pt(20)
                p.font.color.rgb = RGBColor(15, 25, 45)
            else:
                p.font.size = Pt(16)
                p.font.color.rgb = RGBColor(60, 64, 67)
    
    # 页码
    page_box = slide.shapes.add_textbox(Inches(11.5), Inches(7.0), Inches(1.5), Inches(0.3))
    tf = page_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{page_num} / {total_pages}"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(153, 153, 153)
    p.alignment = PP_ALIGN.RIGHT


def main():
    # 读取JSON文件
    json_path = os.path.expanduser("~/Documents/trae_projects/OpenCopilot/output/AI_Agent_深度研究报告_完整版.json")
    
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    # 生成PPT
    output_dir = os.path.expanduser("~/Documents/trae_projects/OpenCopilot/output")
    os.makedirs(output_dir, exist_ok=True)
    
    ppt_path = os.path.join(output_dir, "AI_Agent_研究报告_HTML转PPT.pptx")
    create_ppt_from_html_data(data['slides'], ppt_path)
    
    print(f"✅ PPT已生成: {ppt_path}")
    
    # 打开PPT
    import subprocess
    subprocess.run(['open', ppt_path])


if __name__ == "__main__":
    main()
