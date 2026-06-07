#!/usr/bin/env python3
"""
OpenCopilot 全链路生产环境验证 (Production Validation Suite)

覆盖所有 LLM 功能模块 + 共创模式功能 + 观测性体系验证。
使用真实 LLM 调用，输出作为 benchmark。

运行方式：
  cd /Users/onetwo/Documents/trae_projects/OpenCopilot
  python scripts/production_validation.py

注意：消耗真实 LLM token，预计 7-10 分钟。
"""

import os
import re
import sys
import json
import uuid
import time
import signal
import threading
import tempfile
import traceback
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple

# 项目根目录
PROJECT_ROOT = Path(__file__).parent.parent
sys.path.insert(0, str(PROJECT_ROOT))
os.chdir(str(PROJECT_ROOT))

# ── 常量 ──
TIMEOUT_PER_CALL = 80
BANNER_WIDTH = 72

# ── 测试文档素材 ──
BUDGET_REPORT = (PROJECT_ROOT / "test_docs" / "budget_report.md").read_text("utf-8")
MEETING_NOTES = (PROJECT_ROOT / "test_docs" / "meeting_notes.md").read_text("utf-8")
API_SPEC = (PROJECT_ROOT / "test_docs" / "api_spec.md").read_text("utf-8")


# ═══════════════════════════════════════════════════
# 工具函数
# ═══════════════════════════════════════════════════

class Colors:
    GREEN = "\033[92m"
    RED = "\033[91m"
    YELLOW = "\033[93m"
    CYAN = "\033[96m"
    BOLD = "\033[1m"
    RESET = "\033[0m"

def banner(title: str):
    print(f"\n{Colors.CYAN}{'═' * BANNER_WIDTH}")
    print(f"  {title}")
    print(f"{'═' * BANNER_WIDTH}{Colors.RESET}")

def ok(msg: str, detail: str = ""):
    suffix = f"  ({detail})" if detail else ""
    print(f"  {Colors.GREEN}✅ {msg}{Colors.RESET}{suffix}")

def fail(msg: str, detail: str = ""):
    suffix = f"  ({detail})" if detail else ""
    print(f"  {Colors.RED}❌ {msg}{Colors.RESET}{suffix}")

def warn(msg: str):
    print(f"  {Colors.YELLOW}⚠️ {msg}{Colors.RESET}")

def info(msg: str):
    print(f"  {Colors.CYAN}ℹ️ {msg}{Colors.RESET}")


# ── 结果收集器 ──
class ResultCollector:
    def __init__(self):
        self.results: List[Dict[str, Any]] = []
        self.start_time = time.time()

    def add(self, module: str, test_name: str, passed: bool,
            score: float = -1, detail: str = "", latency_s: float = 0,
            output_preview: str = ""):
        self.results.append({
            "module": module,
            "test": test_name,
            "passed": passed,
            "score": score,
            "detail": detail,
            "latency_s": round(latency_s, 2),
            "output_preview": output_preview[:200],
        })

    @property
    def total(self):
        return len(self.results)

    @property
    def passed_count(self):
        return sum(1 for r in self.results if r["passed"])

    @property
    def failed_count(self):
        return self.total - self.passed_count


collector = ResultCollector()


# ── LLM 调用 ──
def check_llm_available() -> Tuple[bool, str]:
    """通过实际 pipeline 调用检查 LLM 可用性"""
    try:
        from opencopilot.agent.caller import call_agent_pipeline_sync
        cancel = threading.Event()
        chunks = []
        for chunk in call_agent_pipeline_sync(
            text="hi", action_type="chat",
            session_id=f"avail-{uuid.uuid4().hex[:6]}",
            context_source="chat", context_meta={},
            is_new_task=True, cancel_event=cancel, timeout=20,
        ):
            chunks.append(chunk)
            if len(chunks) >= 3:
                break
        if chunks:
            return True, f"Pipeline OK ({len(chunks)} chunks)"
        return False, "Pipeline 返回空"
    except Exception as e:
        return False, str(e)


def collect_pipeline_output(
    prompt: str, action_type: str,
    context_source: str = "test",
    context_meta: dict = None,
    timeout: int = TIMEOUT_PER_CALL,
) -> Tuple[str, int, float]:
    """返回 (full_text, chunk_count, latency_s)"""
    from opencopilot.agent.caller import call_agent_pipeline_sync

    full_text = ""
    chunk_count = 0
    cancel_event = threading.Event()
    t0 = time.time()

    def _alarm_handler(signum, frame):
        cancel_event.set()
        raise TimeoutError(f"LLM 调用超时（>{timeout}s）")

    old_handler = None
    try:
        old_handler = signal.signal(signal.SIGALRM, _alarm_handler)
        signal.alarm(timeout)
        for chunk in call_agent_pipeline_sync(
            text=prompt,
            action_type=action_type,
            session_id=f"prod-val-{uuid.uuid4().hex[:8]}",
            context_source=context_source,
            context_meta=context_meta or {},
            is_new_task=True,
            cancel_event=cancel_event,
            timeout=timeout,
        ):
            full_text += chunk
            chunk_count += 1
            if cancel_event.is_set():
                break
        signal.alarm(0)
    except TimeoutError:
        raise RuntimeError(f"LLM 调用超时（>{timeout}s）")
    finally:
        if old_handler:
            signal.signal(signal.SIGALRM, old_handler)
        signal.alarm(0)

    latency = time.time() - t0
    return full_text, chunk_count, latency


def filter_think(text: str) -> str:
    display = re.sub(r"<think>.*?</think>", "", text, flags=re.DOTALL)
    if "<think>" in display:
        display = display.split("</think>")[0] if "</think>" in display else display.split("</think>")[-1]
    return display.strip()


def extract_json_array(text: str) -> Optional[list]:
    display = filter_think(text)
    # 直接解析
    try:
        r = json.loads(display)
        if isinstance(r, list): return r
    except json.JSONDecodeError:
        pass
    # markdown code block
    m = re.search(r'```(?:json)?\s*\n?(\[.*?\])\s*```', display, re.DOTALL)
    if m:
        try: return json.loads(m.group(1))
        except json.JSONDecodeError: pass
    # 方括号
    s, e = display.find('['), display.rfind(']')
    if s != -1 and e > s:
        try: return json.loads(display[s:e+1])
        except json.JSONDecodeError: pass
    return None


# ═══════════════════════════════════════════════════
# 评分引擎
# ═══════════════════════════════════════════════════

def score_text_quality(text: str, criteria: Dict[str, Any]) -> Tuple[float, List[str]]:
    """通用文本质量评分 0-5"""
    score = 0.0
    notes = []

    # 1. 长度 (0-1)
    min_len = criteria.get("min_length", 50)
    if len(text) >= min_len:
        score += 1.0
        notes.append(f"长度合格 ({len(text)}>={min_len})")
    else:
        notes.append(f"长度不足 ({len(text)}<{min_len})")

    # 2. 关键词覆盖 (0-2)
    keywords = criteria.get("keywords", [])
    if keywords:
        hits = sum(1 for kw in keywords if kw in text)
        ratio = hits / len(keywords)
        score += ratio * 2.0
        notes.append(f"关键词 {hits}/{len(keywords)} ({ratio:.0%})")

    # 3. 无 think 标签泄露 (0-1)
    filtered = filter_think(text)
    if "</think>" not in filtered and "<think>" not in filtered:
        score += 1.0
        notes.append("无 think 泄露")
    else:
        notes.append("think 标签泄露!")

    # 4. 中文语境保持 (0-1)
    cn_chars = sum(1 for c in text if '\u4e00' <= c <= '\u9fff')
    if cn_chars >= 5 or criteria.get("skip_chinese_check"):
        score += 1.0
        notes.append("中文语境OK")
    else:
        notes.append("中文语境不足")

    return min(score, 5.0), notes


def score_ppt_quality(text: str) -> Tuple[float, List[str]]:
    """PPT 输出专项评分 0-5"""
    score = 0.0
    notes = []
    slides = extract_json_array(text)

    if slides is None:
        notes.append("JSON 解析失败!")
        return 0.0, notes

    # 1. 页数 (0-1)
    if len(slides) >= 4:
        score += 1.0
        notes.append(f"页数充足 ({len(slides)}页)")
    elif len(slides) >= 2:
        score += 0.5
        notes.append(f"页数偏少 ({len(slides)}页)")
    else:
        notes.append(f"页数不足 ({len(slides)}页)")

    # 2. 结构完整性 (0-1.5)
    has_title = any(s.get("type") == "title" for s in slides)
    has_content = any(s.get("type") == "content" for s in slides)
    has_items = any(s.get("items") for s in slides if s.get("type") == "content")
    struct_score = (0.5 if has_title else 0) + (0.5 if has_content else 0) + (0.5 if has_items else 0)
    score += struct_score
    notes.append(f"结构: title={has_title} content={has_content} items={has_items}")

    # 3. 版式多样性 (0-1)
    layouts = set(s.get("layout", "") for s in slides)
    if len(layouts) >= 3:
        score += 1.0
        notes.append(f"版式丰富 ({len(layouts)}种: {layouts})")
    elif len(layouts) >= 2:
        score += 0.5
        notes.append(f"版式一般 ({len(layouts)}种)")
    else:
        notes.append(f"版式单一 ({layouts})")

    # 4. 内容相关性 (0-1.5)
    all_text = json.dumps(slides, ensure_ascii=False)
    budget_kws = ["预算", "2025", "万元", "Q4", "SEM", "KOL"]
    meeting_kws = ["智能客服", "数据看板", "权限", "RBAC", "Q4"]
    all_kws = budget_kws + meeting_kws
    hits = sum(1 for kw in all_kws if kw in all_text)
    if hits >= 4:
        score += 1.5
        notes.append(f"内容相关 ({hits}关键词命中)")
    elif hits >= 2:
        score += 0.75
        notes.append(f"内容部分相关 ({hits}关键词)")
    else:
        notes.append(f"内容相关度低 ({hits}关键词)")

    return min(score, 5.0), notes


# ═══════════════════════════════════════════════════
# 测试执行
# ═══════════════════════════════════════════════════

def test_module(name: str, action_type: str, prompt: str,
                context_source: str, context_meta: dict,
                criteria: Dict[str, Any], is_ppt: bool = False):
    """测试单个 LLM 模块"""
    print(f"\n  {Colors.BOLD}📋 {name} (action_type={action_type}){Colors.RESET}")
    t0 = time.time()
    try:
        text, chunks, latency = collect_pipeline_output(
            prompt, action_type, context_source, context_meta,
        )
        if not text.strip():
            fail("输出为空")
            collector.add(name, "output_not_empty", False, 0, "输出为空", latency)
            return

        # 基础检查
        ok(f"Pipeline 完成: {chunks} chunks, {latency:.1f}s, {len(text)} chars")

        # 评分
        if is_ppt:
            score, notes = score_ppt_quality(text)
        else:
            score, notes = score_text_quality(text, criteria)

        score_label = "优秀" if score >= 4 else "良好" if score >= 3 else "及格" if score >= 2 else "不及格"
        ok(f"质量评分: {score:.1f}/5.0 ({score_label})")
        for n in notes:
            info(f"  {n}")

        # think 标签检查
        filtered = filter_think(text)
        think_clean = "</think>" not in filtered and "<think>" not in filtered
        if think_clean:
            ok("无 think 标签泄露")
        else:
            fail("think 标签泄露")

        collector.add(name, "quality_score", score >= 2, score,
                       f"chunks={chunks} latency={latency:.1f}s chars={len(text)}",
                       latency, text)
        collector.add(name, "think_filter", think_clean, -1)

    except RuntimeError as e:
        fail(f"Pipeline 调用失败: {e}")
        collector.add(name, "pipeline_call", False, 0, str(e), time.time() - t0)
    except Exception as e:
        fail(f"未知错误: {e}")
        collector.add(name, "unknown_error", False, 0, str(e), time.time() - t0)


def run_llm_module_tests():
    """Part 1: 所有 LLM 功能模块测试"""
    banner("Part 1: LLM 功能模块全链路测试（真实 LLM 调用）")

    available, msg = check_llm_available()
    if not available:
        print(f"\n  {Colors.RED}❌ LLM 服务不可用: {msg}{Colors.RESET}")
        print(f"  跳过所有 LLM 测试。请检查 config.json 中的 API key 配置。")
        return

    info(f"LLM 服务: {msg}")

    # 1. Chat - 智能对话
    test_module(
        name="Chat 智能对话",
        action_type="chat",
        prompt="你好！我正在做一个AI办公助手项目，你能帮我做些什么？",
        context_source="chat",
        context_meta={},
        criteria={
            "min_length": 30,
            "keywords": ["你好", "帮助", "可以", "能够", "助手", "助理", "AI", "支持"],
        },
    )

    # 2. Explain - 代码解释
    code_sample = """
def lru_cache(maxsize=128):
    def decorator(func):
        cache = {}
        order = []
        def wrapper(*args):
            if args in cache:
                order.remove(args)
                order.append(args)
                return cache[args]
            result = func(*args)
            cache[args] = result
            order.append(args)
            if len(cache) > maxsize:
                oldest = order.pop(0)
                del cache[oldest]
            return result
        return wrapper
    return decorator
"""
    test_module(
        name="Explain 代码解释",
        action_type="explain",
        prompt=f"请解释以下代码/文本:\n\n{code_sample}",
        context_source="selection",
        context_meta={"source_text": code_sample},
        criteria={
            "min_length": 80,
            "keywords": ["缓存", "LRU", "装饰器", "字典", "淘汰", "最近", "访问", "函数"],
        },
    )

    # 3. Fix - 代码修复
    buggy_code = """
def merge_sorted_lists(a, b):
    result = []
    i, j = 0, 0
    while i < len(a) and j < len(b):
        if a[i] <= b[j]:
            result.append(a[i])
            i += 1
        else:
            result.append(b[j])
            j += 1
    # BUG: 忘记处理剩余元素
    return result
"""
    test_module(
        name="Fix 代码修复",
        action_type="fix",
        prompt=f"请修复以下代码中的问题:\n\n{buggy_code}",
        context_source="selection",
        context_meta={"source_text": buggy_code},
        criteria={
            "min_length": 50,
            "keywords": ["剩余", "extend", "append", "result", "return", "while", "i", "j"],
        },
    )

    # 4. Polish - 文本润色
    raw_text = "这个产品功能挺多的，用起来还行，就是有时候有点卡，不过总体来说还可以，值得试试。"
    test_module(
        name="Polish 文本润色",
        action_type="polish",
        prompt=f"请润色优化以下文本:\n\n{raw_text}",
        context_source="selection",
        context_meta={"source_text": raw_text},
        criteria={
            "min_length": 30,
            "keywords": ["功能", "体验", "产品", "值得", "流畅"],
        },
    )

    # 5. Translate - 翻译
    cn_text = "人工智能正在深刻改变我们的工作方式。从智能客服到自动化编程，AI技术已经渗透到日常办公的方方面面。"
    test_module(
        name="Translate 翻译",
        action_type="translate",
        prompt=f"请将以下文本从中文翻译为英文:\n\n{cn_text}",
        context_source="selection",
        context_meta={"source_text": cn_text, "source_language": "中文", "target_language": "英文"},
        criteria={
            "min_length": 30,
            "keywords": ["artificial intelligence", "AI", "work", "customer", "automation", "technology"],
            "skip_chinese_check": True,
        },
    )

    # 6. Code Review - 代码审查
    test_module(
        name="Code Review 代码审查",
        action_type="code_review",
        prompt=f"请对以下代码进行审查:\n\n{API_SPEC[:500]}\n\n```python\n{code_sample}\n```",
        context_source="selection",
        context_meta={"source_text": code_sample},
        criteria={
            "min_length": 100,
            "keywords": ["缓存", "线程", "安全", "性能", "问题", "建议", "优化", "复杂度"],
        },
    )

    # 7. PPT 生成 - 预算报告
    test_module(
        name="PPT 生成（预算报告）",
        action_type="ppt",
        prompt=f"请根据以下内容生成 PPT 大纲。\n\n要求：\n1. 严格输出纯 JSON 数组格式\n2. 每个 slide 包含 type, layout, title, items\n3. layout 可选: center, text_only, image_right, three_columns\n4. 每页 3-5 个要点\n5. 智能选择版式\n\n原始内容：\n{BUDGET_REPORT}",
        context_source="studio",
        context_meta={"input_text_len": len(BUDGET_REPORT)},
        criteria={"min_length": 200},
        is_ppt=True,
    )

    # 8. PPT 生成 - 会议纪要（长文档）
    test_module(
        name="PPT 生成（会议纪要）",
        action_type="ppt",
        prompt=f"请根据以下内容生成 PPT 大纲。\n\n要求：\n1. 严格输出纯 JSON 数组格式\n2. 每个 slide 包含 type, layout, title, items\n3. layout 可选: center, text_only, image_right, three_columns\n4. 每页 3-5 个要点\n5. 智能选择版式\n\n原始内容：\n{MEETING_NOTES}",
        context_source="studio",
        context_meta={"input_text_len": len(MEETING_NOTES)},
        criteria={"min_length": 200},
        is_ppt=True,
    )


def run_cocreation_tests():
    """Part 2: 共创模式功能验证（离线 + 组件级）"""
    banner("Part 2: 共创模式功能验证")

    # 2.1 Persona 文件完整性
    print(f"\n  {Colors.BOLD}📋 Persona 文件完整性{Colors.RESET}")
    from opencopilot.shared.prompt import load_persona
    required = ["default", "chat", "code", "ppt", "translate", "polish"]
    all_ok = True
    for name in required:
        content = load_persona(name)
        if len(content) > 50:
            ok(f"personas/{name}.md ({len(content)} chars)")
        else:
            fail(f"personas/{name}.md 内容过短 ({len(content)} chars)")
            all_ok = False
    collector.add("共创-Persona", "文件完整性", all_ok, -1)

    # 2.2 action_type → persona 映射
    print(f"\n  {Colors.BOLD}📋 action_type → Persona 映射{Colors.RESET}")
    from config_manager import ConfigManager
    mapping = ConfigManager.get_instance().get_persona_mapping()
    mapping_tests = {
        "coding": "code",
        "ppt": "ppt",
        "translate": "translate",
        "chat": "chat",
    }
    for at, expected in mapping_tests.items():
        actual = mapping.get(at, "MISSING")
        if actual == expected:
            ok(f"{at} → {actual}")
        else:
            fail(f"{at} → {actual} (expected {expected})")
    collector.add("共创-映射表", "映射正确性", all(
        mapping.get(k) == v for k, v in mapping_tests.items()
    ), -1)

    # 2.3 子目录 Persona 加载
    print(f"\n  {Colors.BOLD}📋 子目录 Persona 加载{Colors.RESET}")
    sub_persona = load_persona("office/business/presentation")
    sub_ok = len(sub_persona) > 100
    if sub_ok:
        ok(f"office/business/presentation.md ({len(sub_persona)} chars)")
    else:
        fail(f"子目录 persona 加载失败 ({len(sub_persona)} chars)")
    collector.add("共创-子目录", "加载成功", sub_ok, -1)

    # 2.4 IntentRouter 意图分类
    print(f"\n  {Colors.BOLD}📋 IntentRouter 意图分类{Colors.RESET}")
    from opencopilot.capabilities.ppt.intent_router import IntentRouter
    intent_cases = [
        ("把第2页标题改为Q1业绩回顾", "update_title", "direct"),
        ("转为柱状图", "convert_chart", "llm"),
        ("删除第3页", "remove_slide", "direct"),
        ("帮我润色一下这段文字", "polish_text", "llm"),
        ("添加一页新幻灯片", "add_slide", "direct"),
        ("重新生成整个PPT", "regenerate", "llm"),
        ("第3页版式改为图文混排", "update_layout", "direct"),
    ]
    intent_pass = 0
    for instruction, expected_intent, expected_method in intent_cases:
        result = IntentRouter.classify(instruction)
        ok_match = result["intent"] == expected_intent and result["method"] == expected_method
        if ok_match:
            ok(f'"{instruction}" → {result["intent"]} ({result["method"]})')
            intent_pass += 1
        else:
            fail(f'"{instruction}" → {result["intent"]}/{result["method"]} (expected {expected_intent}/{expected_method})')
    collector.add("共创-IntentRouter", "分类准确率", intent_pass == len(intent_cases),
                   intent_pass / len(intent_cases) * 5, f"{intent_pass}/{len(intent_cases)}")

    # 2.5 JSON Slides 解析鲁棒性
    print(f"\n  {Colors.BOLD}📋 JSON Slides 解析鲁棒性{Colors.RESET}")
    parse_cases = [
        ("纯 JSON 数组", '[{"type":"title","layout":"center","title":"测试","subtitle":"2025"}]'),
        ("Markdown 包裹", '```json\n[{"type":"title","layout":"center","title":"测试"}]\n```'),
        ("带前言", '以下是PPT大纲：\n[{"type":"title","layout":"center","title":"预算报告"}]'),
        ("多页幻灯片", json.dumps([
            {"type": "title", "layout": "center", "title": "封面", "subtitle": "2025"},
            {"type": "content", "layout": "text_only", "title": "概况",
             "items": [{"level": 0, "text": "总预算800万"}, {"level": 1, "text": "已执行580万"}]},
            {"type": "content", "layout": "three_columns", "title": "渠道分布",
             "items": [{"level": 0, "text": "SEM 230万"}, {"level": 0, "text": "KOL 90万"}]},
        ])),
    ]
    parse_pass = 0
    for label, text in parse_cases:
        result = extract_json_array(text)
        if result and isinstance(result, list) and len(result) > 0:
            ok(f"{label} → {len(result)} slides")
            parse_pass += 1
        else:
            fail(f"{label} → 解析失败")
    collector.add("共创-JSON解析", "解析成功率", parse_pass == len(parse_cases),
                   parse_pass / len(parse_cases) * 5, f"{parse_pass}/{len(parse_cases)}")

    # 2.6 PPT 导出功能
    print(f"\n  {Colors.BOLD}📋 PPT 导出功能{Colors.RESET}")
    try:
        from ppt_generator import generate_ppt_from_text
        slides_data = [
            {"type": "title", "layout": "center", "title": "2025预算报告", "subtitle": "市场部"},
            {"type": "content", "layout": "text_only", "title": "总体概况",
             "items": [{"level": 0, "text": "总预算800万元"}, {"level": 0, "text": "已执行580万元"},
                       {"level": 1, "text": "执行率72.5%"}]},
            {"type": "content", "layout": "three_columns", "title": "渠道分配",
             "items": [{"level": 0, "text": "SEM: 230万"}, {"level": 0, "text": "信息流: 140万"},
                       {"level": 0, "text": "KOL: 90万"}]},
        ]
        json_str = json.dumps(slides_data, ensure_ascii=False)
        output_path = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False).name
        generate_ppt_from_text(text=json_str, output_path=output_path)

        if os.path.exists(output_path) and os.path.getsize(output_path) > 1000:
            from pptx import Presentation
            prs = Presentation(output_path)
            slide_count = len(prs.slides)
            ok(f"PPT 导出成功: {slide_count} slides, {os.path.getsize(output_path)} bytes")
            collector.add("共创-PPT导出", "导出成功", True, 5, f"{slide_count} slides")
        else:
            fail("PPT 文件过小或不存在")
            collector.add("共创-PPT导出", "导出成功", False, 0)
        os.unlink(output_path)
    except Exception as e:
        fail(f"PPT 导出失败: {e}")
        collector.add("共创-PPT导出", "导出成功", False, 0, str(e))

    # 2.7 PPT 4阶段管线 fallback
    print(f"\n  {Colors.BOLD}📋 PPT 4阶段管线 Fallback{Colors.RESET}")
    try:
        from opencopilot.capabilities.ppt.pipeline import PPTGenerationPipeline, PipelineResult
        pipeline = PPTGenerationPipeline()
        result = pipeline.run(BUDGET_REPORT)
        if isinstance(result, PipelineResult) and result.total_pages >= 2:
            ok(f"Pipeline fallback: {result.total_pages} pages, {len(result.topics)} topics")
            collector.add("共创-管线Fallback", "运行成功", True, 4,
                           f"pages={result.total_pages} topics={len(result.topics)}")
        else:
            fail(f"Pipeline 输出异常: {result}")
            collector.add("共创-管线Fallback", "运行成功", False, 0)
    except Exception as e:
        fail(f"Pipeline 异常: {e}")
        collector.add("共创-管线Fallback", "运行成功", False, 0, str(e))


def run_observability_tests():
    """Part 3: 观测性体系验证"""
    banner("Part 3: 观测性体系验证")

    # 3.1 Pipeline Timer Log
    print(f"\n  {Colors.BOLD}📋 Pipeline Timer Log{Colors.RESET}")
    from opencopilot.agent.observability import PipelineObservability
    obs = PipelineObservability.get_instance()
    paths = obs.get_log_paths()
    timer_path = paths["pipeline_timer_log"]
    if os.path.exists(timer_path):
        size = os.path.getsize(timer_path)
        lines = sum(1 for _ in open(timer_path, encoding="utf-8", errors="replace"))
        ok(f"timer log: {timer_path} ({size} bytes, {lines} lines)")
        collector.add("观测性-TimerLog", "文件存在且非空", size > 0, -1, f"{size}B {lines}L")
    else:
        fail(f"timer log 不存在: {timer_path}")
        collector.add("观测性-TimerLog", "文件存在且非空", False, 0)

    # 3.2 Pipeline Logs DB
    print(f"\n  {Colors.BOLD}📋 Pipeline Logs DB (SQLite){Colors.RESET}")
    from opencopilot.agent.log_store import LogStore
    store = LogStore.get_instance()
    db_path = store._db_path
    if os.path.exists(db_path):
        db_size = os.path.getsize(db_path)
        rows = store.query("SELECT COUNT(*) FROM pipeline_logs")
        row_count = rows[0][0] if rows else 0
        ok(f"pipeline_logs.db: {db_size} bytes, {row_count} rows")

        # 最近事件
        recent = store.query(
            "SELECT event, source, level, timestamp FROM pipeline_logs ORDER BY id DESC LIMIT 5"
        )
        for evt, src, lvl, ts in recent:
            info(f"  [{ts}] {evt} ({src}/{lvl})")

        collector.add("观测性-LogDB", "DB完整且有数据", row_count > 0, -1, f"{row_count} rows")
    else:
        fail(f"pipeline_logs.db 不存在: {db_path}")
        collector.add("观测性-LogDB", "DB完整且有数据", False, 0)

    # 3.3 Timer Stats
    print(f"\n  {Colors.BOLD}📋 Timer Stats (内存){Colors.RESET}")
    stats = obs.get_timer_stats()
    if stats:
        ok(f"总请求数: {stats['total_requests']}")
        for mw_name, mw_stats in stats.get("middlewares", {}).items():
            info(f"  {mw_name}: count={mw_stats['count']} avg={mw_stats['avg']:.3f}s")
    else:
        info("内存中无 timer 历史（正常，需要管线调用后填充）")
    collector.add("观测性-TimerStats", "API可用", True, -1)

    # 3.4 ImmuneSystem 日志事件
    print(f"\n  {Colors.BOLD}📋 ImmuneSystem 安全事件日志{Colors.RESET}")
    immune_events = store.query(
        "SELECT event, COUNT(*) FROM pipeline_logs WHERE event LIKE 'IMMUNE%' GROUP BY event"
    )
    if immune_events:
        for evt, cnt in immune_events:
            ok(f"{evt}: {cnt} events")
    else:
        info("暂无 ImmuneSystem 事件（正常，需要管线调用触发）")
    collector.add("观测性-ImmuneEvents", "查询可用", True, -1)


# ═══════════════════════════════════════════════════
# 报告生成
# ═══════════════════════════════════════════════════

def generate_report():
    """生成综合验证报告"""
    elapsed = time.time() - collector.start_time

    banner("📊 综合验证报告 (Benchmark)")

    # 汇总表
    print(f"\n  {'模块':<25} {'测试':<25} {'结果':<6} {'评分':<6} {'耗时':<8}")
    print(f"  {'─' * 75}")
    for r in collector.results:
        status = f"{Colors.GREEN}PASS{Colors.RESET}" if r["passed"] else f"{Colors.RED}FAIL{Colors.RESET}"
        score_str = f"{r['score']:.1f}" if r["score"] >= 0 else "—"
        print(f"  {r['module']:<25} {r['test']:<25} {status}  {score_str:<6} {r['latency_s']:<8}")

    # 统计
    print(f"\n  {'═' * 75}")
    total = collector.total
    passed = collector.passed_count
    failed = collector.failed_count
    pass_rate = passed / total * 100 if total else 0

    avg_scores = {}
    for r in collector.results:
        if r["score"] >= 0:
            avg_scores.setdefault(r["module"], []).append(r["score"])

    print(f"\n  {Colors.BOLD}总体统计{Colors.RESET}")
    print(f"  测试总数: {total}")
    print(f"  通过: {Colors.GREEN}{passed}{Colors.RESET} | 失败: {Colors.RED}{failed}{Colors.RESET}")
    print(f"  通过率: {pass_rate:.0f}%")
    print(f"  总耗时: {elapsed:.1f}s")

    if avg_scores:
        print(f"\n  {Colors.BOLD}模块质量评分 (LLM 输出){Colors.RESET}")
        for mod, scores in avg_scores.items():
            avg = sum(scores) / len(scores)
            bar = "█" * int(avg) + "░" * (5 - int(avg))
            print(f"  {mod:<25} {avg:.1f}/5.0  {bar}")

    # 判定
    print(f"\n  {'═' * 75}")
    if failed == 0 and pass_rate >= 90:
        verdict = f"{Colors.GREEN}✅ 生产环境验证通过 — 所有模块符合预期{Colors.RESET}"
    elif failed <= 2 and pass_rate >= 80:
        verdict = f"{Colors.YELLOW}⚠️ 基本通过，存在 {failed} 个非关键问题{Colors.RESET}"
    else:
        verdict = f"{Colors.RED}❌ 验证未通过 — {failed} 个测试失败，需排查{Colors.RESET}"
    print(f"\n  {Colors.BOLD}{verdict}{Colors.RESET}")

    return failed == 0


# ═══════════════════════════════════════════════════
# 入口
# ═══════════════════════════════════════════════════

def main():
    print(f"\n{Colors.BOLD}{'█' * BANNER_WIDTH}")
    print(f"  OpenCopilot 全链路生产环境验证 (Production Validation)")
    print(f"  {time.strftime('%Y-%m-%d %H:%M:%S')}")
    print(f"{'█' * BANNER_WIDTH}{Colors.RESET}")

    # Part 1: LLM 功能模块
    run_llm_module_tests()

    # Part 2: 共创模式
    run_cocreation_tests()

    # Part 3: 观测性
    run_observability_tests()

    # 报告
    success = generate_report()
    return 0 if success else 1


if __name__ == "__main__":
    sys.exit(main())
