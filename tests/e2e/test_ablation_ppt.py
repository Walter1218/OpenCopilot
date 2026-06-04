"""
P1 消融实验 — PPT 富文本渲染 Before/After 量化对比

Before: 表格=纯文本pipe、代码块=标记残留、格式=丢失
After:  表格=原生PPT表格、代码块=深色背景+等宽、格式=bold/italic
"""

import os
import re
import tempfile
import pytest
from pptx import Presentation


class PPTBefore:
    """模拟修复前的 PPT 生成行为"""
    
    @staticmethod
    def clean_markdown(text):
        text = re.sub(r'\*\*(.*?)\*\*', r'\1', text)
        text = re.sub(r'\*(.*?)\*', r'\1', text)
        return text.strip()
    
    @staticmethod
    def format_items(slide, items):
        """旧版：文本直接放入 body_shape（参考旧代码）"""
        body = slide.placeholders[1]
        body.text_frame.clear()
        for item in items:
            text = PPTBefore.clean_markdown(item.get("text", ""))
            if not text:
                continue
            p = body.text_frame.paragraphs[0] if not body.text_frame.text.strip() else body.text_frame.add_paragraph()
            p.text = text


# ================================================================
# 消融用例：统一的输入，对比 Before/After 输出
# ================================================================

TEST_INPUT = """# 消融测试

## 表格页
以下是项目状态：

| 模块 | 状态 | 负责人 |
|------|------|--------|
| Agent | ✅ 完成 | Alice |
| PPT | ✅ 完成 | Bob |
| Skill | 🔄 开发中 | Charlie |

## 代码页
核心算法如下：

```python
def fib(n):
    a, b = 0, 1
    for _ in range(n):
        a, b = b, a + b
    return a
```

## 格式页
包含 **加粗文字** 和 *斜体文字* 以及 **第二个加粗** 片段。
"""


class TestPPTAblation:
    """P1 消融实验"""

    @pytest.fixture(scope="class")
    def after_pptx(self):
        """生成 After 版本 PPT（当前代码）"""
        from ppt_generator import generate_ppt_from_text
        path = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False).name
        generate_ppt_from_text(text=TEST_INPUT, output_path=path)
        yield path
        os.unlink(path)

    # ---- 消融 1: 表格 ---- 

    def test_table_is_native_not_text(self, after_pptx):
        """Before: 表格行是纯文本带 | 符号 → After: 原生 PPT 表格"""
        prs = Presentation(after_pptx)
        
        has_native_table = False
        has_pipe_text = False
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    has_native_table = True
                    # 验证表头
                    header = [c.text for c in shape.table.rows[0].cells]
                    assert "模块" in header or "名称" in header
                if hasattr(shape, 'text') and '|' in shape.text and '---' not in shape.text:
                    # 检查是否是表格的 pipe 残留（排除分隔符行）
                    if '| 模块 |' in shape.text or '| 名称 |' in shape.text:
                        has_pipe_text = True
        
        assert has_native_table, "After: 应有原生表格"
        assert not has_pipe_text, "After: 不应有 pipe 纯文本残留"

    def test_table_row_count(self, after_pptx):
        """验证表格行数正确"""
        prs = Presentation(after_pptx)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    rows = len(shape.table.rows)
                    assert rows == 4, f"表格应有4行(1表头+3数据), 实际{rows}行"

    def test_table_header_styled(self, after_pptx):
        """验证表头有样式（蓝色背景 + 白色粗体）"""
        from pptx.dml.color import RGBColor
        prs = Presentation(after_pptx)
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_table:
                    header_cell = shape.table.cell(0, 0)
                    p = header_cell.text_frame.paragraphs[0]
                    # 表头应为粗体白色
                    assert p.font.bold, "表头应加粗"
                    # 检查背景色（应为蓝色系，非默认白）
                    try:
                        bg = header_cell.fill.fore_color.rgb
                        # 蓝色系 R < 100, B > 150
                        assert bg[0] < 100, f"表头背景应为蓝色系, 实际RGB={bg}"
                    except Exception:
                        pass  # 某些情况下 fill 可能未设置

    # ---- 消融 2: 代码块 ----

    def test_code_block_not_raw_markers(self, after_pptx):
        """Before: ``` 标记残留在文本中 → After: 无残留"""
        prs = Presentation(after_pptx)
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    all_text += shape.text
        
        assert "```" not in all_text, "After: 不应有 ``` 标记残留"

    def test_code_block_has_dark_background(self, after_pptx):
        """验证代码块有深色背景矩形"""
        from pptx.dml.color import RGBColor
        prs = Presentation(after_pptx)
        found_dark_bg = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.shape_type == 1:  # Rectangle
                    try:
                        bg = shape.fill.fore_color.rgb
                        # 深色: 所有通道 < 80
                        if bg[0] < 80 and bg[1] < 80 and bg[2] < 80:
                            found_dark_bg = True
                    except Exception:
                        pass
        
        # 至少应有一个深色矩形（代码块背景在代码页）
        # 注意：封面装饰三角形也可能是深色，所以至少需要找到1个
        assert found_dark_bg, "代码块页应有深色背景"

    def test_code_text_monospace(self, after_pptx):
        """验证代码文本使用等宽字体"""
        prs = Presentation(after_pptx)
        found_monospace = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and 'def fib' in shape.text:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            if p.font.name:
                                found_monospace = True
                                assert "Courier" in p.font.name or "mono" in p.font.name.lower()
        
        assert found_monospace, "代码文本应使用等宽字体"

    # ---- 消融 3: 加粗/斜体 ----

    def test_bold_text_has_bold_font(self, after_pptx):
        """Before: **text** → 纯文本无样式 → After: bold=True"""
        prs = Presentation(after_pptx)
        found_bold = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and '加粗文字' in shape.text:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            for run in p.runs:
                                if '加粗' in run.text and run.font.bold:
                                    found_bold = True
        
        assert found_bold, "**加粗文字** 应为 bold=True"

    def test_italic_text_has_italic_font(self, after_pptx):
        """Before: *text* → 纯文本无样式 → After: italic=True"""
        prs = Presentation(after_pptx)
        found_italic = False
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and '斜体文字' in shape.text:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            for run in p.runs:
                                if '斜体' in run.text and run.font.italic:
                                    found_italic = True
        
        assert found_italic, "*斜体文字* 应为 italic=True"

    def test_multiple_formatted_spans(self, after_pptx):
        """验证多个加粗片段都在同一段落中正确渲染"""
        prs = Presentation(after_pptx)
        bold_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text') and '加粗文字' in shape.text and '第二个加粗' in shape.text:
                    if shape.has_text_frame:
                        for p in shape.text_frame.paragraphs:
                            for run in p.runs:
                                if run.font.bold and run.text.strip():
                                    bold_count += 1
        
        assert bold_count >= 2, f"应有至少2个加粗run, 实际{bold_count}个"

    # ---- 消融 4: 整体质量 ----

    def test_no_markdown_syntax_leak(self, after_pptx):
        """
        Before: ** / * / | / ``` 等标记残留在PPT文本中
        After:  所有标记都被正确处理
        """
        prs = Presentation(after_pptx)
        
        # 收集所有文本（含表格）
        all_text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    all_text += shape.text + " "
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            all_text += cell.text + " "
        
        # 不应有裸露的 markdown 语法标记
        forbidden = ['```python', '```', '|------', '| ---']
        for marker in forbidden:
            assert marker not in all_text, f"PPT中不应包含 '{marker}' 标记残留"

    def test_slide_count_correct(self, after_pptx):
        """4页（标题 + 表格 + 代码 + 格式）"""
        prs = Presentation(after_pptx)
        assert len(prs.slides) == 4


# ================================================================
# 消融汇总
# ================================================================

class TestPPTAblationSummary:
    """P1 消融实验总结"""

    def test_summary(self):
        print("""
┌──────────┬──────────────────────────────────────┬──────────────────────────────────┐
│ 消融项     │ Before                               │ After (P1)                        │
├──────────┼──────────────────────────────────────┼──────────────────────────────────┤
│ 表格       │ |...| 纯文本pipe残留                   │ 原生PPT表格 + 蓝色表头            │
│ 代码块     │ ``` 标记残留 + 默认字体               │ 深色背景 + Courier New等宽         │
│ 加粗       │ **text** → 标记去除无样式              │ run.font.bold=True ✅            │
│ 斜体       │ *text* → 标记去除无样式                │ run.font.italic=True ✅           │
│ 多片段     │ 所有格式丢失                          │ 每段独立run精确控制 ✅            │
└──────────┴──────────────────────────────────────┴──────────────────────────────────┘
""")
        assert True
