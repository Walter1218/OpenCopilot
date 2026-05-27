"""
文件处理工具
"""

import os
import asyncio
from typing import Any, Dict, Optional
from .base import BaseTool


class FileReadTool(BaseTool):
    """文件读取工具"""
    
    @property
    def name(self) -> str:
        return "file_read"
    
    @property
    def description(self) -> str:
        return "读取文件内容，支持多种格式"
    
    @property
    def parameters(self) -> Dict[str, Any]:
        return {
            "file_path": {
                "type": "string",
                "description": "文件路径",
                "required": True
            },
            "format": {
                "type": "string",
                "description": "文件格式",
                "enum": ["text", "docx", "pptx", "pdf"],
                "default": "text"
            }
        }
    
    async def execute(self, **kwargs) -> Dict[str, Any]:
        """执行文件读取"""
        file_path = kwargs.get("file_path")
        file_format = kwargs.get("format", "text")
        
        if not file_path:
            return {"error": "file_path is required"}
        
        expanded_path = os.path.expanduser(file_path)
        if not os.path.exists(expanded_path):
            return {"error": f"文件不存在: {file_path}"}
        
        try:
            if file_format == "docx":
                return await self._read_docx(expanded_path)
            elif file_format == "pptx":
                return await self._read_pptx(expanded_path)
            elif file_format == "pdf":
                return await self._read_pdf(expanded_path)
            else:
                return await self._read_text(expanded_path)
        except Exception as e:
            return {"error": f"读取文件失败: {str(e)}"}
    
    async def _read_text(self, file_path: str) -> Dict[str, Any]:
        """读取文本文件"""
        loop = asyncio.get_event_loop()
        content = await loop.run_in_executor(None, self._read_text_sync, file_path)
        return {
            "type": "text",
            "content": content,
            "file_path": file_path
        }
    
    def _read_text_sync(self, file_path: str) -> str:
        """同步读取文本文件"""
        with open(file_path, "r", encoding="utf-8") as f:
            return f.read()
    
    async def _read_docx(self, file_path: str) -> Dict[str, Any]:
        """读取Word文档"""
        try:
            from docx import Document
            loop = asyncio.get_event_loop()
            doc = await loop.run_in_executor(None, Document, file_path)
            
            parts = []
            for p in doc.paragraphs:
                if p.text.strip():
                    parts.append(p.text)
            
            # 提取表格内容
            for table in doc.tables:
                for row in table.rows:
                    row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        parts.append(row_text)
            
            return {
                "type": "docx",
                "content": '\n\n'.join(parts),
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
                "file_path": file_path
            }
        except ImportError:
            return {"error": "需要安装 python-docx: pip install python-docx"}
    
    async def _read_pptx(self, file_path: str) -> Dict[str, Any]:
        """读取PPT文件"""
        try:
            from pptx import Presentation
            loop = asyncio.get_event_loop()
            prs = await loop.run_in_executor(None, Presentation, file_path)
            
            slides = []
            for i, slide in enumerate(prs.slides):
                texts = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                texts.append(t)
                if texts:
                    slides.append(f"--- 幻灯片 {i + 1} ---\n" + '\n'.join(texts))
            
            return {
                "type": "pptx",
                "content": '\n\n'.join(slides),
                "slide_count": len(prs.slides),
                "file_path": file_path
            }
        except ImportError:
            return {"error": "需要安装 python-pptx: pip install python-pptx"}
    
    async def _read_pdf(self, file_path: str) -> Dict[str, Any]:
        """读取PDF文件"""
        try:
            import PyPDF2
            loop = asyncio.get_event_loop()
            
            def read_pdf_sync():
                with open(file_path, "rb") as f:
                    reader = PyPDF2.PdfReader(f)
                    parts = []
                    for page in reader.pages:
                        text = page.extract_text()
                        if text:
                            parts.append(text)
                    return '\n\n'.join(parts), len(reader.pages)
            
            content, page_count = await loop.run_in_executor(None, read_pdf_sync)
            
            return {
                "type": "pdf",
                "content": content,
                "page_count": page_count,
                "file_path": file_path
            }
        except ImportError:
            return {"error": "需要安装 PyPDF2: pip install PyPDF2"}


class FileWriteTool(BaseTool):
    """文件写入工具"""
    
    @property
    def name(self) -> str:
        return "file_write"
    
    @property
    def description(self) -> str:
        return "将内容写入文件"
    
    @property
    def parameters(self) -> Dict[str, Any]:
        return {
            "content": {
                "type": "string",
                "description": "要写入的内容",
                "required": True
            },
            "file_path": {
                "type": "string",
                "description": "文件路径",
                "required": True
            },
            "format": {
                "type": "string",
                "description": "文件格式",
                "enum": ["text", "docx", "pptx"],
                "default": "text"
            }
        }
    
    async def execute(self, **kwargs) -> Dict[str, Any]:
        """执行文件写入"""
        content = kwargs.get("content")
        file_path = kwargs.get("file_path")
        file_format = kwargs.get("format", "text")
        
        if not content or not file_path:
            return {"error": "content and file_path are required"}
        
        expanded_path = os.path.expanduser(file_path)
        
        try:
            if file_format == "docx":
                return await self._write_docx(content, expanded_path)
            elif file_format == "pptx":
                return await self._write_pptx(content, expanded_path)
            else:
                return await self._write_text(content, expanded_path)
        except Exception as e:
            return {"error": f"写入文件失败: {str(e)}"}
    
    async def _write_text(self, content: str, file_path: str) -> Dict[str, Any]:
        """写入文本文件"""
        loop = asyncio.get_event_loop()
        await loop.run_in_executor(None, self._write_text_sync, content, file_path)
        return {"success": True, "file_path": file_path}
    
    def _write_text_sync(self, content: str, file_path: str):
        """同步写入文本文件"""
        os.makedirs(os.path.dirname(file_path), exist_ok=True)
        with open(file_path, "w", encoding="utf-8") as f:
            f.write(content)
    
    async def _write_docx(self, content: str, file_path: str) -> Dict[str, Any]:
        """写入Word文档"""
        try:
            from docx import Document
            loop = asyncio.get_event_loop()
            
            def write_docx_sync():
                doc = Document()
                # 简单的文本写入，实际应用中可能需要解析Markdown
                for line in content.split('\n'):
                    if line.strip():
                        doc.add_paragraph(line)
                os.makedirs(os.path.dirname(file_path), exist_ok=True)
                doc.save(file_path)
            
            await loop.run_in_executor(None, write_docx_sync)
            return {"success": True, "file_path": file_path}
        except ImportError:
            return {"error": "需要安装 python-docx: pip install python-docx"}
    
    async def _write_pptx(self, content: str, file_path: str) -> Dict[str, Any]:
        """写入PPT文件"""
        try:
            from pptx import Presentation
            loop = asyncio.get_event_loop()
            
            def write_pptx_sync():
                prs = Presentation()
                # 简单的文本写入，实际应用中可能需要解析Markdown
                for line in content.split('\n'):
                    if line.strip():
                        slide = prs.slides.add_slide(prs.slide_layouts[1])
                        slide.shapes.title.text = line
                os.makedirs(os.path.dirname(file_path), exist_ok=True)
                prs.save(file_path)
            
            await loop.run_in_executor(None, write_pptx_sync)
            return {"success": True, "file_path": file_path}
        except ImportError:
            return {"error": "需要安装 python-pptx: pip install python-pptx"}


class FileConvertTool(BaseTool):
    """文件格式转换工具"""
    
    @property
    def name(self) -> str:
        return "file_convert"
    
    @property
    def description(self) -> str:
        return "转换文件格式"
    
    @property
    def parameters(self) -> Dict[str, Any]:
        return {
            "input_path": {
                "type": "string",
                "description": "输入文件路径",
                "required": True
            },
            "output_format": {
                "type": "string",
                "description": "输出格式",
                "enum": ["pdf", "docx", "pptx", "txt", "md"],
                "required": True
            },
            "output_path": {
                "type": "string",
                "description": "输出文件路径（可选）"
            }
        }
    
    async def execute(self, **kwargs) -> Dict[str, Any]:
        """执行格式转换"""
        input_path = kwargs.get("input_path")
        output_format = kwargs.get("output_format")
        output_path = kwargs.get("output_path")
        
        if not input_path or not output_format:
            return {"error": "input_path and output_format are required"}
        
        expanded_input = os.path.expanduser(input_path)
        if not os.path.exists(expanded_input):
            return {"error": f"输入文件不存在: {input_path}"}
        
        # 如果没有指定输出路径，自动生成
        if not output_path:
            base_name = os.path.splitext(expanded_input)[0]
            output_path = f"{base_name}.{output_format}"
        
        expanded_output = os.path.expanduser(output_path)
        
        try:
            # 读取输入文件
            read_tool = FileReadTool()
            input_format = self._detect_format(expanded_input)
            read_result = await read_tool.execute(file_path=expanded_input, format=input_format)
            
            if "error" in read_result:
                return read_result
            
            content = read_result.get("content", "")
            
            # 写入输出文件
            write_tool = FileWriteTool()
            write_result = await write_tool.execute(
                content=content,
                file_path=expanded_output,
                format=output_format
            )
            
            return write_result
        except Exception as e:
            return {"error": f"格式转换失败: {str(e)}"}
    
    def _detect_format(self, file_path: str) -> str:
        """检测文件格式"""
        ext = os.path.splitext(file_path)[1].lower()
        format_mapping = {
            ".txt": "text",
            ".md": "text",
            ".docx": "docx",
            ".pptx": "pptx",
            ".pdf": "pdf"
        }
        return format_mapping.get(ext, "text")